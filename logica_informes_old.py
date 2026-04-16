import pandas as pd
from pyathena import connect
from openpyxl import load_workbook
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.utils import get_column_letter
import os
import io
import locale
from datetime import datetime, timedelta
from dateutil.relativedelta import relativedelta
import shutil
import re
import glob
from pathlib import Path
import tempfile
import base64
import time
from dotenv import load_dotenv

from typing import List, Dict, Any, Tuple, Optional

# --- Librerías Fase 2 ---
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import google.generativeai as genai
from PIL import Image


# --- Importaciones para Google Drive API (añadidas aquí para que el archivo sea autocontenido) ---
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload
import streamlit as st # Necesario para st.info/st.warning/st.error en las funciones auxiliares


# --- Manejo de Comtypes para Windows ---
COMTYPES_AVAILABLE = False
try:
    import comtypes.client
    COMTYPES_AVAILABLE = True
except ImportError:
    print("ADVERTENCIA: Librería 'comtypes' no disponible. La captura de imágenes de PowerPoint está desactivada.")

load_dotenv()


def download_file_from_drive(service, file_id: str) -> bytes:
    """Descarga el contenido de un archivo de Drive por su ID."""
    try:
        # Para archivos de Google Sheets, usa export_media para obtener XLSX
        # Comprobar el mimeType del archivo en Drive para decidir si exportar o descargar directamente
        file_metadata = service.files().get(fileId=file_id, fields='mimeType').execute()
        if file_metadata['mimeType'] == 'application/vnd.google-apps.spreadsheet':
            request = service.files().export_media(fileId=file_id, mimeType='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
        else:
            request = service.files().get_media(fileId=file_id)
        
        file_content = io.BytesIO()
        downloader = MediaIoBaseDownload(file_content, request)
        done = False
        while done is False:
            status, done = downloader.next_chunk()
        return file_content.getvalue()
    except HttpError as e:
        st.error(f"Error al descargar archivo de Drive (ID: {file_id}): {e.content.decode('utf-8')}")
        raise
    except Exception as e:
        st.error(f"Error inesperado al descargar archivo de Drive (ID: {file_id}): {e}")
        raise

def upload_file_to_drive(service, filename: str, content_bytes: bytes, mime_type: str, parent_folder_id: Optional[str] = None) -> str:
    """Sube un archivo a Google Drive, actualizando si ya existe por nombre en la carpeta."""
    file_metadata = {'name': filename}
    if parent_folder_id:
        file_metadata['parents'] = [parent_folder_id]

    media_body = io.BytesIO(content_bytes)
    
    existing_file_id = None
    try:
        query = f"name = '{filename}' and trashed = false"
        if parent_folder_id:
            query += f" and '{parent_folder_id}' in parents"
        response = service.files().list(q=query, spaces='drive', fields='files(id, name)').execute()
        files = response.get('files', [])
        if files:
            existing_file_id = files[0]['id']
            st.info(f"Archivo existente '{filename}' encontrado (ID: {existing_file_id}). Se actualizará.")
    except Exception as e:
        st.warning(f"No se pudo verificar la existencia del archivo: {e}")

    try:
        if existing_file_id:
            request = service.files().update(fileId=existing_file_id,
                                             media_body=media_body,
                                             body=file_metadata,
                                             mimeType=mime_type,
                                             fields='id, name')
        else:
            request = service.files().create(media_body=media_body,
                                             body=file_metadata,
                                             mimeType=mime_type,
                                             fields='id, name')

        response = request.execute()
        st.success(f"Archivo '{filename}' {'actualizado' if existing_file_id else 'subido'} a Google Drive con ID: {response.get('id')}")
        return response.get('id')
    except Exception as e:
        st.error(f"Error al subir archivo a Drive ('{filename}'): {e}")
        raise

def find_or_create_folder(service, folder_name: str, parent_folder_id: Optional[str] = None) -> str:
    """Busca una carpeta por nombre dentro de una carpeta padre, o la crea si no existe."""
    query = f"name = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
    if parent_folder_id:
        query += f" and '{parent_folder_id}' in parents"

    try:
        response = service.files().list(q=query, spaces='drive', fields='files(id)').execute()
        folders = response.get('files', [])
        if folders:
            return folders[0]['id']
        else:
            file_metadata = {
                'name': folder_name,
                'mimeType': 'application/vnd.google-apps.folder'
            }
            if parent_folder_id:
                file_metadata['parents'] = [parent_folder_id]
            
            folder = service.files().create(body=file_metadata, fields='id').execute()
            st.info(f"Carpeta '{folder_name}' creada con ID: {folder.get('id')}")
            return folder.get('id')
    except Exception as e:
        st.error(f"Error al buscar o crear carpeta '{folder_name}': {e}")
        raise


##########################################################################################
# ----------------------------------- FASE 1 ------------------------------------------- #
##########################################################################################

def ejecutar_fase_1(load_ids, fecha_cierre_dt, plantillas_bytes, credenciales):
    """
    Ejecuta toda la lógica de la Fase 1: Conexión a Athena, procesamiento de datos
    y actualización de archivos Excel.

    Args:
        load_ids (dict): Diccionario con los IDs de carga.
        fecha_cierre_dt (datetime): Fecha de cierre para el informe.
        plantillas_bytes (dict): Diccionario con el contenido en bytes de las plantillas Excel.
                                 Ej: {'coap': b'...', 'efectos': b'...', 'datos_medios': b'...'}
        credenciales (dict): Diccionario con las credenciales de AWS.

    Returns:
        dict: Un diccionario con los nombres y contenido en bytes de los archivos Excel generados.
    """
    # 1. Configurar locale
    try:
        locale.setlocale(locale.LC_TIME, 'es_ES.UTF-8')
    except locale.Error:
        try:
            locale.setlocale(locale.LC_TIME, 'Spanish_Spain.1252')
        except locale.Error:
            print("Advertencia: No se pudo configurar el locale a español.")

    # 2. Cargar Workbooks desde los bytes en memoria
    try:
        wb = load_workbook(filename=io.BytesIO(plantillas_bytes['coap']))
        wb_1 = load_workbook(filename=io.BytesIO(plantillas_bytes['efectos']))
        wb_2 = load_workbook(filename=io.BytesIO(plantillas_bytes['datos_medios']))
    except Exception as e:
        raise IOError(f"Error al cargar una de las plantillas de Excel: {e}")

    # 3. Cálculo de fechas y nombres
    mes_actual_nombre = fecha_cierre_dt.strftime("%B").capitalize()[:3]
    fecha_mes_anterior_dt = (fecha_cierre_dt.replace(day=1) - timedelta(days=1))
    mes_anterior_nombre = fecha_mes_anterior_dt.strftime("%B").capitalize()[:3]
    mes_anterior_dic = {'mes_anterior': fecha_cierre_dt.replace(day=1).strftime('%Y-%m-01')}

    print(f"mes_anterior_dic: {mes_anterior_dic}")

    # --- Lógica de copiar hojas ---
    # def copiar_hoja_con_nuevo_mes(workbook, nombre_base_hoja, mes_anterior_str, mes_actual_str):
    #     nombre_hoja_anterior = f"{nombre_base_hoja} - {mes_anterior_str}"
    #     nombre_hoja_actual = f"{nombre_base_hoja} - {mes_actual_str}"
    #     if nombre_hoja_actual not in workbook.sheetnames:
    #         if nombre_hoja_anterior in workbook.sheetnames:
    #             hoja_anterior = workbook[nombre_hoja_anterior]
    #             workbook.copy_worksheet(hoja_anterior).title = nombre_hoja_actual
    #         else:
    #             workbook.create_sheet(title=nombre_hoja_actual)

    #     return nombre_hoja_actual

    def copiar_hoja_con_nuevo_mes(workbook, nombre_base_hoja, mes_anterior_str, mes_actual_str):
        """
        Copia una hoja existente (del mes anterior) a una nueva hoja con el nombre del mes actual.
        Si la hoja del mes anterior no existe, crea una hoja vacía con el nombre del mes actual.
        Si la hoja del mes actual ya existe, no hace nada.
        """
        nombre_hoja_anterior = f"{nombre_base_hoja} - {mes_anterior_str}"
        nombre_hoja_actual = f"{nombre_base_hoja} - {mes_actual_str}"

        if nombre_hoja_actual not in workbook.sheetnames:
            if nombre_hoja_anterior in workbook.sheetnames:
                hoja_anterior = workbook[nombre_hoja_anterior]
                workbook.copy_worksheet(hoja_anterior).title = nombre_hoja_actual
                print(f"Hoja '{nombre_hoja_anterior}' copiada a '{nombre_hoja_actual}'")
            else:
                workbook.create_sheet(title=nombre_hoja_actual)
                print(f"Hoja del mes anterior '{nombre_hoja_anterior}' no encontrada. Se creó una nueva hoja '{nombre_hoja_actual}' ")
        else:
            print(f"La hoja '{nombre_hoja_actual}' ya existe. No se realizaron cambios.")
        # return nombre_hoja_actual
    
    nombres_base_hojas_wb = [
        "VE- Sensibilidades", "Sensibilidades VE EBA", "MF - Sensibilidades - Mto",
        "Sensibilidades MF - EBA", "MF - Sensibilidades - Cto"
    ]
    for base_name in nombres_base_hojas_wb:
        copiar_hoja_con_nuevo_mes(wb, base_name, mes_anterior_nombre, mes_actual_nombre)

    # 4. Conexión a Athena
    print("Estableciendo conexión con Amazon Athena...")
    conn = connect(
        aws_access_key_id=credenciales["aws_id"],
        aws_secret_access_key=credenciales["aws_key"],
        s3_staging_dir=credenciales["s3_dir"],
        region_name=credenciales["region"]
    )
    print("Conexión exitosa.")

    # 5. Ejecución de queries y pegado de datos
    print("Ejecutando queries y actualizando hojas...")
    Informe_VE_01= f"""
    WITH jerarquias AS (
        SELECT  1 AS orden, '1.1 BANCOS'                        AS Jerarquia UNION ALL
        SELECT  2,           '1.2 CARTERA MAYORISTA'                         UNION ALL
        SELECT  3,           '1.3 DEPOSITOS OTRAS EECC'                      UNION ALL
        SELECT  4,           '1.4 INVERSION CREDITICIA'                      UNION ALL
        SELECT  5,           '1.5 RENTA FIJA'                                UNION ALL
        SELECT  6,           '1.6 ACTIVOS DUDOSOS'                           UNION ALL
        SELECT  7,           '1.7 ACTIVOS NO SENSIBLE'                       UNION ALL
        SELECT  9,           '2.1 DEPOSITOS OTRAS EECC'                      UNION ALL
        SELECT 10,           '2.2 DEPOSITOS A PLAZO DE CLIENTES'             UNION ALL
        SELECT 11,           '2.3 DEPOSITOS A LA VISTA DE CLIENTES'          UNION ALL
        SELECT 13,           '2.4 PASIVOS NO SENSIBLES'                      UNION ALL
        SELECT 12,           '2.5 CARTERA MAYORISTA'                         UNION ALL
        SELECT 15,           '3.1 DERIVADOS'                                 UNION ALL
        SELECT 17,           '4.1 PN'
    ),

    metrics_agg AS (
        SELECT
            CASE
                WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') THEN '1.1 BANCOS'
                WHEN dim_1 IN ('A_Ata') THEN '1.2 CARTERA MAYORISTA'
                WHEN dim_1 IN ('A_Depositos cedidos') THEN '1.3 DEPOSITOS OTRAS EECC'
                WHEN dim_1 IN (
                    'A_ECORP_Otros_Activo Corto Plazo',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                    'A_EESP_Otros_Activo Corto Plazo',
                    'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                    'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                    'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                    'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                    'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                    'A_PIB_Otros_Descubiertos_Gtia Personal',
                    'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                    'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                    'A_POFI_Otros_Descubiertos_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                    'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                    'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_POFI_Otros_Prest. Origen_Gtia Aval',
                    'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                    'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                    'A_SSCC_Otros_Descubiertos_Gtia Personal',
                    'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                    'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
                ) THEN '1.4 INVERSION CREDITICIA'
                WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') THEN '1.5 RENTA FIJA'
                WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') THEN '1.6 ACTIVOS DUDOSOS'
                WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') THEN '1.7 ACTIVOS NO SENSIBLE'
                WHEN dim_1 IN ('P_Cuentas de otras EECC') THEN '2.1 DEPOSITOS OTRAS EECC'
                WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES'
                WHEN dim_1 IN (
                'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
                'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
                ) THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES'
                WHEN dim_1 IN ('P_Resto no sensible') THEN '2.4 PASIVOS NO SENSIBLES'
                WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') THEN '2.5 CARTERA MAYORISTA'
                WHEN dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo') THEN '3.1 DERIVADOS'
                WHEN dim_1 IN ('Patrimonio Neto') THEN '4.1 PN'
                ELSE 'OTROS'
            END AS Jerarquia,
            SUM(CASE WHEN scenario = 'Dwn200' THEN market_value ELSE 0 END) AS Dwn200,
            SUM(CASE WHEN scenario = 'Dwn100' THEN market_value ELSE 0 END) AS Dwn100,
            SUM(CASE WHEN scenario = 'Dwn75'  THEN market_value ELSE 0 END) AS Dwn75,
            SUM(CASE WHEN scenario = 'Dwn50'  THEN market_value ELSE 0 END) AS Dwn50,
            SUM(CASE WHEN scenario = 'Dwn25'  THEN market_value ELSE 0 END) AS Dwn25,
            SUM(CASE WHEN scenario = 'Base'   THEN market_value ELSE 0 END) AS Base,
            SUM(CASE WHEN scenario = 'Up25'   THEN market_value ELSE 0 END) AS Up25,
            SUM(CASE WHEN scenario = 'Up50'   THEN market_value ELSE 0 END) AS Up50,
            SUM(CASE WHEN scenario = 'Up75'   THEN market_value ELSE 0 END) AS Up75,
            SUM(CASE WHEN scenario = 'Up100'  THEN market_value ELSE 0 END) AS Up100,
            SUM(CASE WHEN scenario = 'Up200'  THEN market_value ELSE 0 END) AS Up200
        FROM pro_pichincha_alquid_old.metric
        WHERE load_id IN (
        '{load_ids['cierre_base']}',
        '{load_ids['cierre_up']}',
        '{load_ids['cierre_dwn']}'
        )
        GROUP BY 1
    )

    -- 1) JERARQUÍAS FIJAS (1.x, 2.x, 3.1, 4.1) CON LEFT JOIN
    SELECT
        j.Jerarquia,
        a.Dwn200,
        a.Dwn100,
        a.Dwn75,
        a.Dwn50,
        a.Dwn25,
        a.Base,
        a.Up25,
        a.Up50,
        a.Up75,
        a.Up100,
        a.Up200
    FROM jerarquias j
    LEFT JOIN metrics_agg a
        ON a.Jerarquia = j.Jerarquia

    UNION ALL

    -- 2) TOTAL GRUPO 1 (igual que antes)
    SELECT 'TOTAL GRUPO 1' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200' THEN market_value ELSE 0 END) AS Dwn200,
        SUM(CASE WHEN scenario = 'Dwn100' THEN market_value ELSE 0 END) AS Dwn100,
        SUM(CASE WHEN scenario = 'Dwn75'  THEN market_value ELSE 0 END) AS Dwn75,
        SUM(CASE WHEN scenario = 'Dwn50'  THEN market_value ELSE 0 END) AS Dwn50,
        SUM(CASE WHEN scenario = 'Dwn25'  THEN market_value ELSE 0 END) AS Dwn25,
        SUM(CASE WHEN scenario = 'Base'   THEN market_value ELSE 0 END) AS Base,
        SUM(CASE WHEN scenario = 'Up25'   THEN market_value ELSE 0 END) AS Up25,
        SUM(CASE WHEN scenario = 'Up50'   THEN market_value ELSE 0 END) AS Up50,
        SUM(CASE WHEN scenario = 'Up75'   THEN market_value ELSE 0 END) AS Up75,
        SUM(CASE WHEN scenario = 'Up100'  THEN market_value ELSE 0 END) AS Up100,
        SUM(CASE WHEN scenario = 'Up200'  THEN market_value ELSE 0 END) AS Up200
    FROM pro_pichincha_alquid_old.metric
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )
    AND dim_1 in ('A_Cuentas en otras EECC', 'A_Tesoreria', 'A_Ata', 'A_Depositos cedidos',
    'A_ECORP_Otros_Activo Corto Plazo', 'A_EESP_Otros_Activo Corto Plazo', 'A_POFI_Otros_Ant. Nomina_Gtia Personal',
    'A_ECORP_Otros_Cred. Empresa_Gtia Aval', 'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
    'A_EESP_Otros_Cred. Empresa_Gtia Aval', 'A_EESP_Otros_Cred. Empresa_Gtia Personal',
    'A_SSCC_Otros_Cred. Empresa_Gtia Aval', 'A_SSCC_Otros_Cred. Empresa_Gtia Hipot.', 'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
    'A_PIB_Otros_Descubiertos_Gtia Personal', 'A_POFI_Otros_Descubiertos_Gtia Personal', 'A_SSCC_Otros_Descubiertos_Gtia Personal',
    'A_ECORP_Otros_Prest. Empresas_Gtia Aval', 'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.', 'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
    'A_ECORP_Otros_Prest. Empresas_Gtia Prenda', 'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',
    'A_EESP_Otros_Prest. Empresas_Gtia Aval', 'A_EESP_Otros_Prest. Empresas_Gtia Personal', 'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
    'A_EESP_Otros_Prest. Promotor_Gtia Hipot.', 'A_PIB_Otros_Prest. Consumo_Gtia Personal', 'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
    'A_PIB_Otros_Prest. Hipotecas_Gtia Personal', 'A_POFI_Otros_Prest. Consumo_Gtia Personal', 'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
    'A_POFI_Otros_Prest. Empresas_Gtia Aval', 'A_POFI_Otros_Prest. Empresas_Gtia Hipot.', 'A_POFI_Otros_Prest. Empresas_Gtia Personal',
    'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
    'A_POFI_Otros_Prest. Origen_Gtia Aval', 'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
    'A_SSCC_BKIA_Prest. Consumo_Gtia Personal', 'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
    'A_SSCC_Otros_Prest. Empresas_Gtia Aval', 'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
    'A_SSCC_Otros_Prest. Empresas_Gtia Personal', 'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
    'A_PIB_Otros_Tarjetas Credito_Gtia Personal', 'A_POFI_Otros_Tarjetas Credito_Gtia Personal',
    'A_SSCC_Otros_Tarjetas Credito_Gtia Personal', 'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal',
    'A_Bonos corporativos', 'A_Bonos soberanos', 'A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso',
    'A_Resto no sensible', 'A_Tesoreria_Admin.', 'A_EESP_Otros_Prest. Promotor_Gtia Aval', 'A_EESP_Otros_Prest. Promotor_Gtia Personal')

    UNION ALL

    -- 3) TOTAL GRUPO 2
    SELECT 'TOTAL GRUPO 2' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200' THEN market_value ELSE 0 END) AS Dwn200,
        SUM(CASE WHEN scenario = 'Dwn100' THEN market_value ELSE 0 END) AS Dwn100,
        SUM(CASE WHEN scenario = 'Dwn75'  THEN market_value ELSE 0 END) AS Dwn75,
        SUM(CASE WHEN scenario = 'Dwn50'  THEN market_value ELSE 0 END) AS Dwn50,
        SUM(CASE WHEN scenario = 'Dwn25'  THEN market_value ELSE 0 END) AS Dwn25,
        SUM(CASE WHEN scenario = 'Base'   THEN market_value ELSE 0 END) AS Base,
        SUM(CASE WHEN scenario = 'Up25'   THEN market_value ELSE 0 END) AS Up25,
        SUM(CASE WHEN scenario = 'Up50'   THEN market_value ELSE 0 END) AS Up50,
        SUM(CASE WHEN scenario = 'Up75'   THEN market_value ELSE 0 END) AS Up75,
        SUM(CASE WHEN scenario = 'Up100'  THEN market_value ELSE 0 END) AS Up100,
        SUM(CASE WHEN scenario = 'Up200'  THEN market_value ELSE 0 END) AS Up200
    FROM pro_pichincha_alquid_old.metric
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )
    AND dim_1 in ('P_Cuentas de otras EECC', 'P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo',
    'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
    'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
    'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas', 'P_Resto no sensible', 'P_Ata', 'P_Emisiones T2')

    UNION ALL

    -- 4) TOTAL GRUPO 3
    SELECT 'TOTAL GRUPO 3' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200' THEN market_value ELSE 0 END) AS Dwn200,
        SUM(CASE WHEN scenario = 'Dwn100' THEN market_value ELSE 0 END) AS Dwn100,
        SUM(CASE WHEN scenario = 'Dwn75'  THEN market_value ELSE 0 END) AS Dwn75,
        SUM(CASE WHEN scenario = 'Dwn50'  THEN market_value ELSE 0 END) AS Dwn50,
        SUM(CASE WHEN scenario = 'Dwn25'  THEN market_value ELSE 0 END) AS Dwn25,
        SUM(CASE WHEN scenario = 'Base'   THEN market_value ELSE 0 END) AS Base,
        SUM(CASE WHEN scenario = 'Up25'   THEN market_value ELSE 0 END) AS Up25,
        SUM(CASE WHEN scenario = 'Up50'   THEN market_value ELSE 0 END) AS Up50,
        SUM(CASE WHEN scenario = 'Up75'   THEN market_value ELSE 0 END) AS Up75,
        SUM(CASE WHEN scenario = 'Up100'  THEN market_value ELSE 0 END) AS Up100,
        SUM(CASE WHEN scenario = 'Up200'  THEN market_value ELSE 0 END) AS Up200
    FROM pro_pichincha_alquid_old.metric
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )
    AND dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo')

    UNION ALL

    -- 5) TOTAL GRUPO 4
    SELECT 'TOTAL GRUPO 4' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200' THEN market_value ELSE 0 END) AS Dwn200,
        SUM(CASE WHEN scenario = 'Dwn100' THEN market_value ELSE 0 END) AS Dwn100,
        SUM(CASE WHEN scenario = 'Dwn75'  THEN market_value ELSE 0 END) AS Dwn75,
        SUM(CASE WHEN scenario = 'Dwn50'  THEN market_value ELSE 0 END) AS Dwn50,
        SUM(CASE WHEN scenario = 'Dwn25'  THEN market_value ELSE 0 END) AS Dwn25,
        SUM(CASE WHEN scenario = 'Base'   THEN market_value ELSE 0 END) AS Base,
        SUM(CASE WHEN scenario = 'Up25'   THEN market_value ELSE 0 END) AS Up25,
        SUM(CASE WHEN scenario = 'Up50'   THEN market_value ELSE 0 END) AS Up50,
        SUM(CASE WHEN scenario = 'Up75'   THEN market_value ELSE 0 END) AS Up75,
        SUM(CASE WHEN scenario = 'Up100'  THEN market_value ELSE 0 END) AS Up100,
        SUM(CASE WHEN scenario = 'Up200'  THEN market_value ELSE 0 END) AS Up200
    FROM pro_pichincha_alquid_old.metric
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )
    AND dim_1 IN ('Patrimonio Neto')

    ORDER BY CASE Jerarquia
        WHEN '1.1 BANCOS'                         THEN 1
        WHEN '1.2 CARTERA MAYORISTA'             THEN 2
        WHEN '1.3 DEPOSITOS OTRAS EECC'          THEN 3
        WHEN '1.4 INVERSION CREDITICIA'          THEN 4
        WHEN '1.5 RENTA FIJA'                    THEN 5
        WHEN '1.6 ACTIVOS DUDOSOS'               THEN 6
        WHEN '1.7 ACTIVOS NO SENSIBLE'           THEN 7
        WHEN 'TOTAL GRUPO 1'                     THEN 8
        WHEN '2.1 DEPOSITOS OTRAS EECC'          THEN 9
        WHEN '2.2 DEPOSITOS A PLAZO DE CLIENTES' THEN 10
        WHEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES' THEN 11
        WHEN '2.5 CARTERA MAYORISTA'             THEN 12
        WHEN '2.4 PASIVOS NO SENSIBLES'          THEN 13
        WHEN 'TOTAL GRUPO 2'                     THEN 14
        WHEN '3.1 DERIVADOS'                     THEN 15
        WHEN 'TOTAL GRUPO 3'                     THEN 16
        WHEN '4.1 PN'                            THEN 17
        WHEN 'TOTAL GRUPO 4'                     THEN 18
        ELSE 999
    END;
    """

    Informe_VE_02= f"""SELECT
                        dim_1 AS Jerarquia,
                    SUM(CASE WHEN SCENARIO = 'Base' THEN effective_duration * market_value ELSE 0 END) /
                        NULLIF(SUM(CASE WHEN SCENARIO = 'Base' THEN market_value ELSE 0 END), 0)/100 AS effective_duration_PONDERADA,
                    SUM(CASE WHEN scenario = 'Dwn200' THEN market_value ELSE 0 END) AS Dwn200,
                    SUM(CASE WHEN scenario = 'Dwn100' THEN market_value ELSE 0 END) AS Dwn100,
                    SUM(CASE WHEN scenario = 'Dwn75' THEN market_value ELSE 0 END) AS Dwn75,
                    SUM(CASE WHEN scenario = 'Dwn50' THEN market_value ELSE 0 END) AS Dwn50,
                    SUM(CASE WHEN scenario = 'Dwn25' THEN market_value ELSE 0 END) AS Dwn25,
                    SUM(CASE WHEN scenario = 'Base' THEN market_value ELSE 0 END) AS Base,
                    SUM(CASE WHEN scenario = 'Up25' THEN market_value ELSE 0 END) AS Up25,
                    SUM(CASE WHEN scenario = 'Up50' THEN market_value ELSE 0 END) AS Up50,
                    SUM(CASE WHEN scenario = 'Up75' THEN market_value ELSE 0 END) AS Up75,
                    SUM(CASE WHEN scenario = 'Up100' THEN market_value ELSE 0 END) AS Up100,
                    SUM(CASE WHEN scenario = 'Up200' THEN market_value ELSE 0 END) AS Up200
                    FROM pro_pichincha_alquid_old.metric
                    WHERE load_id IN (
                            '{load_ids['cierre_base']}',
                            '{load_ids['cierre_up']}',
                            '{load_ids['cierre_dwn']}'
                            )
                    AND dim_1 IN (
                        'A_Bonos corporativos',
                        'A_Bonos soberanos',
                        'IRS_pago',
                        'IRS_recibo',
                        'FXSWAP_pago',
                        'FXSWAP_recibo'
                    )
                    GROUP BY dim_1
                    ORDER BY case
                        when dim_1 = 'A_Bonos corporativos' then 1
                        when dim_1 = 'A_Bonos soberanos' then 2
                        when dim_1 = 'IRS_pago' then 3
                        when dim_1 = 'IRS_recibo' then 4
                        when dim_1 = 'FXSWAP_pagos' then 5
                        when dim_1 = 'FXSWAP_recibo' then 6
                    end; """

    Informe_VE_03= f"""SELECT
                        dim_1 AS Jerarquia,
                        SUM(CASE WHEN SCENARIO = 'Base' THEN macaulay_duration * market_value ELSE 0 END) /
                        NULLIF(SUM(CASE WHEN SCENARIO = 'Base' THEN market_value ELSE 0 END), 0) AS DURACION_MACAULAY_PONDERADA,
                    SUM(CASE WHEN scenario = 'Dwn200' THEN market_value ELSE 0 END) AS Dwn200,
                    SUM(CASE WHEN scenario = 'Dwn100' THEN market_value ELSE 0 END) AS Dwn100,
                    SUM(CASE WHEN scenario = 'Dwn75' THEN market_value ELSE 0 END) AS Dwn75,
                    SUM(CASE WHEN scenario = 'Dwn50' THEN market_value ELSE 0 END) AS Dwn50,
                    SUM(CASE WHEN scenario = 'Dwn25' THEN market_value ELSE 0 END) AS Dwn25,
                    SUM(CASE WHEN scenario = 'Base' THEN market_value ELSE 0 END) AS Base,
                    SUM(CASE WHEN scenario = 'Up25' THEN market_value ELSE 0 END) AS Up25,
                    SUM(CASE WHEN scenario = 'Up50' THEN market_value ELSE 0 END) AS Up50,
                    SUM(CASE WHEN scenario = 'Up75' THEN market_value ELSE 0 END) AS Up75,
                    SUM(CASE WHEN scenario = 'Up100' THEN market_value ELSE 0 END) AS Up100,
                    SUM(CASE WHEN scenario = 'Up200' THEN market_value ELSE 0 END) AS Up200
                        FROM pro_pichincha_alquid_old.metric
                            WHERE load_id IN (
                            '{load_ids['cierre_base']}',
                            '{load_ids['cierre_up']}',
                            '{load_ids['cierre_dwn']}'
                            )
                        AND dim_1 IN (
                        'A_Bonos corporativos',
                        'A_Bonos soberanos',
                        'IRS_pago',
                        'IRS_recibo',
                        'FXSWAP_pago',
                        'FXSWAP_recibo')
                        GROUP BY dim_1
                        ORDER BY case
                        when dim_1 = 'A_Bonos corporativos' then 1
                        when dim_1 = 'A_Bonos soberanos' then 2
                        when dim_1 = 'IRS_pago' then 3
                        when dim_1 = 'IRS_recibo' then 4
                        when dim_1 = 'FXSWAP_pagos' then 5
                        when dim_1 = 'FXSWAP_recibo' then 6
                        end; """

    Informe_VE_04= f"""
WITH jerarquias AS (
    SELECT '1.1 BANCOS' AS jerarquia, 1 AS orden UNION ALL
    SELECT '1.2 CARTERA MAYORISTA', 2 UNION ALL
    SELECT '1.3 DEPOSITOS OTRAS EECC', 3 UNION ALL
    SELECT '1.4 INVERSION CREDITICIA', 4 UNION ALL
    SELECT '1.6 ACTIVOS DUDOSOS', 6 UNION ALL
    SELECT '1.7 ACTIVOS NO SENSIBLE', 7 UNION ALL
    SELECT '2.1 DEPOSITOS OTRAS EECC', 8 UNION ALL
    SELECT '2.2 DEPOSITOS A PLAZO DE CLIENTES', 9 UNION ALL
    SELECT '2.5 CARTERA MAYORISTA', 10 UNION ALL
    SELECT '2.4 PASIVOS NO SENSIBLES', 11 UNION ALL
    SELECT '3.1 DERIVADOS', 12 UNION ALL
    SELECT '4.1 PN', 13
),
agg AS (
    SELECT
        CASE
            WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') THEN '1.1 BANCOS'
            WHEN dim_1 IN ('A_Ata') THEN '1.2 CARTERA MAYORISTA'
            WHEN dim_1 IN ('A_Depositos cedidos') THEN '1.3 DEPOSITOS OTRAS EECC'
            WHEN dim_1 IN (
                'A_ECORP_Otros_Activo Corto Plazo',
                'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                'A_EESP_Otros_Activo Corto Plazo',
                'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                'A_PIB_Otros_Descubiertos_Gtia Personal',
                'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                'A_POFI_Otros_Descubiertos_Gtia Personal',
                'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                'A_POFI_Otros_Prest. Origen_Gtia Aval',
                'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                'A_SSCC_Otros_Descubiertos_Gtia Personal',
                'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
            ) THEN '1.4 INVERSION CREDITICIA'
            WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') THEN '1.5 RENTA FIJA'
            WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') THEN '1.6 ACTIVOS DUDOSOS'
            WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') THEN '1.7 ACTIVOS NO SENSIBLE'
            WHEN dim_1 IN ('P_Cuentas de otras EECC') THEN '2.1 DEPOSITOS OTRAS EECC'
            WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES'
            WHEN dim_1 IN (
                'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
                'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
            ) THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES'
            WHEN dim_1 IN ('P_Resto no sensible') THEN '2.4 PASIVOS NO SENSIBLES'
            WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') THEN '2.5 CARTERA MAYORISTA'
            WHEN dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo') THEN '3.1 DERIVADOS'
            WHEN dim_1 IN ('Patrimonio Neto') THEN '4.1 PN'
            ELSE 'OTROS'
        END AS jerarquia,

        SUM(CASE WHEN SCENARIO = 'Base' THEN ABS(effective_duration * market_value) ELSE 0 END)
        / NULLIF(SUM(CASE WHEN SCENARIO = 'Base' THEN ABS(market_value) ELSE 0 END), 0)
        / 100 AS effective_duration_ponderada,

        SUM(CASE WHEN SCENARIO = 'Base' THEN ABS(macaulay_duration * market_value) ELSE 0 END)
        / NULLIF(SUM(CASE WHEN SCENARIO = 'Base' THEN ABS(market_value) ELSE 0 END), 0)
        AS duracion_macaulay_ponderada

    FROM pro_pichincha_alquid_old.metric
    WHERE load_id IN (
        '{load_ids['cierre_base']}',
        '{load_ids['cierre_up']}',
        '{load_ids['cierre_dwn']}'
    )
    GROUP BY 1
)

SELECT
    j.jerarquia,
    COALESCE(a.effective_duration_ponderada, 0) AS effective_duration_ponderada,
    COALESCE(a.duracion_macaulay_ponderada, 0) AS duracion_macaulay_ponderada
FROM jerarquias j
LEFT JOIN agg a
    ON j.jerarquia = a.jerarquia
ORDER BY j.orden;
"""


    # Se define nuestro data frame con los datos extrídos de la base de datos
    df_VE_01 = pd.read_sql(Informe_VE_01, conn)
    df_VE_02 = pd.read_sql(Informe_VE_02, conn)
    df_VE_03 = pd.read_sql(Informe_VE_03, conn)
    df_VE_04 = pd.read_sql(Informe_VE_04, conn)

    # Nombre de la hoja donde quieres pegar los datos. Ahora se construye dinámicamente.
    nombre_hoja_VE = f"VE- Sensibilidades - {mes_actual_nombre}"


    # Acceder a la hoja
    ws_VE = wb[nombre_hoja_VE]


    # Se definen las coordenadas donde queremos que se peguen los datos
    start_row_VE_01 = 4
    start_col_VE_01 = 33

    start_row_VE_02 = 33
    start_col_VE_02 = 32

    start_row_VE_03 = 44
    start_col_VE_03 = 32

    start_row_VE_04 = 4
    start_col_VE_04 = 45


    # Límite de filas donde se pegarán los datos (de la 6 a la 21)
    end_row_VE_01 = 21
    end_row_VE_02 = 38
    end_row_VE_03 = 49
    end_row_VE_04 = 20

    # Filas a evitar
    filas_a_saltar_VE_04 = {8, 11, 14, 17, 19}
    fila_actual_VE_04 = start_row_VE_04
    # Escribir los datos en las coordenadas anteriores

    # Query de pegado para el informe 1
    for i, row in enumerate(df_VE_01.values, start=start_row_VE_01):
        if i > end_row_VE_01:
            break  # Detener la escritura si superamos la fila 21
        for j, value_VE_01 in enumerate(row[1:], start=start_col_VE_01):
            ws_VE.cell(row=i, column=j).value = value_VE_01

    # Query de pegado para el informe 2
    for i, row in enumerate(df_VE_02.values, start=start_row_VE_02):
        if i > end_row_VE_02:
            break  # Detener la escritura si superamos la fila 21
        for j, value_VE_02 in enumerate(row[1:], start=start_col_VE_02):
            ws_VE.cell(row=i, column=j).value = value_VE_02

    # Query de pegado para el informe 3
    for i, row in enumerate(df_VE_03.values, start=start_row_VE_03):
        if i > end_row_VE_03:
            break  # Detener la escritura si superamos la fila 21
        for j, value_VE_03 in enumerate(row[1:], start=start_col_VE_03):
            ws_VE.cell(row=i, column=j).value = value_VE_03

    # Query de pegado para el informe 4
    for row in df_VE_04.values:
        # Saltar las filas especificadas
        while fila_actual_VE_04 in filas_a_saltar_VE_04:
            fila_actual_VE_04 += 1

        if fila_actual_VE_04 > end_row_VE_04:
            break

        for j, value_VE_04 in enumerate(row[1:], start=start_col_VE_04):
            ws_VE.cell(row=fila_actual_VE_04, column=j).value = value_VE_04

        fila_actual_VE_04 += 1  # Mover a la siguiente fila en la hoja



    ###################################### PROCESO DE PEGADO DE INFORMES DE VALOR ECONÓMICO EBA #####################################
    # Se define la query de SQL en función de los filtros y los satos que queremos extraer
    Informe_VE_EBA_01= f"""
        WITH jerarquias AS (
        SELECT '1.1 BANCOS'                         AS Jerarquia UNION ALL
        SELECT '1.2 CARTERA MAYORISTA'                         UNION ALL
        SELECT '1.3 DEPOSITOS OTRAS EECC'                      UNION ALL
        SELECT '1.4 INVERSION CREDITICIA'                      UNION ALL
        SELECT '1.5 RENTA FIJA'                                UNION ALL
        SELECT '1.6 ACTIVOS DUDOSOS'                           UNION ALL
        SELECT '1.7 ACTIVOS NO SENSIBLE'                       UNION ALL
        SELECT '2.1 DEPOSITOS OTRAS EECC'                      UNION ALL
        SELECT '2.2 DEPOSITOS A PLAZO DE CLIENTES'             UNION ALL
        SELECT '2.3 DEPOSITOS A LA VISTA DE CLIENTES'          UNION ALL
        SELECT '2.4 PASIVOS NO SENSIBLES'                      UNION ALL
        SELECT '2.5 CARTERA MAYORISTA'                         UNION ALL
        SELECT '3.1 DERIVADOS'                                 UNION ALL
        SELECT '4.1 PN'
    ),

    metrics_base AS (
        SELECT
        CASE
            WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') THEN '1.1 BANCOS'
            WHEN dim_1 IN ('A_Ata') THEN '1.2 CARTERA MAYORISTA'
            WHEN dim_1 IN ('A_Depositos cedidos') THEN '1.3 DEPOSITOS OTRAS EECC'
            WHEN dim_1 IN (
                    'A_ECORP_Otros_Activo Corto Plazo',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                    'A_EESP_Otros_Activo Corto Plazo',
                    'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                    'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                    'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                    'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                    'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                    'A_PIB_Otros_Descubiertos_Gtia Personal',
                    'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                    'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                    'A_POFI_Otros_Descubiertos_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                    'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                    'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_POFI_Otros_Prest. Origen_Gtia Aval',
                    'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                    'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                    'A_SSCC_Otros_Descubiertos_Gtia Personal',
                    'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                    'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
                ) THEN '1.4 INVERSION CREDITICIA'
            WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') THEN '1.5 RENTA FIJA'
            WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') THEN '1.6 ACTIVOS DUDOSOS'
            WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') THEN '1.7 ACTIVOS NO SENSIBLE'
            WHEN dim_1 IN ('P_Cuentas de otras EECC') THEN '2.1 DEPOSITOS OTRAS EECC'
            WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES'
            WHEN dim_1 IN (
            'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
            'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
            'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
            ) THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES'
            WHEN dim_1 IN ('P_Resto no sensible') THEN '2.4 PASIVOS NO SENSIBLES'
            WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') THEN '2.5 CARTERA MAYORISTA'
            WHEN dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo') THEN '3.1 DERIVADOS'
            WHEN dim_1 IN ('Patrimonio Neto') THEN '4.1 PN'
            ELSE 'OTROS'
        END AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200'   THEN market_value ELSE 0 END) AS Dwn200,
        SUM(CASE WHEN scenario = 'ShortRDwn' THEN market_value ELSE 0 END) AS Short_Rates_Down,
        SUM(CASE WHEN scenario = 'Flat'     THEN market_value ELSE 0 END) AS Flattening,
        SUM(CASE WHEN scenario = 'Base'     THEN market_value ELSE 0 END) AS Base,
        SUM(CASE WHEN scenario = 'Steep'    THEN market_value ELSE 0 END) AS Steepening,
        SUM(CASE WHEN scenario = 'ShortRUP' THEN market_value ELSE 0 END) AS Short_Rates_Up,
        SUM(CASE WHEN scenario = 'Up200'    THEN market_value ELSE 0 END) AS Up200
        FROM pro_pichincha_alquid_old.metric
        WHERE load_id IN (
        '{load_ids["cierre_base"]}',
        '{load_ids["cierre_up"]}',
        '{load_ids["cierre_dwn"]}'
        )
        GROUP BY 1
    ),

    base_con_todas AS (
        SELECT
            j.Jerarquia,
            b.Dwn200,
            b.Short_Rates_Down,
            b.Flattening,
            b.Base,
            b.Steepening,
            b.Short_Rates_Up,
            b.Up200
        FROM jerarquias j
        LEFT JOIN metrics_base b
        ON b.Jerarquia = j.Jerarquia
    )

    -- SELECT principal completo
    SELECT
        Jerarquia,
        Dwn200,
        Short_Rates_Down,
        Flattening,
        Base,
        Steepening,
        Short_Rates_Up,
        Up200
    FROM base_con_todas

    UNION ALL

    SELECT 'TOTAL GRUPO 1' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200'   THEN market_value ELSE 0 END) AS Dwn200,
        SUM(CASE WHEN scenario = 'ShortRDwn' THEN market_value ELSE 0 END) AS Short_Rates_Down,
        SUM(CASE WHEN scenario = 'Flat'     THEN market_value ELSE 0 END) AS Flattening,
        SUM(CASE WHEN scenario = 'Base'     THEN market_value ELSE 0 END) AS Base,
        SUM(CASE WHEN scenario = 'Steep'    THEN market_value ELSE 0 END) AS Steepening,
        SUM(CASE WHEN scenario = 'ShortRUP' THEN market_value ELSE 0 END) AS Short_Rates_Up,
        SUM(CASE WHEN scenario = 'Up200'    THEN market_value ELSE 0 END) AS Up200
    FROM pro_pichincha_alquid_old.metric
    WHERE load_id IN (
    '{load_ids["cierre_base"]}',
    '{load_ids["cierre_up"]}',
    '{load_ids["cierre_dwn"]}'
    )
    AND dim_1 in (
        'A_Cuentas en otras EECC', 'A_Tesoreria', 'A_Ata', 'A_Depositos cedidos',
        'A_ECORP_Otros_Activo Corto Plazo', 'A_EESP_Otros_Activo Corto Plazo', 'A_POFI_Otros_Ant. Nomina_Gtia Personal',
        'A_ECORP_Otros_Cred. Empresa_Gtia Aval', 'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
        'A_EESP_Otros_Cred. Empresa_Gtia Aval', 'A_EESP_Otros_Cred. Empresa_Gtia Personal',
        'A_SSCC_Otros_Cred. Empresa_Gtia Aval', 'A_SSCC_Otros_Cred. Empresa_Gtia Hipot.', 'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
        'A_PIB_Otros_Descubiertos_Gtia Personal', 'A_POFI_Otros_Descubiertos_Gtia Personal', 'A_SSCC_Otros_Descubiertos_Gtia Personal',
        'A_ECORP_Otros_Prest. Empresas_Gtia Aval', 'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.', 'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
        'A_ECORP_Otros_Prest. Empresas_Gtia Prenda', 'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',
        'A_EESP_Otros_Prest. Empresas_Gtia Aval', 'A_EESP_Otros_Prest. Empresas_Gtia Personal', 'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
        'A_EESP_Otros_Prest. Promotor_Gtia Hipot.', 'A_PIB_Otros_Prest. Consumo_Gtia Personal', 'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
        'A_PIB_Otros_Prest. Hipotecas_Gtia Personal', 'A_POFI_Otros_Prest. Consumo_Gtia Personal', 'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
        'A_POFI_Otros_Prest. Empresas_Gtia Aval', 'A_POFI_Otros_Prest. Empresas_Gtia Hipot.', 'A_POFI_Otros_Prest. Empresas_Gtia Personal',
        'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
        'A_POFI_Otros_Prest. Origen_Gtia Aval', 'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
        'A_SSCC_BKIA_Prest. Consumo_Gtia Personal', 'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
        'A_SSCC_Otros_Prest. Empresas_Gtia Aval', 'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
        'A_SSCC_Otros_Prest. Empresas_Gtia Personal', 'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
        'A_PIB_Otros_Tarjetas Credito_Gtia Personal', 'A_POFI_Otros_Tarjetas Credito_Gtia Personal',
        'A_SSCC_Otros_Tarjetas Credito_Gtia Personal', 'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal',
        'A_Bonos corporativos', 'A_Bonos soberanos', 'A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso',
        'A_Resto no sensible', 'A_Tesoreria_Admin.', 'A_EESP_Otros_Prest. Promotor_Gtia Aval', 'A_EESP_Otros_Prest. Promotor_Gtia Personal'
    )

    UNION ALL

    SELECT 'TOTAL GRUPO 2' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200'   THEN market_value ELSE 0 END) AS Dwn200,
        SUM(CASE WHEN scenario = 'ShortRDwn' THEN market_value ELSE 0 END) AS Short_Rates_Down,
        SUM(CASE WHEN scenario = 'Flat'     THEN market_value ELSE 0 END) AS Flattening,
        SUM(CASE WHEN scenario = 'Base'     THEN market_value ELSE 0 END) AS Base,
        SUM(CASE WHEN scenario = 'Steep'    THEN market_value ELSE 0 END) AS Steepening,
        SUM(CASE WHEN scenario = 'ShortRUP' THEN market_value ELSE 0 END) AS Short_Rates_Up,
        SUM(CASE WHEN scenario = 'Up200'    THEN market_value ELSE 0 END) AS Up200
    FROM pro_pichincha_alquid_old.metric
    WHERE load_id IN (
    '{load_ids["cierre_base"]}',
    '{load_ids["cierre_up"]}',
    '{load_ids["cierre_dwn"]}'
    )
    AND dim_1 in (
        'P_Cuentas de otras EECC', 'P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo',
        'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
        'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
        'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas',
        'P_Resto no sensible', 'P_Ata', 'P_Emisiones T2'
    )

    UNION ALL

    SELECT 'TOTAL GRUPO 3' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200'   THEN market_value ELSE 0 END) AS Dwn200,
        SUM(CASE WHEN scenario = 'ShortRDwn' THEN market_value ELSE 0 END) AS Short_Rates_Down,
        SUM(CASE WHEN scenario = 'Flat'     THEN market_value ELSE 0 END) AS Flattening,
        SUM(CASE WHEN scenario = 'Base'     THEN market_value ELSE 0 END) AS Base,
        SUM(CASE WHEN scenario = 'Steep'    THEN market_value ELSE 0 END) AS Steepening,
        SUM(CASE WHEN scenario = 'ShortRUP' THEN market_value ELSE 0 END) AS Short_Rates_Up,
        SUM(CASE WHEN scenario = 'Up200'    THEN market_value ELSE 0 END) AS Up200
    FROM pro_pichincha_alquid_old.metric
    WHERE load_id IN (
    '{load_ids["cierre_base"]}',
    '{load_ids["cierre_up"]}',
    '{load_ids["cierre_dwn"]}'
    )
    AND dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo')

    UNION ALL

    SELECT 'TOTAL GRUPO 4' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200'   THEN market_value ELSE 0 END) AS Dwn200,
        SUM(CASE WHEN scenario = 'ShortRDwn' THEN market_value ELSE 0 END) AS Short_Rates_Down,
        SUM(CASE WHEN scenario = 'Flat'     THEN market_value ELSE 0 END) AS Flattening,
        SUM(CASE WHEN scenario = 'Base'     THEN market_value ELSE 0 END) AS Base,
        SUM(CASE WHEN scenario = 'Steep'    THEN market_value ELSE 0 END) AS Steepening,
        SUM(CASE WHEN scenario = 'ShortRUP' THEN market_value ELSE 0 END) AS Short_Rates_Up,
        SUM(CASE WHEN scenario = 'Up200'    THEN market_value ELSE 0 END) AS Up200
    FROM pro_pichincha_alquid_old.metric
    WHERE load_id IN (
    '{load_ids["cierre_base"]}',
    '{load_ids["cierre_up"]}',
    '{load_ids["cierre_dwn"]}'
    )
    AND dim_1 IN ('Patrimonio Neto')

    ORDER BY CASE Jerarquia
        WHEN '1.1 BANCOS'                          THEN 1
        WHEN '1.2 CARTERA MAYORISTA'               THEN 2
        WHEN '1.3 DEPOSITOS OTRAS EECC'            THEN 3
        WHEN '1.4 INVERSION CREDITICIA'            THEN 4
        WHEN '1.5 RENTA FIJA'                      THEN 5
        WHEN '1.6 ACTIVOS DUDOSOS'                 THEN 6
        WHEN '1.7 ACTIVOS NO SENSIBLE'             THEN 7
        WHEN 'TOTAL GRUPO 1'                       THEN 8
        WHEN '2.1 DEPOSITOS OTRAS EECC'            THEN 9
        WHEN '2.2 DEPOSITOS A PLAZO DE CLIENTES'   THEN 10
        WHEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES' THEN 11
        WHEN '2.4 PASIVOS NO SENSIBLES'            THEN 13
        WHEN '2.5 CARTERA MAYORISTA'               THEN 12
        WHEN 'TOTAL GRUPO 2'                       THEN 14
        WHEN '3.1 DERIVADOS'                       THEN 15
        WHEN 'TOTAL GRUPO 3'                       THEN 16
        WHEN '4.1 PN'                              THEN 17
        WHEN 'TOTAL GRUPO 4'                       THEN 18
        ELSE 999
    END;
    """

    Informe_VE_EBA_02= f"""SELECT
                            dim_1 AS Jerarquia,
                        SUM(CASE WHEN SCENARIO = 'Base' THEN effective_duration * market_value ELSE 0 END) /
                            NULLIF(SUM(CASE WHEN SCENARIO = 'Base' THEN market_value ELSE 0 END), 0)/100 AS effective_duration_PONDERADA,
                            SUM(CASE WHEN scenario = 'Dwn200' THEN market_value ELSE 0 END) AS Dwn200,
                            SUM(CASE WHEN scenario = 'ShortRDwn' THEN market_value ELSE 0 END) AS Short_Rates_Down,
                            SUM(CASE WHEN scenario = 'Flat' THEN market_value ELSE 0 END) AS Flattening,
                            SUM(CASE WHEN scenario = 'Base' THEN market_value ELSE 0 END) AS Base,
                            SUM(CASE WHEN scenario = 'Steep' THEN market_value ELSE 0 END) AS Steepening,
                            SUM(CASE WHEN scenario = 'ShortRUP' THEN market_value ELSE 0 END) AS Short_Rates_Up,
                            SUM(CASE WHEN scenario = 'Up200' THEN market_value ELSE 0 END) AS Up200
                        FROM pro_pichincha_alquid_old.metric
                            WHERE load_id IN (
                            '{load_ids['cierre_base']}',
                            '{load_ids['cierre_up']}',
                            '{load_ids['cierre_dwn']}'
                            )
                        AND dim_1 IN (
                            'A_Bonos corporativos',
                            'A_Bonos soberanos',
                            'IRS_pago',
                            'IRS_recibo',
                            'FXSWAP_pago',
                            'FXSWAP_recibo'
                        )
                        GROUP BY dim_1
                        ORDER BY case
                            when dim_1 = 'A_Bonos corporativos' then 1
                            when dim_1 = 'A_Bonos soberanos' then 2
                            when dim_1 = 'IRS_pago' then 3
                            when dim_1 = 'IRS_recibo' then 4
                            when dim_1 = 'FXSWAP_pagos' then 5
                            when dim_1 = 'FXSWAP_recibo' then 6
                        end; """

    Informe_VE_EBA_03= f"""SELECT
                            dim_1 AS Jerarquia,
                            SUM(CASE WHEN SCENARIO = 'Base' THEN macaulay_duration * market_value ELSE 0 END) /
                                NULLIF(SUM(CASE WHEN SCENARIO = 'Base' THEN market_value ELSE 0 END), 0) AS DURACION_MACAULAY_PONDERADA,
                            SUM(CASE WHEN scenario = 'Dwn200' THEN market_value ELSE 0 END) AS Dwn200,
                            SUM(CASE WHEN scenario = 'ShortRDwn' THEN market_value ELSE 0 END) AS Short_Rates_Down,
                            SUM(CASE WHEN scenario = 'Flat' THEN market_value ELSE 0 END) AS Flattening,
                            SUM(CASE WHEN scenario = 'Base' THEN market_value ELSE 0 END) AS Base,
                            SUM(CASE WHEN scenario = 'Steep' THEN market_value ELSE 0 END) AS Steepening,
                            SUM(CASE WHEN scenario = 'ShortRUP' THEN market_value ELSE 0 END) AS Short_Rates_Up,
                            SUM(CASE WHEN scenario = 'Up200' THEN market_value ELSE 0 END) AS Up200
                            FROM pro_pichincha_alquid_old.metric
                            WHERE load_id IN (
                            '{load_ids['cierre_base']}',
                            '{load_ids['cierre_up']}',
                            '{load_ids['cierre_dwn']}'
                            )
                            AND dim_1 IN (
                            'A_Bonos corporativos',
                            'A_Bonos soberanos',
                            'IRS_pago',
                            'IRS_recibo',
                                'FXSWAP_pago',
                                'FXSWAP_recibo')
                            GROUP BY dim_1
                            ORDER BY case
                                when dim_1 = 'A_Bonos corporativos' then 1
                                when dim_1 = 'A_Bonos soberanos' then 2
                                when dim_1 = 'IRS_pago' then 3
                                when dim_1 = 'IRS_recibo' then 4
                                when dim_1 = 'FXSWAP_pagos' then 5
                                when dim_1 = 'FXSWAP_recibo' then 6
                            end; """

    Informe_VE_EBA_04= f"""
WITH jerarquias AS (
    SELECT '1.1 BANCOS' AS jerarquia, 1 AS orden UNION ALL
    SELECT '1.2 CARTERA MAYORISTA', 2 UNION ALL
    SELECT '1.3 DEPOSITOS OTRAS EECC', 3 UNION ALL
    SELECT '1.4 INVERSION CREDITICIA', 4 UNION ALL
    SELECT '1.6 ACTIVOS DUDOSOS', 6 UNION ALL
    SELECT '1.7 ACTIVOS NO SENSIBLE', 7 UNION ALL
    SELECT '2.1 DEPOSITOS OTRAS EECC', 8 UNION ALL
    SELECT '2.2 DEPOSITOS A PLAZO DE CLIENTES', 9 UNION ALL
    SELECT '2.5 CARTERA MAYORISTA', 10 UNION ALL
    SELECT '2.4 PASIVOS NO SENSIBLES', 11 UNION ALL
    SELECT '3.1 DERIVADOS', 12 UNION ALL
    SELECT '4.1 PN', 13
),
agg AS (
SELECT
    CASE
        WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') THEN '1.1 BANCOS'
        WHEN dim_1 IN ('A_Ata') THEN '1.2 CARTERA MAYORISTA'
        WHEN dim_1 IN ('A_Depositos cedidos') THEN '1.3 DEPOSITOS OTRAS EECC'
        WHEN dim_1 IN (
                    'A_ECORP_Otros_Activo Corto Plazo',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                    'A_EESP_Otros_Activo Corto Plazo',
                    'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                    'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                    'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                    'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                    'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                    'A_PIB_Otros_Descubiertos_Gtia Personal',
                    'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                    'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                    'A_POFI_Otros_Descubiertos_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                    'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                    'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_POFI_Otros_Prest. Origen_Gtia Aval',
                    'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                    'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                    'A_SSCC_Otros_Descubiertos_Gtia Personal',
                    'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                    'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
                ) THEN '1.4 INVERSION CREDITICIA'
        WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') THEN '1.5 RENTA FIJA'
        WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') THEN '1.6 ACTIVOS DUDOSOS'
        WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') THEN '1.7 ACTIVOS NO SENSIBLE'
        WHEN dim_1 IN ('P_Cuentas de otras EECC') THEN '2.1 DEPOSITOS OTRAS EECC'
        WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES'
        WHEN dim_1 IN (
        'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
        'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
        'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
        ) THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES'
        WHEN dim_1 IN ('P_Resto no sensible') THEN '2.4 PASIVOS NO SENSIBLES'
        WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') THEN '2.5 CARTERA MAYORISTA'
        WHEN dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo') THEN '3.1 DERIVADOS'
        WHEN dim_1 IN ('Patrimonio Neto') THEN '4.1 PN'
        ELSE 'OTROS'
    END AS Jerarquia,
    SUM(CASE WHEN SCENARIO = 'Base' THEN ABS(effective_duration * market_value) ELSE 0 END) /
        NULLIF(SUM(CASE WHEN SCENARIO = 'Base' THEN ABS(market_value) ELSE 0 END), 0)/100 AS effective_duration_PONDERADA,
    SUM(CASE WHEN SCENARIO = 'Base' THEN ABS(macaulay_duration * market_value) ELSE 0 END) /
        NULLIF(SUM(CASE WHEN SCENARIO = 'Base' THEN ABS(market_value) ELSE 0 END), 0) AS DURACION_MACAULAY_PONDERADA
    FROM pro_pichincha_alquid_old.metric
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )
    GROUP BY 1
)


SELECT
    j.jerarquia,
    ROUND(COALESCE(a.effective_duration_ponderada, 0), 4) AS effective_duration_ponderada,
    ROUND(COALESCE(a.duracion_macaulay_ponderada, 0), 4) AS duracion_macaulay_ponderada
FROM jerarquias j
LEFT JOIN agg a
    ON j.jerarquia = a.jerarquia
ORDER BY j.orden;

"""

    # Se define nuestro data frame con los datos extrídos de la base de datos
    df_VE_EBA_01 = pd.read_sql(Informe_VE_EBA_01, conn)
    df_VE_EBA_02 = pd.read_sql(Informe_VE_EBA_02, conn)
    df_VE_EBA_03 = pd.read_sql(Informe_VE_EBA_03, conn)
    df_VE_EBA_04= pd.read_sql(Informe_VE_EBA_04, conn)

    # Nombre de la hoja donde quieres pegar los datos. Ahora se construye dinámicamente.
    nombre_hoja_VE_EBA = f"Sensibilidades VE EBA - {mes_actual_nombre}"


    # Acceder a la hoja
    ws_VE_EBA = wb[nombre_hoja_VE_EBA]


    # Se definen las coordenadas donde queremos que se peguen los datos
    start_row_VE_EBA_01 = 4
    start_col_VE_EBA_01 = 33

    start_row_VE_EBA_02 = 33
    start_col_VE_EBA_02 = 32


    start_row_VE_EBA_03 = 44
    start_col_VE_EBA_03 = 32

    start_row_VE_EBA_04 = 4
    start_col_VE_EBA_04 = 47


    # Límite de filas donde se pegarán los datos (de la 6 a la 22)
    end_row_VE_EBA_01 = 22
    end_row_VE_EBA_02 = 38
    end_row_VE_EBA_03 = 49
    end_row_VE_EBA_04 = 20

    # Columnas de Excel donde NO queremos escribir
    skip_excel_columns_EBA= [36, 37, 39, 40]  # Números de columnas de Excel a evitar

    # Filas a evitar
    filas_a_saltar_EBA_04 = {8, 11, 14, 17, 19}
    fila_actual_VE_EBA_04 = start_row_VE_EBA_04

    # Query de pegado para el informe 1
    for i, row in enumerate(df_VE_EBA_01.values, start=start_row_VE_EBA_01):
        if i > end_row_VE_EBA_01:
            break  # Limita hasta fila 22

        excel_col_VE_EBA_01 = start_col_VE_EBA_01  # Puntero en Excel
        for value_EBA_01 in row[1:]:
            while excel_col_VE_EBA_01 in skip_excel_columns_EBA:
                excel_col_VE_EBA_01 += 1  # Saltamos columnas prohibidas

            ws_VE_EBA.cell(row=i, column=excel_col_VE_EBA_01).value = value_EBA_01
            excel_col_VE_EBA_01 += 1  # Avanzamos a la siguiente columna de Excel

    # Query de pegado para el informe 2
    for i, row in enumerate(df_VE_EBA_02.values, start=start_row_VE_EBA_02):
        if i > end_row_VE_EBA_02:
            break  # Limita hasta fila 22

        excel_col_VE_EBA_02 = start_col_VE_EBA_02  # Puntero en Excel
        for value_EBA_02 in row[1:]:
            while excel_col_VE_EBA_02 in skip_excel_columns_EBA:
                excel_col_VE_EBA_02 += 1  # Saltamos columnas prohibidas

            ws_VE_EBA.cell(row=i, column=excel_col_VE_EBA_02).value = value_EBA_02
            excel_col_VE_EBA_02 += 1  # Avanzamos a la siguiente columna de Excel

    # Query de pegado para el informe 3
    for i, row in enumerate(df_VE_EBA_03.values, start=start_row_VE_EBA_03):
        if i > end_row_VE_EBA_03:
            break  # Limita hasta fila 22

        excel_col_VE_EBA_03 = start_col_VE_EBA_03  # Puntero en Excel
        for value_EBA_03 in row[1:]:
            while excel_col_VE_EBA_03 in skip_excel_columns_EBA:
                excel_col_VE_EBA_03 += 1  # Saltamos columnas prohibidas

            ws_VE_EBA.cell(row=i, column=excel_col_VE_EBA_03).value = value_EBA_03
            excel_col_VE_EBA_03 += 1  # Avanzamos a la siguiente columna de Excel

    # Query de pegado para el informe 4
    for row in df_VE_EBA_04.values:
        # Saltar las filas especificadas
        while fila_actual_VE_EBA_04 in filas_a_saltar_EBA_04:
            fila_actual_VE_EBA_04 += 1

        if fila_actual_VE_EBA_04 > end_row_VE_EBA_04:
            break

        for j, value_EBA_04 in enumerate(row[1:], start=start_col_VE_EBA_04):
            ws_VE_EBA.cell(row=fila_actual_VE_EBA_04, column=j).value = value_EBA_04

        fila_actual_VE_EBA_04 += 1  # Mover a la siguiente fila en la hoja





    ###################################### PROCESO DE PEGADO DE INFORMES DE MARGEN FINANCIERO MTO #####################################
    Informe_MF_MTO= f"""
        WITH jerarquias AS (
        -- GRUPO 1
        SELECT 1  AS orden, '1.1 BANCOS'                         AS Jerarquia UNION ALL
        SELECT 2,             '1.2 CARTERA MAYORISTA'                         UNION ALL
        SELECT 3,             '1.3 DEPOSITOS OTRAS EECC'                      UNION ALL
        SELECT 4,             '1.4 INVERSION CREDITICIA'                      UNION ALL
        SELECT 5,             '1.5 RENTA FIJA'                                UNION ALL
        SELECT 6,             '1.6 ACTIVOS DUDOSOS'                           UNION ALL
        SELECT 7,             '1.7 ACTIVOS NO SENSIBLE'                       UNION ALL
        -- Línea en blanco (se pinta como ' ' y quedará con NULLs)
        SELECT 9,             ' '                                             UNION ALL
        -- GRUPO 2
        SELECT 10,            '2.1 DEPOSITOS OTRAS EECC'                      UNION ALL
        SELECT 11,            '2.2 DEPOSITOS A PLAZO DE CLIENTES'             UNION ALL
        SELECT 12,            '2.3 DEPOSITOS A LA VISTA DE CLIENTES'          UNION ALL
        SELECT 13,            '2.4 PASIVOS NO SENSIBLES'                      UNION ALL
        SELECT 14,            '2.5 CARTERA MAYORISTA'                         UNION ALL
        -- GRUPO 3
        SELECT 16,            '3.1 DERIVADOS'                                 UNION ALL
        -- GRUPO 4
        SELECT 18,            '4.1 PN'
    ),

    metrics_agg AS (
        SELECT
            CASE
            WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') THEN '1.1 BANCOS'
            WHEN dim_1 IN ('A_Ata') THEN '1.2 CARTERA MAYORISTA'
            WHEN dim_1 IN ('A_Depositos cedidos') THEN '1.3 DEPOSITOS OTRAS EECC'
            WHEN dim_1 IN (
                    'A_ECORP_Otros_Activo Corto Plazo',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                    'A_EESP_Otros_Activo Corto Plazo',
                    'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                    'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                    'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                    'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                    'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                    'A_PIB_Otros_Descubiertos_Gtia Personal',
                    'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                    'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                    'A_POFI_Otros_Descubiertos_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                    'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                    'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_POFI_Otros_Prest. Origen_Gtia Aval',
                    'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                    'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                    'A_SSCC_Otros_Descubiertos_Gtia Personal',
                    'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                    'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
                ) THEN '1.4 INVERSION CREDITICIA'
            WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') THEN '1.5 RENTA FIJA'
            WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') THEN '1.6 ACTIVOS DUDOSOS'
            WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') THEN '1.7 ACTIVOS NO SENSIBLE'
            WHEN dim_1 IN ('P_Cuentas de otras EECC') THEN '2.1 DEPOSITOS OTRAS EECC'
            WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES'
            WHEN dim_1 IN (
                'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
                'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
            ) THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES'
            WHEN dim_1 IN ('P_Resto no sensible') THEN '2.4 PASIVOS NO SENSIBLES'
            WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') THEN '2.5 CARTERA MAYORISTA'
            WHEN dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo') THEN '3.1 DERIVADOS'
            WHEN dim_1 IN ('Patrimonio Neto') THEN '4.1 PN'
            ELSE 'OTROS'
            END AS Jerarquia,
            SUM(CASE WHEN scenario = 'Dwn200' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_stck,
            SUM(CASE WHEN scenario = 'Dwn200' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_Nb,
            SUM(CASE WHEN scenario = 'Dwn100' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_stck,
            SUM(CASE WHEN scenario = 'Dwn100' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_Nb,
            SUM(CASE WHEN scenario = 'Dwn75'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_stck,
            SUM(CASE WHEN scenario = 'Dwn75'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_Nb,
            SUM(CASE WHEN scenario = 'Dwn50'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_stck,
            SUM(CASE WHEN scenario = 'Dwn50'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_Nb,
            SUM(CASE WHEN scenario = 'Dwn25'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_stck,
            SUM(CASE WHEN scenario = 'Dwn25'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_Nb,
            SUM(CASE WHEN scenario = 'Base'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_stck,
            SUM(CASE WHEN scenario = 'Base'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_Nb,
            SUM(CASE WHEN scenario = 'Up25'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_stck,
            SUM(CASE WHEN scenario = 'Up25'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_Nb,
            SUM(CASE WHEN scenario = 'Up50'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_stck,
            SUM(CASE WHEN scenario = 'Up50'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_Nb,
            SUM(CASE WHEN scenario = 'Up75'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_stck,
            SUM(CASE WHEN scenario = 'Up75'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_Nb,
            SUM(CASE WHEN scenario = 'Up100'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_stck,
            SUM(CASE WHEN scenario = 'Up100'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_Nb,
            SUM(CASE WHEN scenario = 'Up200'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_stck,
            SUM(CASE WHEN scenario = 'Up200'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_Nb
        FROM pro_pichincha_alquid_old."result"
        WHERE load_id IN (
        '{load_ids['cierre_base']}',
        '{load_ids['cierre_up']}',
        '{load_ids['cierre_dwn']}'
        )
        GROUP BY 1
    )

    -- 1) Jerarquías fijas (1.x, 2.x, 3.1, 4.1 y la línea en blanco) con LEFT JOIN
    SELECT
        j.Jerarquia,
        a.Dwn200_stck,
        a.Dwn200_Nb,
        a.Dwn100_stck,
        a.Dwn100_Nb,
        a.Dwn75_stck,
        a.Dwn75_Nb,
        a.Dwn50_stck,
        a.Dwn50_Nb,
        a.Dwn25_stck,
        a.Dwn25_Nb,
        a.Base_stck,
        a.Base_Nb,
        a.Up25_stck,
        a.Up25_Nb,
        a.Up50_stck,
        a.Up50_Nb,
        a.Up75_stck,
        a.Up75_Nb,
        a.Up100_stck,
        a.Up100_Nb,
        a.Up200_stck,
        a.Up200_Nb
    FROM jerarquias j
    LEFT JOIN metrics_agg a
        ON a.Jerarquia = j.Jerarquia

    UNION ALL

    -- 2) TOTAL GRUPO 1
    SELECT 'TOTAL GRUPO 1' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_stck,
        SUM(CASE WHEN scenario = 'Dwn200' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_Nb,
        SUM(CASE WHEN scenario = 'Dwn100' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_stck,
        SUM(CASE WHEN scenario = 'Dwn100' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_Nb,
        SUM(CASE WHEN scenario = 'Dwn75'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_stck,
        SUM(CASE WHEN scenario = 'Dwn75'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_Nb,
        SUM(CASE WHEN scenario = 'Dwn50'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_stck,
        SUM(CASE WHEN scenario = 'Dwn50'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_Nb,
        SUM(CASE WHEN scenario = 'Dwn25'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_stck,
        SUM(CASE WHEN scenario = 'Dwn25'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_Nb,
        SUM(CASE WHEN scenario = 'Base'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_stck,
        SUM(CASE WHEN scenario = 'Base'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_Nb,
        SUM(CASE WHEN scenario = 'Up25'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_stck,
        SUM(CASE WHEN scenario = 'Up25'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_Nb,
        SUM(CASE WHEN scenario = 'Up50'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_stck,
        SUM(CASE WHEN scenario = 'Up50'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_Nb,
        SUM(CASE WHEN scenario = 'Up75'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_stck,
        SUM(CASE WHEN scenario = 'Up75'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_Nb,
        SUM(CASE WHEN scenario = 'Up100'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_stck,
        SUM(CASE WHEN scenario = 'Up100'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_Nb,
        SUM(CASE WHEN scenario = 'Up200'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_stck,
        SUM(CASE WHEN scenario = 'Up200'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_Nb
    FROM pro_pichincha_alquid_old."result"
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )
    AND dim_1 IN (
    'A_Cuentas en otras EECC', 'A_Tesoreria', 'A_Ata', 'A_Depositos cedidos',
    'A_ECORP_Otros_Activo Corto Plazo', 'A_EESP_Otros_Activo Corto Plazo', 'A_POFI_Otros_Ant. Nomina_Gtia Personal',
    'A_ECORP_Otros_Cred. Empresa_Gtia Aval', 'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
    'A_EESP_Otros_Cred. Empresa_Gtia Aval', 'A_EESP_Otros_Cred. Empresa_Gtia Personal',
    'A_SSCC_Otros_Cred. Empresa_Gtia Aval', 'A_SSCC_Otros_Cred. Empresa_Gtia Hipot.', 'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
    'A_PIB_Otros_Descubiertos_Gtia Personal', 'A_POFI_Otros_Descubiertos_Gtia Personal', 'A_SSCC_Otros_Descubiertos_Gtia Personal',
    'A_ECORP_Otros_Prest. Empresas_Gtia Aval', 'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.', 'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
    'A_ECORP_Otros_Prest. Empresas_Gtia Prenda', 'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',
    'A_EESP_Otros_Prest. Empresas_Gtia Aval', 'A_EESP_Otros_Prest. Empresas_Gtia Personal', 'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
    'A_EESP_Otros_Prest. Promotor_Gtia Hipot.', 'A_PIB_Otros_Prest. Consumo_Gtia Personal', 'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
    'A_PIB_Otros_Prest. Hipotecas_Gtia Personal', 'A_POFI_Otros_Prest. Consumo_Gtia Personal', 'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
    'A_POFI_Otros_Prest. Empresas_Gtia Aval', 'A_POFI_Otros_Prest. Empresas_Gtia Hipot.', 'A_POFI_Otros_Prest. Empresas_Gtia Personal',
    'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
    'A_POFI_Otros_Prest. Origen_Gtia Aval', 'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
    'A_SSCC_BKIA_Prest. Consumo_Gtia Personal', 'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
    'A_SSCC_Otros_Prest. Empresas_Gtia Aval', 'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
    'A_SSCC_Otros_Prest. Empresas_Gtia Personal', 'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
    'A_PIB_Otros_Tarjetas Credito_Gtia Personal', 'A_POFI_Otros_Tarjetas Credito_Gtia Personal',
    'A_SSCC_Otros_Tarjetas Credito_Gtia Personal', 'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal',
    'A_Bonos corporativos', 'A_Bonos soberanos', 'A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso',
    'A_Resto no sensible', 'A_Tesoreria_Admin.', 'A_EESP_Otros_Prest. Promotor_Gtia Aval', 'A_EESP_Otros_Prest. Promotor_Gtia Personal'
    )

    UNION ALL
    -- 3) TOTAL GRUPO 2
    SELECT 'TOTAL GRUPO 2' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_stck,
        SUM(CASE WHEN scenario = 'Dwn200' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_Nb,
        SUM(CASE WHEN scenario = 'Dwn100' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_stck,
        SUM(CASE WHEN scenario = 'Dwn100' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_Nb,
        SUM(CASE WHEN scenario = 'Dwn75'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_stck,
        SUM(CASE WHEN scenario = 'Dwn75'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_Nb,
        SUM(CASE WHEN scenario = 'Dwn50'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_stck,
        SUM(CASE WHEN scenario = 'Dwn50'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_Nb,
        SUM(CASE WHEN scenario = 'Dwn25'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_stck,
        SUM(CASE WHEN scenario = 'Dwn25'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_Nb,
        SUM(CASE WHEN scenario = 'Base'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_stck,
        SUM(CASE WHEN scenario = 'Base'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_Nb,
        SUM(CASE WHEN scenario = 'Up25'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_stck,
        SUM(CASE WHEN scenario = 'Up25'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_Nb,
        SUM(CASE WHEN scenario = 'Up50'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_stck,
        SUM(CASE WHEN scenario = 'Up50'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_Nb,
        SUM(CASE WHEN scenario = 'Up75'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_stck,
        SUM(CASE WHEN scenario = 'Up75'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_Nb,
        SUM(CASE WHEN scenario = 'Up100'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_stck,
        SUM(CASE WHEN scenario = 'Up100'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_Nb,
        SUM(CASE WHEN scenario = 'Up200'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_stck,
        SUM(CASE WHEN scenario = 'Up200'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_Nb
    FROM pro_pichincha_alquid_old."result"
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )
    AND dim_1 IN (
    'P_Cuentas de otras EECC', 'P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo',
    'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
    'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
    'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas',
    'P_Resto no sensible', 'P_Ata', 'P_Emisiones T2'
    )

    UNION ALL
    -- 4) TOTAL GRUPO 3
    SELECT 'TOTAL GRUPO 3' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_stck,
        SUM(CASE WHEN scenario = 'Dwn200' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_Nb,
        SUM(CASE WHEN scenario = 'Dwn100' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_stck,
        SUM(CASE WHEN scenario = 'Dwn100' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_Nb,
        SUM(CASE WHEN scenario = 'Dwn75'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_stck,
        SUM(CASE WHEN scenario = 'Dwn75'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_Nb,
        SUM(CASE WHEN scenario = 'Dwn50'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_stck,
        SUM(CASE WHEN scenario = 'Dwn50'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_Nb,
        SUM(CASE WHEN scenario = 'Dwn25'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_stck,
        SUM(CASE WHEN scenario = 'Dwn25'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_Nb,
        SUM(CASE WHEN scenario = 'Base'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_stck,
        SUM(CASE WHEN scenario = 'Base'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_Nb,
        SUM(CASE WHEN scenario = 'Up25'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_stck,
        SUM(CASE WHEN scenario = 'Up25'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_Nb,
        SUM(CASE WHEN scenario = 'Up50'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_stck,
        SUM(CASE WHEN scenario = 'Up50'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_Nb,
        SUM(CASE WHEN scenario = 'Up75'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_stck,
        SUM(CASE WHEN scenario = 'Up75'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_Nb,
        SUM(CASE WHEN scenario = 'Up100'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_stck,
        SUM(CASE WHEN scenario = 'Up100'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_Nb,
        SUM(CASE WHEN scenario = 'Up200'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_stck,
        SUM(CASE WHEN scenario = 'Up200'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_Nb
    FROM pro_pichincha_alquid_old."result"
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )
    AND dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo')

    UNION ALL
    -- 5) TOTAL GRUPO 4
    SELECT 'TOTAL GRUPO 4' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_stck,
        SUM(CASE WHEN scenario = 'Dwn200' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_Nb,
        SUM(CASE WHEN scenario = 'Dwn100' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_stck,
        SUM(CASE WHEN scenario = 'Dwn100' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_Nb,
        SUM(CASE WHEN scenario = 'Dwn75'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_stck,
        SUM(CASE WHEN scenario = 'Dwn75'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_Nb,
        SUM(CASE WHEN scenario = 'Dwn50'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_stck,
        SUM(CASE WHEN scenario = 'Dwn50'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_Nb,
        SUM(CASE WHEN scenario = 'Dwn25'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_stck,
        SUM(CASE WHEN scenario = 'Dwn25'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_Nb,
        SUM(CASE WHEN scenario = 'Base'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_stck,
        SUM(CASE WHEN scenario = 'Base'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_Nb,
        SUM(CASE WHEN scenario = 'Up25'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_stck,
        SUM(CASE WHEN scenario = 'Up25'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_Nb,
        SUM(CASE WHEN scenario = 'Up50'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_stck,
        SUM(CASE WHEN scenario = 'Up50'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_Nb,
        SUM(CASE WHEN scenario = 'Up75'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_stck,
        SUM(CASE WHEN scenario = 'Up75'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_Nb,
        SUM(CASE WHEN scenario = 'Up100'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_stck,
        SUM(CASE WHEN scenario = 'Up100'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_Nb,
        SUM(CASE WHEN scenario = 'Up200'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_stck,
        SUM(CASE WHEN scenario = 'Up200'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_Nb
    FROM pro_pichincha_alquid_old."result"
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )
    AND dim_1 IN ('Patrimonio Neto')

    ORDER BY CASE Jerarquia
    WHEN '1.1 BANCOS'                         THEN 1
    WHEN '1.2 CARTERA MAYORISTA'             THEN 2
    WHEN '1.3 DEPOSITOS OTRAS EECC'          THEN 3
    WHEN '1.4 INVERSION CREDITICIA'          THEN 4
    WHEN '1.5 RENTA FIJA'                    THEN 5
    WHEN '1.6 ACTIVOS DUDOSOS'               THEN 6
    WHEN '1.7 ACTIVOS NO SENSIBLE'           THEN 7
    WHEN 'TOTAL GRUPO 1'                     THEN 8
    WHEN ' '                                 THEN 9
    WHEN '2.1 DEPOSITOS OTRAS EECC'          THEN 10
    WHEN '2.2 DEPOSITOS A PLAZO DE CLIENTES' THEN 11
    WHEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES' THEN 12
    WHEN '2.4 PASIVOS NO SENSIBLES'          THEN 13
    WHEN '2.5 CARTERA MAYORISTA'             THEN 14
    WHEN 'TOTAL GRUPO 2'                     THEN 15
    WHEN '3.1 DERIVADOS'                     THEN 16
    WHEN 'TOTAL GRUPO 3'                     THEN 17
    WHEN '4.1 PN'                            THEN 18
    WHEN 'TOTAL GRUPO 4'                     THEN 19
    ELSE 999
    END;
    """

    # Se define nuestro data frame con los datos extrídos de la base de datos
    df_MF_MTO = pd.read_sql(Informe_MF_MTO, conn)

    # Nombre de la hoja donde quieres pegar los datos. Ahora se construye dinámicamente.
    nombre_hoja_MF_MTO = f"MF - Sensibilidades - Mto - {mes_actual_nombre}"

    # Acceder a la hoja
    ws_MF_MTO = wb[nombre_hoja_MF_MTO]


    # Se definen las coordenadas donde queremos que se peguen los datos
    start_row_MF_MTO = 6
    start_col_MF_MTO = 52

    # Límite de filas donde se pegarán los datos (de la 6 a la 24)
    end_row_MF_MTO= 24

    # Escribir los datos en las coordenadas anteriores
    for i, row in enumerate(df_MF_MTO.values, start=start_row_MF_MTO):
        if i > end_row_MF_MTO:
            break  # Detener la escritura si superamos la fila 24
        for j, value_MF_MTO in enumerate(row[1:], start=start_col_MF_MTO):
            ws_MF_MTO.cell(row=i, column=j).value = value_MF_MTO








    ###################################### PROCESO DE PEGADO DE INFORMES DE MARGEN FINANCIERO EBA #####################################
    Informe_MF_EBA= f"""
        WITH metrics_agg AS (
        SELECT
        CASE
            WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') THEN '1.1 BANCOS'
            WHEN dim_1 IN ('A_Ata') THEN '1.2 CARTERA MAYORISTA'
            WHEN dim_1 IN ('A_Depositos cedidos') THEN '1.3 DEPOSITOS OTRAS EECC'
            WHEN dim_1 IN (
                    'A_ECORP_Otros_Activo Corto Plazo',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                    'A_EESP_Otros_Activo Corto Plazo',
                    'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                    'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                    'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                    'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                    'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                    'A_PIB_Otros_Descubiertos_Gtia Personal',
                    'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                    'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                    'A_POFI_Otros_Descubiertos_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                    'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                    'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_POFI_Otros_Prest. Origen_Gtia Aval',
                    'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                    'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                    'A_SSCC_Otros_Descubiertos_Gtia Personal',
                    'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                    'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
                ) THEN '1.4 INVERSION CREDITICIA'
            WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') THEN '1.5 RENTA FIJA'
            WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') THEN '1.6 ACTIVOS DUDOSOS'
            WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') THEN '1.7 ACTIVOS NO SENSIBLE'
            WHEN dim_1 IN ('P_Cuentas de otras EECC') THEN '2.1 DEPOSITOS OTRAS EECC'
            WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES'
            WHEN dim_1 IN (
            'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
            'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
            'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
            ) THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES'
            WHEN dim_1 IN ('P_Resto no sensible') THEN '2.4 PASIVOS NO SENSIBLES'
            WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') THEN '2.5 CARTERA MAYORISTA'
            WHEN dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo') THEN '3.1 DERIVADOS'
            WHEN dim_1 IN ('Patrimonio Neto') THEN '4.1 PN'
            ELSE 'OTROS'
        END AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_stck,
        SUM(CASE WHEN scenario = 'Dwn200'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_Nb,
        SUM(CASE WHEN scenario = 'ShortRDwn' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS ShortRDwn_stck,
        SUM(CASE WHEN scenario = 'ShortRDwn' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS ShortRDwn_Nb,
        SUM(CASE WHEN scenario = 'Steep'    AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Steep_stck,
        SUM(CASE WHEN scenario = 'Steep'    AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Steep_Nb,
        SUM(CASE WHEN scenario = 'Base'     AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_stck,
        SUM(CASE WHEN scenario = 'Base'     AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_Nb,
        SUM(CASE WHEN scenario = 'Flat'     AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Flat_stck,
        SUM(CASE WHEN scenario = 'Flat'     AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Flat_Nb,
        SUM(CASE WHEN scenario = 'ShortRUP' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS ShortRUP_stck,
        SUM(CASE WHEN scenario = 'ShortRUP' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS ShortRUP_Nb,
        SUM(CASE WHEN scenario = 'Up200'    AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_stck,
        SUM(CASE WHEN scenario = 'Up200'    AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_Nb
        FROM pro_pichincha_alquid_old."result"
        WHERE load_id IN (
        '{load_ids['cierre_base']}',
        '{load_ids['cierre_up']}',
        '{load_ids['cierre_dwn']}'
        )
        GROUP BY 1
    ),

    -- lista fija de jerarquías que quieres ver siempre
    jerarquias AS (
        SELECT '1.1 BANCOS'                   AS Jerarquia UNION ALL
        SELECT '1.2 CARTERA MAYORISTA'                    UNION ALL
        SELECT '1.3 DEPOSITOS OTRAS EECC'                 UNION ALL
        SELECT '1.4 INVERSION CREDITICIA'                 UNION ALL
        SELECT '1.5 RENTA FIJA'                           UNION ALL
        SELECT '1.6 ACTIVOS DUDOSOS'                      UNION ALL
        SELECT '1.7 ACTIVOS NO SENSIBLE'                  UNION ALL
        SELECT '2.1 DEPOSITOS OTRAS EECC'                 UNION ALL
        SELECT '2.2 DEPOSITOS A PLAZO DE CLIENTES'        UNION ALL
        SELECT '2.3 DEPOSITOS A LA VISTA DE CLIENTES'     UNION ALL
        SELECT '2.4 PASIVOS NO SENSIBLES'                 UNION ALL
        SELECT '2.5 CARTERA MAYORISTA'                    UNION ALL
        SELECT '3.1 DERIVADOS'                            UNION ALL
        SELECT '4.1 PN'
    )

    SELECT *
    FROM (
        -- línea en blanco (después del bloque 1.x en el orden final)
        SELECT
            '  ' AS Jerarquia,
            NULL AS Dwn200_stck, NULL AS Dwn200_Nb,
            NULL AS ShortRDwn_stck, NULL AS ShortRDwn_Nb,
            NULL AS Steep_stck, NULL AS Steep_Nb,
            NULL AS Base_stck, NULL AS Base_Nb,
            NULL AS Flat_stck, NULL AS Flat_Nb,
            NULL AS ShortRUP_stck, NULL AS ShortRUP_Nb,
            NULL AS Up200_stck, NULL AS Up200_Nb
        UNION ALL
        -- otra línea en blanco (entre bloques de pasivo, como tenías)
        SELECT
            ' ' AS Jerarquia,
            NULL AS Dwn200_stck, NULL AS Dwn200_Nb,
            NULL AS ShortRDwn_stck, NULL AS ShortRDwn_Nb,
            NULL AS Steep_stck, NULL AS Steep_Nb,
            NULL AS Base_stck, NULL AS Base_Nb,
            NULL AS Flat_stck, NULL AS Flat_Nb,
            NULL AS ShortRUP_stck, NULL AS ShortRUP_Nb,
            NULL AS Up200_stck, NULL AS Up200_Nb
        UNION ALL
        -- todas las jerarquías 1.x, 2.x, 3.1, 4.1 con LEFT JOIN para que salgan siempre
        SELECT
            j.Jerarquia,
            m.Dwn200_stck,
            m.Dwn200_Nb,
            m.ShortRDwn_stck,
            m.ShortRDwn_Nb,
            m.Steep_stck,
            m.Steep_Nb,
            m.Base_stck,
            m.Base_Nb,
            m.Flat_stck,
            m.Flat_Nb,
            m.ShortRUP_stck,
            m.ShortRUP_Nb,
            m.Up200_stck,
            m.Up200_Nb
        FROM jerarquias j
        LEFT JOIN metrics_agg m
        ON m.Jerarquia = j.Jerarquia
    ) AS tabla_union

    UNION ALL

    -- TOTAL GRUPO 1: todos los 1.x
    SELECT 'TOTAL GRUPO 1' AS Jerarquia,
        SUM(Dwn200_stck)    AS Dwn200_stck,
        SUM(Dwn200_Nb)      AS Dwn200_Nb,
        SUM(ShortRDwn_stck) AS ShortRDwn_stck,
        SUM(ShortRDwn_Nb)   AS ShortRDwn_Nb,
        SUM(Steep_stck)     AS Steep_stck,
        SUM(Steep_Nb)       AS Steep_Nb,
        SUM(Base_stck)      AS Base_stck,
        SUM(Base_Nb)        AS Base_Nb,
        SUM(Flat_stck)      AS Flat_stck,
        SUM(Flat_Nb)        AS Flat_Nb,
        SUM(ShortRUP_stck)  AS ShortRUP_stck,
        SUM(ShortRUP_Nb)    AS ShortRUP_Nb,
        SUM(Up200_stck)     AS Up200_stck,
        SUM(Up200_Nb)       AS Up200_Nb
    FROM metrics_agg
    WHERE Jerarquia IN (
        '1.1 BANCOS',
        '1.2 CARTERA MAYORISTA',
        '1.3 DEPOSITOS OTRAS EECC',
        '1.4 INVERSION CREDITICIA',
        '1.5 RENTA FIJA',
        '1.6 ACTIVOS DUDOSOS',
        '1.7 ACTIVOS NO SENSIBLE'
    )

    UNION ALL

    -- TOTAL GRUPO 2: todas las 2.x
    SELECT 'TOTAL GRUPO 2' AS Jerarquia,
        SUM(Dwn200_stck)    AS Dwn200_stck,
        SUM(Dwn200_Nb)      AS Dwn200_Nb,
        SUM(ShortRDwn_stck) AS ShortRDwn_stck,
        SUM(ShortRDwn_Nb)   AS ShortRDwn_Nb,
        SUM(Steep_stck)     AS Steep_stck,
        SUM(Steep_Nb)       AS Steep_Nb,
        SUM(Base_stck)      AS Base_stck,
        SUM(Base_Nb)        AS Base_Nb,
        SUM(Flat_stck)      AS Flat_stck,
        SUM(Flat_Nb)        AS Flat_Nb,
        SUM(ShortRUP_stck)  AS ShortRUP_stck,
        SUM(ShortRUP_Nb)    AS ShortRUP_Nb,
        SUM(Up200_stck)     AS Up200_stck,
        SUM(Up200_Nb)       AS Up200_Nb
    FROM metrics_agg
    WHERE Jerarquia IN (
        '2.1 DEPOSITOS OTRAS EECC',
        '2.2 DEPOSITOS A PLAZO DE CLIENTES',
        '2.3 DEPOSITOS A LA VISTA DE CLIENTES',
        '2.4 PASIVOS NO SENSIBLES',
        '2.5 CARTERA MAYORISTA'
    )

    UNION ALL

    -- TOTAL GRUPO 3: derivados
    SELECT 'TOTAL GRUPO 3' AS Jerarquia,
        SUM(Dwn200_stck)    AS Dwn200_stck,
        SUM(Dwn200_Nb)      AS Dwn200_Nb,
        SUM(ShortRDwn_stck) AS ShortRDwn_stck,
        SUM(ShortRDwn_Nb)   AS ShortRDwn_Nb,
        SUM(Steep_stck)     AS Steep_stck,
        SUM(Steep_Nb)       AS Steep_Nb,
        SUM(Base_stck)      AS Base_stck,
        SUM(Base_Nb)        AS Base_Nb,
        SUM(Flat_stck)      AS Flat_stck,
        SUM(Flat_Nb)        AS Flat_Nb,
        SUM(ShortRUP_stck)  AS ShortRUP_stck,
        SUM(ShortRUP_Nb)    AS ShortRUP_Nb,
        SUM(Up200_stck)     AS Up200_stck,
        SUM(Up200_Nb)       AS Up200_Nb
    FROM metrics_agg
    WHERE Jerarquia = '3.1 DERIVADOS'

    UNION ALL

    -- TOTAL GRUPO 4: patrimonio neto
    SELECT 'TOTAL GRUPO 4' AS Jerarquia,
        SUM(Dwn200_stck)    AS Dwn200_stck,
        SUM(Dwn200_Nb)      AS Dwn200_Nb,
        SUM(ShortRDwn_stck) AS ShortRDwn_stck,
        SUM(ShortRDwn_Nb)   AS ShortRDwn_Nb,
        SUM(Steep_stck)     AS Steep_stck,
        SUM(Steep_Nb)       AS Steep_Nb,
        SUM(Base_stck)      AS Base_stck,
        SUM(Base_Nb)        AS Base_Nb,
        SUM(Flat_stck)      AS Flat_stck,
        SUM(Flat_Nb)        AS Flat_Nb,
        SUM(ShortRUP_stck)  AS ShortRUP_stck,
        SUM(ShortRUP_Nb)    AS ShortRUP_Nb,
        SUM(Up200_stck)     AS Up200_stck,
        SUM(Up200_Nb)       AS Up200_Nb
    FROM metrics_agg
    WHERE Jerarquia = '4.1 PN'

    ORDER BY CASE Jerarquia
        WHEN '1.1 BANCOS'                      THEN 1
        WHEN '1.2 CARTERA MAYORISTA'          THEN 2
        WHEN '1.3 DEPOSITOS OTRAS EECC'       THEN 3
        WHEN '1.4 INVERSION CREDITICIA'       THEN 4
        WHEN '1.5 RENTA FIJA'                 THEN 5
        WHEN '1.6 ACTIVOS DUDOSOS'            THEN 6
        WHEN '1.7 ACTIVOS NO SENSIBLE'        THEN 7
        WHEN 'TOTAL GRUPO 1'                  THEN 8
        WHEN ' '                              THEN 9
        WHEN '2.1 DEPOSITOS OTRAS EECC'       THEN 10
        WHEN '2.2 DEPOSITOS A PLAZO DE CLIENTES' THEN 11
        WHEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES' THEN 12
        WHEN '  '                             THEN 13
        WHEN '2.4 PASIVOS NO SENSIBLES'       THEN 14
        WHEN '2.5 CARTERA MAYORISTA'          THEN 15
        WHEN 'TOTAL GRUPO 2'                  THEN 16
        WHEN '3.1 DERIVADOS'                  THEN 17
        WHEN 'TOTAL GRUPO 3'                  THEN 18
        WHEN '4.1 PN'                         THEN 19
        WHEN 'TOTAL GRUPO 4'                  THEN 20
        ELSE 999
    END;
    """


    # Se define nuestro data frame con los datos extrídos de la base de datos
    df_MF_EBA = pd.read_sql(Informe_MF_EBA, conn)

    # Nombre de la hoja donde quieres pegar los datos. Ahora se construye dinámicamente.
    nombre_hoja_MF_EBA = f"Sensibilidades MF - EBA - {mes_actual_nombre}"

    # Acceder a la hoja
    ws_MF_EBA = wb[nombre_hoja_MF_EBA]


    # Se definen las coordenadas donde queremos que se peguen los datos
    start_row_MF_EBA = 6
    start_col_MF_EBA = 36

    # Límite de filas donde se pegarán los datos (de la 6 a la 25)
    end_row_MF_EBA = 25

    # Escribir los datos en las coordenadas anteriores
    for i, row in enumerate(df_MF_EBA.values, start=start_row_MF_EBA):
        if i > end_row_MF_EBA:
            break  # Detener la escritura si superamos la fila 25
        for j, value_MF_EBA in enumerate(row[1:], start=start_col_MF_EBA):
            ws_MF_EBA.cell(row=i, column=j).value = value_MF_EBA








    ###################################### PROCESO DE PEGADO DE INFORMES DE MARGEN FINANCIERO CTO #####################################
    Informe_MF_CTO= f"""
        WITH jerarquias AS (
        -- GRUPO 1
        SELECT 1  AS orden, '1.1 BANCOS'                         AS Jerarquia UNION ALL
        SELECT 2,             '1.2 CARTERA MAYORISTA'                         UNION ALL
        SELECT 3,             '1.3 DEPOSITOS OTRAS EECC'                      UNION ALL
        SELECT 4,             '1.4 INVERSION CREDITICIA'                      UNION ALL
        SELECT 5,             '1.5 RENTA FIJA'                                UNION ALL
        SELECT 6,             '1.6 ACTIVOS DUDOSOS'                           UNION ALL
        SELECT 7,             '1.7 ACTIVOS NO SENSIBLE'                       UNION ALL
        -- GRUPO 2
        SELECT 9,            '2.1 DEPOSITOS OTRAS EECC'                      UNION ALL
        SELECT 10,            '2.2 DEPOSITOS A PLAZO DE CLIENTES'             UNION ALL
        SELECT 11,            '2.3 DEPOSITOS A LA VISTA DE CLIENTES'          UNION ALL
        SELECT 12,            '2.4 PASIVOS NO SENSIBLES'                      UNION ALL
        SELECT 13,            '2.5 CARTERA MAYORISTA'                         UNION ALL
        -- GRUPO 3
        SELECT 15,            '3.1 DERIVADOS'                                 UNION ALL
        -- GRUPO 4
        SELECT 17,            '4.1 PN'
    ),

    metrics_agg AS (
        SELECT
            CASE
            WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') THEN '1.1 BANCOS'
            WHEN dim_1 IN ('A_Ata') THEN '1.2 CARTERA MAYORISTA'
            WHEN dim_1 IN ('A_Depositos cedidos') THEN '1.3 DEPOSITOS OTRAS EECC'
            WHEN dim_1 IN (
                    'A_ECORP_Otros_Activo Corto Plazo',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                    'A_EESP_Otros_Activo Corto Plazo',
                    'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                    'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                    'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                    'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                    'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                    'A_PIB_Otros_Descubiertos_Gtia Personal',
                    'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                    'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                    'A_POFI_Otros_Descubiertos_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                    'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                    'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_POFI_Otros_Prest. Origen_Gtia Aval',
                    'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                    'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                    'A_SSCC_Otros_Descubiertos_Gtia Personal',
                    'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                    'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
                ) THEN '1.4 INVERSION CREDITICIA'
            WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') THEN '1.5 RENTA FIJA'
            WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') THEN '1.6 ACTIVOS DUDOSOS'
            WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') THEN '1.7 ACTIVOS NO SENSIBLE'
            WHEN dim_1 IN ('P_Cuentas de otras EECC') THEN '2.1 DEPOSITOS OTRAS EECC'
            WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES'
            WHEN dim_1 IN (
                'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
                'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
            ) THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES'
            WHEN dim_1 IN ('P_Resto no sensible') THEN '2.4 PASIVOS NO SENSIBLES'
            WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') THEN '2.5 CARTERA MAYORISTA'
            WHEN dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo') THEN '3.1 DERIVADOS'
            WHEN dim_1 IN ('Patrimonio Neto') THEN '4.1 PN'
            ELSE 'OTROS'
            END AS Jerarquia,
            SUM(CASE WHEN scenario = 'Dwn200_C' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_stck,
            SUM(CASE WHEN scenario = 'Dwn200_C' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_Nb,
            SUM(CASE WHEN scenario = 'Dwn100_C' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_stck,
            SUM(CASE WHEN scenario = 'Dwn100_C' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_Nb,
            SUM(CASE WHEN scenario = 'Dwn75_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_stck,
            SUM(CASE WHEN scenario = 'Dwn75_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_Nb,
            SUM(CASE WHEN scenario = 'Dwn50_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_stck,
            SUM(CASE WHEN scenario = 'Dwn50_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_Nb,
            SUM(CASE WHEN scenario = 'Dwn25_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_stck,
            SUM(CASE WHEN scenario = 'Dwn25_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_Nb,
            SUM(CASE WHEN scenario = 'Base_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_stck,
            SUM(CASE WHEN scenario = 'Base_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_Nb,
            SUM(CASE WHEN scenario = 'Up25_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_stck,
            SUM(CASE WHEN scenario = 'Up25_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_Nb,
            SUM(CASE WHEN scenario = 'Up50_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_stck,
            SUM(CASE WHEN scenario = 'Up50_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_Nb,
            SUM(CASE WHEN scenario = 'Up75_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_stck,
            SUM(CASE WHEN scenario = 'Up75_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_Nb,
            SUM(CASE WHEN scenario = 'Up100_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_stck,
            SUM(CASE WHEN scenario = 'Up100_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_Nb,
            SUM(CASE WHEN scenario = 'Up200_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_stck,
            SUM(CASE WHEN scenario = 'Up200_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_Nb
        FROM pro_pichincha_alquid_old."result"
        WHERE load_id IN (
        '{load_ids['cierre_base']}',
        '{load_ids['cierre_up']}',
        '{load_ids['cierre_dwn']}'
        )
        GROUP BY 1
    )

    -- 1) Jerarquías fijas (1.x, 2.x, 3.1, 4.1 y la línea en blanco) con LEFT JOIN
    SELECT
        j.Jerarquia,
        a.Dwn200_stck,
        a.Dwn200_Nb,
        a.Dwn100_stck,
        a.Dwn100_Nb,
        a.Dwn75_stck,
        a.Dwn75_Nb,
        a.Dwn50_stck,
        a.Dwn50_Nb,
        a.Dwn25_stck,
        a.Dwn25_Nb,
        a.Base_stck,
        a.Base_Nb,
        a.Up25_stck,
        a.Up25_Nb,
        a.Up50_stck,
        a.Up50_Nb,
        a.Up75_stck,
        a.Up75_Nb,
        a.Up100_stck,
        a.Up100_Nb,
        a.Up200_stck,
        a.Up200_Nb
    FROM jerarquias j
    LEFT JOIN metrics_agg a
        ON a.Jerarquia = j.Jerarquia

    UNION ALL

    -- 2) TOTAL GRUPO 1
    SELECT 'TOTAL GRUPO 1' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200_C' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_stck,
        SUM(CASE WHEN scenario = 'Dwn200_C' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_Nb,
        SUM(CASE WHEN scenario = 'Dwn100_C' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_stck,
        SUM(CASE WHEN scenario = 'Dwn100_C' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_Nb,
        SUM(CASE WHEN scenario = 'Dwn75_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_stck,
        SUM(CASE WHEN scenario = 'Dwn75_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_Nb,
        SUM(CASE WHEN scenario = 'Dwn50_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_stck,
        SUM(CASE WHEN scenario = 'Dwn50_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_Nb,
        SUM(CASE WHEN scenario = 'Dwn25_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_stck,
        SUM(CASE WHEN scenario = 'Dwn25_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_Nb,
        SUM(CASE WHEN scenario = 'Base_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_stck,
        SUM(CASE WHEN scenario = 'Base_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_Nb,
        SUM(CASE WHEN scenario = 'Up25_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_stck,
        SUM(CASE WHEN scenario = 'Up25_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_Nb,
        SUM(CASE WHEN scenario = 'Up50_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_stck,
        SUM(CASE WHEN scenario = 'Up50_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_Nb,
        SUM(CASE WHEN scenario = 'Up75_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_stck,
        SUM(CASE WHEN scenario = 'Up75_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_Nb,
        SUM(CASE WHEN scenario = 'Up100_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_stck,
        SUM(CASE WHEN scenario = 'Up100_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_Nb,
        SUM(CASE WHEN scenario = 'Up200_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_stck,
        SUM(CASE WHEN scenario = 'Up200_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_Nb
    FROM pro_pichincha_alquid_old."result"
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )
    AND dim_1 IN (
    'A_Cuentas en otras EECC', 'A_Tesoreria', 'A_Ata', 'A_Depositos cedidos',
    'A_ECORP_Otros_Activo Corto Plazo', 'A_EESP_Otros_Activo Corto Plazo', 'A_POFI_Otros_Ant. Nomina_Gtia Personal',
    'A_ECORP_Otros_Cred. Empresa_Gtia Aval', 'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
    'A_EESP_Otros_Cred. Empresa_Gtia Aval', 'A_EESP_Otros_Cred. Empresa_Gtia Personal',
    'A_SSCC_Otros_Cred. Empresa_Gtia Aval', 'A_SSCC_Otros_Cred. Empresa_Gtia Hipot.', 'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
    'A_PIB_Otros_Descubiertos_Gtia Personal', 'A_POFI_Otros_Descubiertos_Gtia Personal', 'A_SSCC_Otros_Descubiertos_Gtia Personal',
    'A_ECORP_Otros_Prest. Empresas_Gtia Aval', 'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.', 'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
    'A_ECORP_Otros_Prest. Empresas_Gtia Prenda', 'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',
    'A_EESP_Otros_Prest. Empresas_Gtia Aval', 'A_EESP_Otros_Prest. Empresas_Gtia Personal', 'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
    'A_EESP_Otros_Prest. Promotor_Gtia Hipot.', 'A_PIB_Otros_Prest. Consumo_Gtia Personal', 'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
    'A_PIB_Otros_Prest. Hipotecas_Gtia Personal', 'A_POFI_Otros_Prest. Consumo_Gtia Personal', 'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
    'A_POFI_Otros_Prest. Empresas_Gtia Aval', 'A_POFI_Otros_Prest. Empresas_Gtia Hipot.', 'A_POFI_Otros_Prest. Empresas_Gtia Personal',
    'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
    'A_POFI_Otros_Prest. Origen_Gtia Aval', 'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
    'A_SSCC_BKIA_Prest. Consumo_Gtia Personal', 'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
    'A_SSCC_Otros_Prest. Empresas_Gtia Aval', 'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
    'A_SSCC_Otros_Prest. Empresas_Gtia Personal', 'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
    'A_PIB_Otros_Tarjetas Credito_Gtia Personal', 'A_POFI_Otros_Tarjetas Credito_Gtia Personal',
    'A_SSCC_Otros_Tarjetas Credito_Gtia Personal', 'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal',
    'A_Bonos corporativos', 'A_Bonos soberanos', 'A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso',
    'A_Resto no sensible', 'A_Tesoreria_Admin.', 'A_EESP_Otros_Prest. Promotor_Gtia Aval', 'A_EESP_Otros_Prest. Promotor_Gtia Personal'
    )

    UNION ALL
    -- 3) TOTAL GRUPO 2
    SELECT 'TOTAL GRUPO 2' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200_C' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_stck,
        SUM(CASE WHEN scenario = 'Dwn200_C' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_Nb,
        SUM(CASE WHEN scenario = 'Dwn100_C' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_stck,
        SUM(CASE WHEN scenario = 'Dwn100_C' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_Nb,
        SUM(CASE WHEN scenario = 'Dwn75_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_stck,
        SUM(CASE WHEN scenario = 'Dwn75_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_Nb,
        SUM(CASE WHEN scenario = 'Dwn50_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_stck,
        SUM(CASE WHEN scenario = 'Dwn50_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_Nb,
        SUM(CASE WHEN scenario = 'Dwn25_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_stck,
        SUM(CASE WHEN scenario = 'Dwn25_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_Nb,
        SUM(CASE WHEN scenario = 'Base_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_stck,
        SUM(CASE WHEN scenario = 'Base_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_Nb,
        SUM(CASE WHEN scenario = 'Up25_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_stck,
        SUM(CASE WHEN scenario = 'Up25_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_Nb,
        SUM(CASE WHEN scenario = 'Up50_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_stck,
        SUM(CASE WHEN scenario = 'Up50_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_Nb,
        SUM(CASE WHEN scenario = 'Up75_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_stck,
        SUM(CASE WHEN scenario = 'Up75_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_Nb,
        SUM(CASE WHEN scenario = 'Up100_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_stck,
        SUM(CASE WHEN scenario = 'Up100_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_Nb,
        SUM(CASE WHEN scenario = 'Up200_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_stck,
        SUM(CASE WHEN scenario = 'Up200_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_Nb
    FROM pro_pichincha_alquid_old."result"
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )
    AND dim_1 IN (
    'P_Cuentas de otras EECC', 'P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo',
    'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
    'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
    'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas',
    'P_Resto no sensible', 'P_Ata', 'P_Emisiones T2'
    )

    UNION ALL
    -- 4) TOTAL GRUPO 3
    SELECT 'TOTAL GRUPO 3' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200_C' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_stck,
        SUM(CASE WHEN scenario = 'Dwn200_C' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_Nb,
        SUM(CASE WHEN scenario = 'Dwn100_C' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_stck,
        SUM(CASE WHEN scenario = 'Dwn100_C' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_Nb,
        SUM(CASE WHEN scenario = 'Dwn75_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_stck,
        SUM(CASE WHEN scenario = 'Dwn75_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_Nb,
        SUM(CASE WHEN scenario = 'Dwn50_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_stck,
        SUM(CASE WHEN scenario = 'Dwn50_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_Nb,
        SUM(CASE WHEN scenario = 'Dwn25_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_stck,
        SUM(CASE WHEN scenario = 'Dwn25_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_Nb,
        SUM(CASE WHEN scenario = 'Base_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_stck,
        SUM(CASE WHEN scenario = 'Base_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_Nb,
        SUM(CASE WHEN scenario = 'Up25_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_stck,
        SUM(CASE WHEN scenario = 'Up25_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_Nb,
        SUM(CASE WHEN scenario = 'Up50_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_stck,
        SUM(CASE WHEN scenario = 'Up50_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_Nb,
        SUM(CASE WHEN scenario = 'Up75_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_stck,
        SUM(CASE WHEN scenario = 'Up75_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_Nb,
        SUM(CASE WHEN scenario = 'Up100_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_stck,
        SUM(CASE WHEN scenario = 'Up100_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_Nb,
        SUM(CASE WHEN scenario = 'Up200_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_stck,
        SUM(CASE WHEN scenario = 'Up200_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_Nb
    FROM pro_pichincha_alquid_old."result"
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )
    AND dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo')

    UNION ALL
    -- 5) TOTAL GRUPO 4
    SELECT 'TOTAL GRUPO 4' AS Jerarquia,
        SUM(CASE WHEN scenario = 'Dwn200_C' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_stck,
        SUM(CASE WHEN scenario = 'Dwn200_C' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_Nb,
        SUM(CASE WHEN scenario = 'Dwn100_C' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_stck,
        SUM(CASE WHEN scenario = 'Dwn100_C' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn100_Nb,
        SUM(CASE WHEN scenario = 'Dwn75_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_stck,
        SUM(CASE WHEN scenario = 'Dwn75_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn75_Nb,
        SUM(CASE WHEN scenario = 'Dwn50_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_stck,
        SUM(CASE WHEN scenario = 'Dwn50_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn50_Nb,
        SUM(CASE WHEN scenario = 'Dwn25_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_stck,
        SUM(CASE WHEN scenario = 'Dwn25_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn25_Nb,
        SUM(CASE WHEN scenario = 'Base_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_stck,
        SUM(CASE WHEN scenario = 'Base_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_Nb,
        SUM(CASE WHEN scenario = 'Up25_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_stck,
        SUM(CASE WHEN scenario = 'Up25_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up25_Nb,
        SUM(CASE WHEN scenario = 'Up50_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_stck,
        SUM(CASE WHEN scenario = 'Up50_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up50_Nb,
        SUM(CASE WHEN scenario = 'Up75_C'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_stck,
        SUM(CASE WHEN scenario = 'Up75_C'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up75_Nb,
        SUM(CASE WHEN scenario = 'Up100_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_stck,
        SUM(CASE WHEN scenario = 'Up100_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up100_Nb,
        SUM(CASE WHEN scenario = 'Up200_C'  AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_stck,
        SUM(CASE WHEN scenario = 'Up200_C'  AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Up200_Nb
    FROM pro_pichincha_alquid_old."result"
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )
    AND dim_1 IN ('Patrimonio Neto')

    ORDER BY CASE Jerarquia
    WHEN '1.1 BANCOS'                         THEN 1
    WHEN '1.2 CARTERA MAYORISTA'             THEN 2
    WHEN '1.3 DEPOSITOS OTRAS EECC'          THEN 3
    WHEN '1.4 INVERSION CREDITICIA'          THEN 4
    WHEN '1.5 RENTA FIJA'                    THEN 5
    WHEN '1.6 ACTIVOS DUDOSOS'               THEN 6
    WHEN '1.7 ACTIVOS NO SENSIBLE'           THEN 7
    WHEN 'TOTAL GRUPO 1'                     THEN 8
    WHEN '2.1 DEPOSITOS OTRAS EECC'          THEN 9
    WHEN '2.2 DEPOSITOS A PLAZO DE CLIENTES' THEN 10
    WHEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES' THEN 11
    WHEN '2.4 PASIVOS NO SENSIBLES'          THEN 12
    WHEN '2.5 CARTERA MAYORISTA'             THEN 13
    WHEN 'TOTAL GRUPO 2'                     THEN 14
    WHEN '3.1 DERIVADOS'                     THEN 15
    WHEN 'TOTAL GRUPO 3'                     THEN 16
    WHEN '4.1 PN'                            THEN 17
    WHEN 'TOTAL GRUPO 4'                     THEN 18
    ELSE 999
    END;
    """

    # Se define nuestro data frame con los datos extrídos de la base de datos
    df_MF_CTO = pd.read_sql(Informe_MF_CTO, conn)

    # Nombre de la hoja donde quieres pegar los datos. Ahora se construye dinámicamente.
    nombre_hoja_MF_CTO = f"MF - Sensibilidades - Cto - {mes_actual_nombre}"

    # Acceder a la hoja
    ws_MF_CTO = wb[nombre_hoja_MF_CTO]


    # Se definen las coordenadas donde queremos que se peguen los datos
    start_row_MF_CTO = 6
    start_col_MF_CTO = 50

    # Límite de filas donde se pegarán los datos (de la 6 a la 23)
    end_row_MF_CTO = 23

    # Escribir los datos en las coordenadas anteriores
    for i, row in enumerate(df_MF_CTO.values, start=start_row_MF_CTO):
        if i > end_row_MF_CTO:
            break  # Detener la escritura si superamos la fila 23
        for j, value_MF_CTO in enumerate(row[1:], start=start_col_MF_CTO):
            ws_MF_CTO.cell(row=i, column=j).value = value_MF_CTO




    ###################################### PROCESO DE PEGADO DE INFORMES DE DetalleVF #####################################
    Informe_Detalle_VF= f"""
        WITH jerarquias AS (
        SELECT  1 AS orden, '1.1 BANCOS'                        AS Jerarquia UNION ALL
        SELECT  2,           '1.2 CARTERA MAYORISTA'                        UNION ALL
        SELECT  3,           '1.3 DEPOSITOS OTRAS EECC'                     UNION ALL
        SELECT  4,           '1.4 INVERSION CREDITICIA'                     UNION ALL
        SELECT  5,           '1.5 RENTA FIJA'                               UNION ALL
        SELECT  6,           '1.6 ACTIVOS DUDOSOS'                          UNION ALL
        SELECT  7,           '1.7 ACTIVOS NO SENSIBLE'                      UNION ALL
        SELECT  8,           '2.1 DEPOSITOS OTRAS EECC'                     UNION ALL
        SELECT  9,           '2.2 DEPOSITOS A PLAZO DE CLIENTES'            UNION ALL
        SELECT 10,           '2.3 DEPOSITOS A LA VISTA DE CLIENTES'         UNION ALL
        SELECT 11,           '2.4 PASIVOS NO SENSIBLES'                     UNION ALL
        SELECT 12,           '2.5 CARTERA MAYORISTA'                        UNION ALL
        SELECT 13,           '3.1 DERIVADOS'                                UNION ALL
        SELECT 14,           '4.1 PN'
    ),

    metrics_agg AS (
        SELECT
            CASE
            WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') THEN '1.1 BANCOS'
            WHEN dim_1 IN ('A_Ata') THEN '1.2 CARTERA MAYORISTA'
            WHEN dim_1 IN ('A_Depositos cedidos') THEN '1.3 DEPOSITOS OTRAS EECC'
            WHEN dim_1 IN (
                    'A_ECORP_Otros_Activo Corto Plazo',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                    'A_EESP_Otros_Activo Corto Plazo',
                    'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                    'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                    'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                    'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                    'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                    'A_PIB_Otros_Descubiertos_Gtia Personal',
                    'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                    'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                    'A_POFI_Otros_Descubiertos_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                    'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                    'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_POFI_Otros_Prest. Origen_Gtia Aval',
                    'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                    'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                    'A_SSCC_Otros_Descubiertos_Gtia Personal',
                    'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                    'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
                ) THEN '1.4 INVERSION CREDITICIA'
            WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') THEN '1.5 RENTA FIJA'
            WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') THEN '1.6 ACTIVOS DUDOSOS'
            WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') THEN '1.7 ACTIVOS NO SENSIBLE'
            WHEN dim_1 IN ('P_Cuentas de otras EECC') THEN '2.1 DEPOSITOS OTRAS EECC'
            WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES'
            WHEN dim_1 IN (
                'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
                'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
            ) THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES'
            WHEN dim_1 IN ('P_Resto no sensible') THEN '2.4 PASIVOS NO SENSIBLES'
            WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') THEN '2.5 CARTERA MAYORISTA'
            WHEN dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo') THEN '3.1 DERIVADOS'
            WHEN dim_1 IN ('Patrimonio Neto') THEN '4.1 PN'
            ELSE 'OTROS'
            END AS Jerarquia,

            -- Dwn200
            SUM(CASE WHEN scenario = 'Dwn200' THEN ABS(effective_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Dwn200' THEN ABS(market_value) ELSE 0 END), 0) / 100 AS effective_duration_pond_Dwn200,
            SUM(CASE WHEN scenario = 'Dwn200' THEN ABS(macaulay_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Dwn200' THEN ABS(market_value) ELSE 0 END), 0) AS duracion_macaulay_pond_Dwn200,
            SUM(CASE WHEN scenario = 'Dwn200' THEN market_value ELSE 0 END) AS Dwn200,

            -- Dwn100
            SUM(CASE WHEN scenario = 'Dwn100' THEN ABS(effective_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Dwn100' THEN ABS(market_value) ELSE 0 END), 0) / 100 AS effective_duration_pond_Dwn100,
            SUM(CASE WHEN scenario = 'Dwn100' THEN ABS(macaulay_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Dwn100' THEN ABS(market_value) ELSE 0 END), 0) AS duracion_macaulay_pond_Dwn100,
            SUM(CASE WHEN scenario = 'Dwn100' THEN market_value ELSE 0 END) AS Dwn100,

            -- Dwn75
            SUM(CASE WHEN scenario = 'Dwn75' THEN ABS(effective_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Dwn75' THEN ABS(market_value) ELSE 0 END), 0) / 100 AS effective_duration_pond_Dwn75,
            SUM(CASE WHEN scenario = 'Dwn75' THEN ABS(macaulay_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Dwn75' THEN ABS(market_value) ELSE 0 END), 0) AS duracion_macaulay_pond_Dwn75,
            SUM(CASE WHEN scenario = 'Dwn75' THEN market_value ELSE 0 END) AS Dwn75,

            -- Dwn50
            SUM(CASE WHEN scenario = 'Dwn50' THEN ABS(effective_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Dwn50' THEN ABS(market_value) ELSE 0 END), 0) / 100 AS effective_duration_pond_Dwn50,
            SUM(CASE WHEN scenario = 'Dwn50' THEN ABS(macaulay_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Dwn50' THEN ABS(market_value) ELSE 0 END), 0) AS duracion_macaulay_pond_Dwn50,
            SUM(CASE WHEN scenario = 'Dwn50' THEN market_value ELSE 0 END) AS Dwn50,

            -- Dwn25
            SUM(CASE WHEN scenario = 'Dwn25' THEN ABS(effective_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Dwn25' THEN ABS(market_value) ELSE 0 END), 0) / 100 AS effective_duration_pond_Dwn25,
            SUM(CASE WHEN scenario = 'Dwn25' THEN ABS(macaulay_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Dwn25' THEN ABS(market_value) ELSE 0 END), 0) AS duracion_macaulay_pond_Dwn25,
            SUM(CASE WHEN scenario = 'Dwn25' THEN market_value ELSE 0 END) AS Dwn25,

            -- Base
            SUM(CASE WHEN scenario = 'Base' THEN ABS(effective_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Base' THEN ABS(market_value) ELSE 0 END), 0) / 100 AS effective_duration_pond_Base,
            SUM(CASE WHEN scenario = 'Base' THEN ABS(macaulay_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Base' THEN ABS(market_value) ELSE 0 END), 0) AS duracion_macaulay_pond_Base,
            SUM(CASE WHEN scenario = 'Base' THEN market_value ELSE 0 END) AS Base,

            -- Up25
            SUM(CASE WHEN scenario = 'Up25' THEN ABS(effective_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Up25' THEN ABS(market_value) ELSE 0 END), 0) / 100 AS effective_duration_pond_Up25,
            SUM(CASE WHEN scenario = 'Up25' THEN ABS(macaulay_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Up25' THEN ABS(market_value) ELSE 0 END), 0) AS duracion_macaulay_pond_Up25,
            SUM(CASE WHEN scenario = 'Up25' THEN market_value ELSE 0 END) AS Up25,

            -- Up50
            SUM(CASE WHEN scenario = 'Up50' THEN ABS(effective_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Up50' THEN ABS(market_value) ELSE 0 END), 0) / 100 AS effective_duration_pond_Up50,
            SUM(CASE WHEN scenario = 'Up50' THEN ABS(macaulay_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Up50' THEN ABS(market_value) ELSE 0 END), 0) AS duracion_macaulay_pond_Up50,
            SUM(CASE WHEN scenario = 'Up50' THEN market_value ELSE 0 END) AS Up50,

            -- Up75
            SUM(CASE WHEN scenario = 'Up75' THEN ABS(effective_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Up75' THEN ABS(market_value) ELSE 0 END), 0) / 100 AS effective_duration_pond_Up75,
            SUM(CASE WHEN scenario = 'Up75' THEN ABS(macaulay_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Up75' THEN ABS(market_value) ELSE 0 END), 0) AS duracion_macaulay_pond_Up75,
            SUM(CASE WHEN scenario = 'Up75' THEN market_value ELSE 0 END) AS Up75,

            -- Up100
            SUM(CASE WHEN scenario = 'Up100' THEN ABS(effective_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Up100' THEN ABS(market_value) ELSE 0 END), 0) / 100 AS effective_duration_pond_Up100,
            SUM(CASE WHEN scenario = 'Up100' THEN ABS(macaulay_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Up100' THEN ABS(market_value) ELSE 0 END), 0) AS duracion_macaulay_pond_Up100,
            SUM(CASE WHEN scenario = 'Up100' THEN market_value ELSE 0 END) AS Up100,

            -- Up200
            SUM(CASE WHEN scenario = 'Up200' THEN ABS(effective_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Up200' THEN ABS(market_value) ELSE 0 END), 0) / 100 AS effective_duration_pond_Up200,
            SUM(CASE WHEN scenario = 'Up200' THEN ABS(macaulay_duration * market_value) ELSE 0 END)
            / NULLIF(SUM(CASE WHEN scenario = 'Up200' THEN ABS(market_value) ELSE 0 END), 0) AS duracion_macaulay_pond_Up200,
            SUM(CASE WHEN scenario = 'Up200' THEN market_value ELSE 0 END) AS Up200

        FROM pro_pichincha_alquid_old.metric
        WHERE load_id IN (
            '{load_ids['cierre_base']}',
            '{load_ids['cierre_up']}',
            '{load_ids['cierre_dwn']}'
        )
        GROUP BY 1
    )

    SELECT
        j.Jerarquia,
        a.effective_duration_pond_Dwn200,
        a.duracion_macaulay_pond_Dwn200,
        a.Dwn200,
        a.effective_duration_pond_Dwn100,
        a.duracion_macaulay_pond_Dwn100,
        a.Dwn100,
        a.effective_duration_pond_Dwn75,
        a.duracion_macaulay_pond_Dwn75,
        a.Dwn75,
        a.effective_duration_pond_Dwn50,
        a.duracion_macaulay_pond_Dwn50,
        a.Dwn50,
        a.effective_duration_pond_Dwn25,
        a.duracion_macaulay_pond_Dwn25,
        a.Dwn25,
        a.effective_duration_pond_Base,
        a.duracion_macaulay_pond_Base,
        a.Base,
        a.effective_duration_pond_Up25,
        a.duracion_macaulay_pond_Up25,
        a.Up25,
        a.effective_duration_pond_Up50,
        a.duracion_macaulay_pond_Up50,
        a.Up50,
        a.effective_duration_pond_Up75,
        a.duracion_macaulay_pond_Up75,
        a.Up75,
        a.effective_duration_pond_Up100,
        a.duracion_macaulay_pond_Up100,
        a.Up100,
        a.effective_duration_pond_Up200,
        a.duracion_macaulay_pond_Up200,
        a.Up200
    FROM jerarquias j
    LEFT JOIN metrics_agg a
        ON a.Jerarquia = j.Jerarquia
    ORDER BY j.orden;
    """


    # Se define nuestro data frame con los datos extrídos de la base de datos
    df_Detalle_VF = pd.read_sql(Informe_Detalle_VF, conn)

    # Nombre de la hoja donde quieres pegar los datos. Se mantiene 'detalleVF' según el original.
    # Si esta hoja también debe ser versionada por mes, la lógica debe cambiar.
    nombre_hoja_Detalle_VF = "detalleVF"

    # Acceder a la hoja
    ws_Detalle_VF = wb[nombre_hoja_Detalle_VF] # Asumo que 'detalleVF' está en 'Plantilla_COAP.xlsx'


    # Se definen las coordenadas donde queremos que se peguen los datos
    start_row_Detalle_VF = 4
    start_col_Detalle_VF = 3

    # Límite de filas donde se pegarán los datos (de la 4 a la 17)
    end_row_Detalle_VF = 17

    # Escribir los datos en las coordenadas anteriores
    for i, row in enumerate(df_Detalle_VF.values, start=start_row_Detalle_VF):
        if i > end_row_Detalle_VF:
            break  # Detener la escritura si superamos la fila 17
        for j, value_Detalle_VF in enumerate(row[1:], start=start_col_Detalle_VF):
            ws_Detalle_VF.cell(row=i, column=j).value = value_Detalle_VF




    ###################################### PROCESO DE PEGADO DE INFORMES DE DetalleMF #####################################
    Informe_Detalle_MF= f"""
        WITH jerarquias AS (
        SELECT 16 AS orden, ' '   AS JER_PRODUCTO UNION ALL   -- línea en blanco 1
        SELECT 24 AS orden, '  '  AS JER_PRODUCTO UNION ALL   -- línea en blanco 2
        SELECT 30 AS orden, '   ' AS JER_PRODUCTO UNION ALL   -- línea en blanco 3

        SELECT  1 AS orden, 'BANCOS'                          UNION ALL
        SELECT  2,          'CARTERA MAYORISTA'              UNION ALL
        SELECT  3,          'DEPOSITOS OTRAS EECC'           UNION ALL
        SELECT  4,          'A_Balancer'                     UNION ALL
        SELECT  5,          'ACTIVO CORTO PLAZO'             UNION ALL
        SELECT  6,          'ANTICIPO NOMINA'                UNION ALL
        SELECT  7,          'CREDITOS'                       UNION ALL
        SELECT  8,          'DESCUBIERTOS'                   UNION ALL
        SELECT  9,          'PRESTAMOS'                      UNION ALL
        SELECT  10,          'TARJETAS CREDITO'               UNION ALL
        SELECT 11,          'A_Bonos corporativos'           UNION ALL
        SELECT 12,          'A_Bonos soberanos'              UNION ALL
        SELECT 13,          'ACT DUDOSOS'                    UNION ALL
        SELECT 14,          'A_Resto no sensible'            UNION ALL
        SELECT 15,          'A_Tesoreria_Admin.'             UNION ALL

        SELECT 17,          'P_BANCOS'                       UNION ALL
        SELECT 18,          'P_Balancer'                     UNION ALL
        SELECT 19,          'PLAZO'                          UNION ALL
        SELECT 20,          'CTAS NO REMUNERADAS'            UNION ALL
        SELECT 21,          'CTAS REMUNERADAS'               UNION ALL
        SELECT 22,          'P_Resto no sensible'            UNION ALL
        SELECT 23,          'P_Ata'                          UNION ALL
        SELECT 24,          'P_Emisiones T2'                 UNION ALL

        SELECT 26,          'FXSWAP_pago'                    UNION ALL
        SELECT 27,          'FXSWAP_recibo'                  UNION ALL
        SELECT 28,          'IRS_pago'                       UNION ALL
        SELECT 29,          'IRS_recibo'                     UNION ALL
        SELECT 30,          'Patrimonio Neto'                
    ),

    metrics_agg AS (
        SELECT
        CASE
            WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') THEN 'BANCOS'
            WHEN dim_1 = 'A_Ata' THEN 'CARTERA MAYORISTA'
            WHEN dim_1 = 'A_Depositos cedidos' THEN 'DEPOSITOS OTRAS EECC'
            WHEN dim_1 = 'A_Balancer' THEN 'A_Balancer'
            WHEN dim_1 LIKE '%Activo Corto Plazo%' THEN 'ACTIVO CORTO PLAZO'
            WHEN dim_1 LIKE '%Ant. Nomina%' THEN 'ANTICIPO NOMINA'
            -- Importante que estas dos siguientes se mantengan en este orden!!!!
            WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') THEN 'ACT DUDOSOS'
            WHEN dim_1 LIKE '%_Cred.%' THEN 'CREDITOS'
            WHEN dim_1 LIKE '%Descubiertos%' THEN 'DESCUBIERTOS'
            WHEN dim_1 LIKE '%_Prest.%' THEN 'PRESTAMOS'
            WHEN dim_1 LIKE '%Tarjetas Credito%' THEN 'TARJETAS CREDITO'
            WHEN dim_1 = 'A_Bonos corporativos' THEN 'A_Bonos corporativos'
            WHEN dim_1 = 'A_Bonos soberanos' THEN 'A_Bonos soberanos'
            WHEN dim_1 = 'A_Resto no sensible' THEN 'A_Resto no sensible'
            WHEN dim_1 = 'A_Tesoreria_Admin.' THEN 'A_Tesoreria_Admin.'
            WHEN dim_1 = 'P_Cuentas de otras EECC' THEN 'P_BANCOS'
            WHEN dim_1 = 'P_Balancer' THEN 'P_Balancer'
            WHEN dim_1 LIKE 'P_%Plazo%' THEN 'PLAZO'
            WHEN dim_1 LIKE '%Ctas No Remuneradas%' THEN 'CTAS NO REMUNERADAS'
            WHEN dim_1 LIKE '%Ctas Remuneradas%' THEN 'CTAS REMUNERADAS'
            WHEN dim_1 = 'P_Resto no sensible' THEN 'P_Resto no sensible'
            WHEN dim_1 = 'P_Ata' THEN 'P_Ata'
            WHEN dim_1 = 'P_Emisiones T2' THEN 'P_Emisiones T2'
            WHEN dim_1 = 'FXSWAP_pago' THEN 'FXSWAP_pago'
            WHEN dim_1 = 'FXSWAP_recibo' THEN 'FXSWAP_recibo'
            WHEN dim_1 = 'IRS_pago' THEN 'IRS_pago'
            WHEN dim_1 = 'IRS_recibo' THEN 'IRS_recibo'
            WHEN dim_1 = 'Patrimonio Neto' THEN 'Patrimonio Neto'
            WHEN dim_1 = 'A_ECORP_Otros_Tarjetas Credito_Gtia Personal'
                THEN 'A_ECORP_Otros_Tarjetas Credito_Gtia Personal'
            ELSE 'OTROS'
        END AS JER_PRODUCTO,
        SUM(CASE WHEN scenario = 'Dwn200' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_stck,
        SUM(CASE WHEN scenario = 'Dwn200' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_Nb,
        SUM(CASE WHEN scenario = 'Base'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_stck,
        SUM(CASE WHEN scenario = 'Base'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_Nb
        FROM pro_pichincha_alquid_old."result"
        WHERE load_id IN (
        '{load_ids['cierre_base']}',
        '{load_ids['cierre_up']}',
        '{load_ids['cierre_dwn']}'
        )
        GROUP BY 1
    )

    -- 1) Todas las jerarquías + líneas en blanco, siempre presentes
    SELECT
        j.JER_PRODUCTO,
        CASE WHEN j.JER_PRODUCTO = 'A_Balancer' THEN 0 ELSE a.Dwn200_stck END AS Dwn200_stck,
        CASE WHEN j.JER_PRODUCTO = 'A_Balancer' THEN 0 ELSE a.Dwn200_Nb   END AS Dwn200_Nb,
        CASE WHEN j.JER_PRODUCTO = 'A_Balancer' THEN 0 ELSE a.Base_stck   END AS Base_stck,
        CASE WHEN j.JER_PRODUCTO = 'A_Balancer' THEN 0 ELSE a.Base_Nb     END AS Base_Nb
       
    FROM jerarquias j
    LEFT JOIN metrics_agg a
    ON a.JER_PRODUCTO = j.JER_PRODUCTO

    UNION ALL

    -- 2) TOTAL GENERAL (igual que en tu query original)
    SELECT
    'TOTAL GENERAL' AS JER_PRODUCTO,
    SUM(CASE WHEN scenario = 'Dwn200' AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_stck,
    SUM(CASE WHEN scenario = 'Dwn200' AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Dwn200_Nb,
    SUM(CASE WHEN scenario = 'Base'   AND NOT is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_stck,
    SUM(CASE WHEN scenario = 'Base'   AND     is_nb THEN (incexp + diff_incexp_yield) ELSE 0 END) AS Base_Nb
    FROM pro_pichincha_alquid_old."result"
    WHERE load_id IN (
    '{load_ids['cierre_base']}',
    '{load_ids['cierre_up']}',
    '{load_ids['cierre_dwn']}'
    )

    ORDER BY CASE JER_PRODUCTO
    WHEN 'BANCOS' THEN 1
    WHEN 'CARTERA MAYORISTA' THEN 2
    WHEN 'DEPOSITOS OTRAS EECC' THEN 3
    WHEN 'A_Balancer' THEN 4
    WHEN 'ACTIVO CORTO PLAZO' THEN 5
    WHEN 'ANTICIPO NOMINA' THEN 6
    WHEN 'CREDITOS' THEN 7
    WHEN 'DESCUBIERTOS' THEN 8
    WHEN 'PRESTAMOS' THEN 9
    WHEN 'TARJETAS CREDITO' THEN 10
    WHEN 'A_Bonos corporativos' THEN 11
    WHEN 'A_Bonos soberanos' THEN 12
    WHEN 'ACT DUDOSOS' THEN 13
    WHEN 'A_Resto no sensible' THEN 14
    WHEN 'A_Tesoreria_Admin.' THEN 15
    WHEN ' ' THEN 16
    WHEN 'P_BANCOS' THEN 17
    WHEN 'P_Balancer' THEN 18
    WHEN 'PLAZO' THEN 19
    WHEN 'CTAS NO REMUNERADAS' THEN 20
    WHEN 'CTAS REMUNERADAS' THEN 21
    WHEN 'P_Resto no sensible' THEN 22
    WHEN 'P_Ata' THEN 23
    WHEN 'P_Emisiones T2' THEN 24
    WHEN '  ' THEN 25
    WHEN 'FXSWAP_pago' THEN 26
    WHEN 'FXSWAP_recibo' THEN 27
    WHEN 'IRS_pago' THEN 28
    WHEN 'IRS_recibo' THEN 29
    WHEN 'Patrimonio Neto' THEN 30
    WHEN 'TOTAL GENERAL' THEN 31
    ELSE 999
    END;
    """
    
    # Se define nuestro data frame con los datos extrídos de la base de datos
    df_Detalle_MF = pd.read_sql(Informe_Detalle_MF, conn)

    # Nombre de la hoja donde quieres pegar los datos. Se mantiene 'detalleMF' según el original.
    # Si esta hoja también debe ser versionada por mes, la lógica debe cambiar.
    nombre_hoja_Detalle_MF = "detalleMF"

    # Acceder a la hoja
    ws_Detalle_MF = wb[nombre_hoja_Detalle_MF] # Asumo que 'detalleMF' está en 'Plantilla_COAP.xlsx'


    # Se definen las coordenadas donde queremos que se peguen los datos
    start_row_Detalle_MF = 5
    start_col_Detalle_MF = 4


    # Límite de filas donde se pegarán los datos (de la 5 a la 36, saltando la fila 8)
    end_row_Detalle_MF = 37


    # Escribir los datos en las coordenadas anteriores
    for i, row in enumerate(df_Detalle_MF.values, start=start_row_Detalle_MF):
        if i > end_row_Detalle_MF:
            break  # Detener la escritura si superamos la fila 37
        for j, value_Detalle_MF in enumerate(row[1:], start=start_col_Detalle_MF):
            ws_Detalle_MF.cell(row=i, column=j).value = value_Detalle_MF




    ###################################### PROCESO DE PEGADO DE INFORMES DE EFECTO BALANCE #####################################
    Informe_Efecto_Balance= f"""
        WITH metrics_agg AS (
        SELECT
            CASE
                WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') THEN '1.1 BANCOS'
                WHEN dim_1 IN ('A_Ata') THEN '1.2 CARTERA MAYORISTA'
                WHEN dim_1 IN ('A_Depositos cedidos') THEN '1.3 DEPOSITOS OTRAS EECC'
                WHEN dim_1 IN (
                    'A_ECORP_Otros_Activo Corto Plazo',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_EESP_Otros_Activo Corto Plazo',
                    'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                    'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                    'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                    /* ELIMINADO SEGÚN IMAGEN: 'A_EESP_Otros_Prest. Empresas_Gtia Personal', */
                    'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Personal', /* AÑADIDO SEGÚN IMAGEN */
                    'A_PIB_Otros_Descubiertos_Gtia Personal',
                    'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_PIB_Otros_Tarjetas Credito_Gtia Personal',
                    'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                    'A_POFI_Otros_Descubiertos_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                    'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                    'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_POFI_Otros_Prest. Origen_Gtia Aval',
                    'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                    'A_POFI_Otros_Tarjetas Credito_Gtia Personal',
                    'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Hipot.',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                    'A_SSCC_Otros_Descubiertos_Gtia Personal',
                    'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                    'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',
                    'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
                ) THEN '1.4 INVERSION CREDITICIA'
                WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') THEN 'RENTA FIJA'
                WHEN dim_1 = 'IRS_pago'      THEN 'IRS_pago'
                WHEN dim_1 = 'IRS_recibo'    THEN 'IRS_recibo'
                WHEN dim_1 = 'FXSWAP_pago'   THEN 'FXSWAP_pago'
                WHEN dim_1 = 'FXSWAP_recibo' THEN 'FXSWAP_recibo'
                WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') THEN '1.6 ACTIVOS DUDOSOS'
                WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') THEN '1.7 ACTIVOS NO SENSIBLE'
                WHEN dim_1 = 'P_Cuentas de otras EECC' THEN '2.1 DEPOSITOS OTRAS EECC'
                WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES'
                WHEN dim_1 IN (
                    'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                    'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
                    'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
                ) THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES'
                WHEN dim_1 = 'P_Resto no sensible' THEN '2.4 PASIVOS NO SENSIBLES'
                WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') THEN '2.5 CARTERA MAYORISTA'
                ELSE 'OTROS'
            END AS Jerarquia,
            COALESCE(ROUND(SUM(CASE WHEN scenario = 'Base'
                    AND load_id = '{load_ids['cierre_base_efecto_balance']}'
                    THEN market_value ELSE 0 END) / 1000000, 2), 0) AS Base,
            COALESCE(ROUND(SUM(CASE WHEN scenario = 'Up200'
                    AND load_id = '{load_ids['cierre_up_efecto_balance']}'
                    THEN market_value ELSE 0 END) / 1000000, 2), 0) AS Up200
        FROM pro_pichincha_alquid_old.metric
        WHERE load_id IN (
            '{load_ids['cierre_base_efecto_balance']}',
            '{load_ids['cierre_up_efecto_balance']}'
        )
        GROUP BY 1
    ),

    jerarquias AS (
        SELECT '1.1 BANCOS'                      AS Jerarquia UNION ALL
        SELECT '1.2 CARTERA MAYORISTA'                      UNION ALL
        SELECT '1.3 DEPOSITOS OTRAS EECC'                   UNION ALL
        SELECT '1.4 INVERSION CREDITICIA'                   UNION ALL
        SELECT 'RENTA FIJA'                                 UNION ALL
        SELECT 'FXSWAP_pago'                                UNION ALL
        SELECT 'FXSWAP_recibo'                              UNION ALL
        SELECT 'IRS_pago'                                   UNION ALL
        SELECT 'IRS_recibo'                                 UNION ALL
        SELECT '1.6 ACTIVOS DUDOSOS'                        UNION ALL
        SELECT '1.7 ACTIVOS NO SENSIBLE'                    UNION ALL
        SELECT '2.1 DEPOSITOS OTRAS EECC'                   UNION ALL
        SELECT '2.2 DEPOSITOS A PLAZO DE CLIENTES'          UNION ALL
        SELECT '2.3 DEPOSITOS A LA VISTA DE CLIENTES'       UNION ALL
        SELECT '2.4 PASIVOS NO SENSIBLES'                   UNION ALL
        SELECT '2.5 CARTERA MAYORISTA'
    )

    , total_renta_fija AS (
        SELECT
            COALESCE(SUM(Base), 0)  AS Base,
            COALESCE(SUM(Up200), 0) AS Up200
        FROM metrics_agg
        WHERE Jerarquia IN (
            'RENTA FIJA',
            'FXSWAP_pago',
            'FXSWAP_recibo',
            'IRS_pago',
            'IRS_recibo'
        )
    )

    -- filas por jerarquía (si no hay saldo, salen Base/Up200 = NULL)
    SELECT
        j.Jerarquia,
        m.Base,
        m.Up200
    FROM jerarquias j
    LEFT JOIN metrics_agg m
    ON m.Jerarquia = j.Jerarquia

    UNION ALL
    -- TOTAL GRUPO 1: activos + derivados (igual que el filtro original)
        SELECT
            'TOTAL GRUPO 1' AS Jerarquia,
        (
            COALESCE(SUM(Base), 0)
            + (SELECT Base FROM total_renta_fija)
        ) AS Base,
        (
            COALESCE(SUM(Up200), 0)
            + (SELECT Up200 FROM total_renta_fija)
        ) AS Up200
    FROM metrics_agg
    WHERE Jerarquia IN (
        '1.1 BANCOS',
        '1.2 CARTERA MAYORISTA',
        '1.3 DEPOSITOS OTRAS EECC',
        '1.4 INVERSION CREDITICIA',
        '1.6 ACTIVOS DUDOSOS'
    )

    UNION ALL
    -- TOTAL GRUPO 2: pasivos (como en tu filtro: 2.1, 2.2, 2.3 y 2.5)
    SELECT
        'TOTAL GRUPO 2' AS Jerarquia,
        SUM(Base) AS Base,
        SUM(Up200) AS Up200
    FROM metrics_agg
    WHERE Jerarquia IN (
        '2.1 DEPOSITOS OTRAS EECC',
        '2.2 DEPOSITOS A PLAZO DE CLIENTES',
        '2.3 DEPOSITOS A LA VISTA DE CLIENTES',
        '2.5 CARTERA MAYORISTA'
    )

    UNION ALL
    -- TOTAL RENTA FIJA: bonos + derivados (como tu TOTAL RENTA FIJA original)
    SELECT
        'TOTAL RENTA FIJA' AS Jerarquia,
        SUM(Base) AS Base,
        SUM(Up200) AS Up200
    FROM metrics_agg
    WHERE Jerarquia IN (
        'RENTA FIJA',
        'FXSWAP_pago',
        'FXSWAP_recibo',
        'IRS_pago',
        'IRS_recibo'
    )

    ORDER BY CASE Jerarquia
        WHEN 'TOTAL GRUPO 1' THEN 1
        WHEN '1.1 BANCOS' THEN 2
        WHEN '1.2 CARTERA MAYORISTA' THEN 3
        WHEN '1.3 DEPOSITOS OTRAS EECC' THEN 4
        WHEN '1.4 INVERSION CREDITICIA' THEN 5
        WHEN 'TOTAL RENTA FIJA' THEN 6
        WHEN 'RENTA FIJA' THEN 7
        WHEN 'FXSWAP_pago' THEN 8
        WHEN 'FXSWAP_recibo' THEN 9
        WHEN 'IRS_pago' THEN 10
        WHEN 'IRS_recibo' THEN 11
        WHEN '1.6 ACTIVOS DUDOSOS' THEN 12
        WHEN '1.7 ACTIVOS NO SENSIBLE' THEN 13
        WHEN 'TOTAL GRUPO 2' THEN 14
        WHEN '2.5 CARTERA MAYORISTA' THEN 15
        WHEN '2.1 DEPOSITOS OTRAS EECC' THEN 16
        WHEN '2.2 DEPOSITOS A PLAZO DE CLIENTES' THEN 17
        WHEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES' THEN 18
        WHEN '2.4 PASIVOS NO SENSIBLES' THEN 19
        ELSE 999
    END;
    """


    # Se define nuestro data frame con los datos extrídos de la base de datos
    df_Efecto_Balance = pd.read_sql(Informe_Efecto_Balance, conn)

    # Nombre de la hoja donde quieres pegar los datos.
    # El nombre original no incluye el mes. Si necesitas versionarlo, ajusta la lógica.
    nombre_hoja_Efecto_Balance = "VE_Efecto_balance"

    # Acceder a la hoja
    ws_Efecto_Balance = wb_1[nombre_hoja_Efecto_Balance]


    # Se definen las coordenadas donde queremos que se peguen los datos
    start_row_Efecto_Balance = 5
    start_col_Efecto_Balance = 2

    # Límite de filas donde se pegarán los datos (de la 5 a la 23)
    end_row_Efecto_Balance = 23

    # Escribir los datos en las coordenadas anteriores
    for i, row in enumerate(df_Efecto_Balance.values, start=start_row_Efecto_Balance):
        if i > end_row_Efecto_Balance:
            break  # Detener la escritura si superamos la fila 23
        for j, value_Efecto_Balance in enumerate(row[1:], start=start_col_Efecto_Balance):
            ws_Efecto_Balance.cell(row=i, column=j).value = value_Efecto_Balance




    ###################################### PROCESO DE PEGADO DE INFORMES DE EFECTO CURVA #####################################
    Informe_Efecto_Curva = f"""
        WITH metrics_agg AS (
        SELECT
            CASE
                WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') THEN '1.1 BANCOS'
                WHEN dim_1 = 'A_Ata' THEN '1.2 CARTERA MAYORISTA'
                WHEN dim_1 = 'A_Depositos cedidos' THEN '1.3 DEPOSITOS OTRAS EECC'
                WHEN dim_1 IN (
                    'A_ECORP_Otros_Activo Corto Plazo',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                    'A_EESP_Otros_Activo Corto Plazo',
                    'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                    'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                    'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                    'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                    'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                    'A_PIB_Otros_Descubiertos_Gtia Personal',
                    'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                    'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                    'A_POFI_Otros_Descubiertos_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                    'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                    'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_POFI_Otros_Prest. Origen_Gtia Aval',
                    'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                    'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                    'A_SSCC_Otros_Descubiertos_Gtia Personal',
                    'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                    'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
                ) THEN '1.4 INVERSION CREDITICIA'
                WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') THEN 'RENTA FIJA'
                WHEN dim_1 = 'IRS_pago'      THEN 'IRS_pago'
                WHEN dim_1 = 'IRS_recibo'    THEN 'IRS_recibo'
                WHEN dim_1 = 'FXSWAP_pago'   THEN 'FXSWAP_pago'
                WHEN dim_1 = 'FXSWAP_recibo' THEN 'FXSWAP_recibo'
                WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') THEN '1.6 ACTIVOS DUDOSOS'
                WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') THEN '1.7 ACTIVOS NO SENSIBLE'
                WHEN dim_1 = 'P_Cuentas de otras EECC' THEN '2.1 DEPOSITOS OTRAS EECC'
                WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES'
                WHEN dim_1 IN (
                    'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                    'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
                    'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
                ) THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES'
                WHEN dim_1 = 'P_Resto no sensible' THEN '2.4 PASIVOS NO SENSIBLES'
                WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') THEN '2.5 CARTERA MAYORISTA'
                ELSE 'OTROS'
            END AS Jerarquia,
            SUM(CASE WHEN scenario = 'Base'
                    AND load_id = '{load_ids['cierre_base_efecto_curva']}'
                    THEN market_value ELSE 0 END) / 1000000 AS Base,
            SUM(CASE WHEN scenario = 'Up200'
                    AND load_id = '{load_ids['cierre_up_efecto_curva']}'
                    THEN market_value ELSE 0 END) / 1000000 AS Up200
        FROM pro_pichincha_alquid_old.metric
        WHERE load_id IN (
            '{load_ids['cierre_base_efecto_curva']}',
            '{load_ids['cierre_up_efecto_curva']}'
        )
        GROUP BY 1
    ),

    jerarquias AS (
        SELECT '1.1 BANCOS'                      AS Jerarquia UNION ALL
        SELECT '1.2 CARTERA MAYORISTA'                      UNION ALL
        SELECT '1.3 DEPOSITOS OTRAS EECC'                   UNION ALL
        SELECT '1.4 INVERSION CREDITICIA'                   UNION ALL
        SELECT 'RENTA FIJA'                                 UNION ALL
        SELECT 'FXSWAP_pago'                                UNION ALL
        SELECT 'FXSWAP_recibo'                              UNION ALL
        SELECT 'IRS_pago'                                   UNION ALL
        SELECT 'IRS_recibo'                                 UNION ALL
        SELECT '1.6 ACTIVOS DUDOSOS'                        UNION ALL
        SELECT '1.7 ACTIVOS NO SENSIBLE'                    UNION ALL
        SELECT '2.1 DEPOSITOS OTRAS EECC'                   UNION ALL
        SELECT '2.2 DEPOSITOS A PLAZO DE CLIENTES'          UNION ALL
        SELECT '2.3 DEPOSITOS A LA VISTA DE CLIENTES'       UNION ALL
        SELECT '2.4 PASIVOS NO SENSIBLES'                   UNION ALL
        SELECT '2.5 CARTERA MAYORISTA'
    )
    , total_renta_fija AS (
        SELECT
            COALESCE(SUM(Base), 0)  AS Base,
            COALESCE(SUM(Up200), 0) AS Up200
        FROM metrics_agg
        WHERE Jerarquia IN (
            'RENTA FIJA',
            'FXSWAP_pago',
            'FXSWAP_recibo',
            'IRS_pago',
            'IRS_recibo'
        )
    )


    -- filas por jerarquía (si no hay saldo, salen Base/Up200 = NULL)
    SELECT
        j.Jerarquia,
        COALESCE(ROUND(m.Base, 2), 0),
        COALESCE(ROUND(m.Up200, 2), 0)
    FROM jerarquias j
    LEFT JOIN metrics_agg m
    ON m.Jerarquia = j.Jerarquia

    UNION ALL
    -- TOTAL GRUPO 1: SOLO BANCOS + CARTERA MAYORISTA + DEP. OTRAS EECC + INV. CREDITICIA + TOTAL RENTA FIJA + DUDOSOS
    SELECT
        'TOTAL GRUPO 1' AS Jerarquia,
        (
            COALESCE(SUM(Base), 0)
            + (SELECT Base FROM total_renta_fija)
        ) AS Base,
        (
            COALESCE(SUM(Up200), 0)
            + (SELECT Up200 FROM total_renta_fija)
        ) AS Up200
    FROM metrics_agg
    WHERE Jerarquia IN (
        '1.1 BANCOS',
        '1.2 CARTERA MAYORISTA',
        '1.3 DEPOSITOS OTRAS EECC',
        '1.4 INVERSION CREDITICIA',
        '1.6 ACTIVOS DUDOSOS'
    )


    UNION ALL
    -- TOTAL GRUPO 2: pasivos (sin 2.4, igual que tu filtro original)
    SELECT
        'TOTAL GRUPO 2' AS Jerarquia,
        SUM(Base) AS Base,
        SUM(Up200) AS Up200
    FROM metrics_agg
    WHERE Jerarquia IN (
        '2.1 DEPOSITOS OTRAS EECC',
        '2.2 DEPOSITOS A PLAZO DE CLIENTES',
        '2.3 DEPOSITOS A LA VISTA DE CLIENTES',
        '2.5 CARTERA MAYORISTA'
    )

    UNION ALL
    -- TOTAL RENTA FIJA: renta fija + derivados de tipos
    SELECT
        'TOTAL RENTA FIJA' AS Jerarquia,
        SUM(Base) AS Base,
        SUM(Up200) AS Up200
    FROM metrics_agg
    WHERE Jerarquia IN (
        'RENTA FIJA',
        'FXSWAP_pago',
        'FXSWAP_recibo',
        'IRS_pago',
        'IRS_recibo'
    )

    ORDER BY CASE Jerarquia
        WHEN 'TOTAL GRUPO 1' THEN 1
        WHEN '1.1 BANCOS' THEN 2
        WHEN '1.2 CARTERA MAYORISTA' THEN 3
        WHEN '1.3 DEPOSITOS OTRAS EECC' THEN 4
        WHEN '1.4 INVERSION CREDITICIA' THEN 5
        WHEN 'TOTAL RENTA FIJA' THEN 6
        WHEN 'RENTA FIJA' THEN 7
        WHEN 'FXSWAP_pago' THEN 8
        WHEN 'FXSWAP_recibo' THEN 9
        WHEN 'IRS_pago' THEN 10
        WHEN 'IRS_recibo' THEN 11
        WHEN '1.6 ACTIVOS DUDOSOS' THEN 12
        WHEN '1.7 ACTIVOS NO SENSIBLE' THEN 13
        WHEN 'TOTAL GRUPO 2' THEN 14
        WHEN '2.5 CARTERA MAYORISTA' THEN 15
        WHEN '2.1 DEPOSITOS OTRAS EECC' THEN 16
        WHEN '2.2 DEPOSITOS A PLAZO DE CLIENTES' THEN 17
        WHEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES' THEN 18
        WHEN '2.4 PASIVOS NO SENSIBLES' THEN 19
        ELSE 999
    END;
    """


    Informe_Efecto_Curva_Mes_Actual = f"""
        WITH metrics_agg AS (
        SELECT
            CASE
                WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') THEN '1.1 BANCOS'
                WHEN dim_1 IN ('A_Ata') THEN '1.2 CARTERA MAYORISTA'
                WHEN dim_1 IN ('A_Depositos cedidos') THEN '1.3 DEPOSITOS OTRAS EECC'
                WHEN dim_1 IN (
                    'A_ECORP_Otros_Activo Corto Plazo',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                    'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                    'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                    'A_EESP_Otros_Activo Corto Plazo',
                    'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                    'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                    'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                    'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                    'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                    'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                    'A_PIB_Otros_Descubiertos_Gtia Personal',
                    'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                    'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                    'A_POFI_Otros_Descubiertos_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                    'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                    'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                    'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                    'A_POFI_Otros_Prest. Origen_Gtia Aval',
                    'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                    'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                    'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                    'A_SSCC_Otros_Descubiertos_Gtia Personal',
                    'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                    'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                    'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                    'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                    'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
                ) THEN '1.4 INVERSION CREDITICIA'
                WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') THEN 'RENTA FIJA'
                WHEN dim_1 IN ('IRS_pago') THEN 'IRS_pago'
                WHEN dim_1 IN ('IRS_recibo') THEN 'IRS_recibo'
                WHEN dim_1 IN ('FXSWAP_pago') THEN 'FXSWAP_pago'
                WHEN dim_1 IN ('FXSWAP_recibo') THEN 'FXSWAP_recibo'
                WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') THEN '1.6 ACTIVOS DUDOSOS'
                WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') THEN '1.7 ACTIVOS NO SENSIBLE'
                WHEN dim_1 IN ('P_Cuentas de otras EECC') THEN '2.1 DEPOSITOS OTRAS EECC'
                WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES'
                WHEN dim_1 IN (
                    'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                    'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
                    'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
                ) THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES'
                WHEN dim_1 IN ('P_Resto no sensible') THEN '2.4 PASIVOS NO SENSIBLES'
                WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') THEN '2.5 CARTERA MAYORISTA'
                ELSE 'OTROS'
            END AS Jerarquia,
            SUM(CASE WHEN scenario = 'Base'  AND load_id = '{load_ids['cierre_base']}' THEN market_value ELSE 0 END) / 1000000 AS Base,
            SUM(CASE WHEN scenario = 'Up200' AND load_id = '{load_ids['cierre_up']}'   THEN market_value ELSE 0 END) / 1000000 AS Up200
        FROM pro_pichincha_alquid_old.metric
        GROUP BY 1
    ),

    jerarquias AS (
        SELECT '1.1 BANCOS'                     AS Jerarquia UNION ALL
        SELECT '1.2 CARTERA MAYORISTA'                      UNION ALL
        SELECT '1.3 DEPOSITOS OTRAS EECC'                   UNION ALL
        SELECT '1.4 INVERSION CREDITICIA'                   UNION ALL
        SELECT 'RENTA FIJA'                                 UNION ALL
        SELECT 'FXSWAP_pago'                                UNION ALL
        SELECT 'FXSWAP_recibo'                              UNION ALL
        SELECT 'IRS_pago'                                   UNION ALL
        SELECT 'IRS_recibo'                                 UNION ALL
        SELECT '1.6 ACTIVOS DUDOSOS'                        UNION ALL
        SELECT '1.7 ACTIVOS NO SENSIBLE'                    UNION ALL
        SELECT '2.1 DEPOSITOS OTRAS EECC'                   UNION ALL
        SELECT '2.2 DEPOSITOS A PLAZO DE CLIENTES'          UNION ALL
        SELECT '2.3 DEPOSITOS A LA VISTA DE CLIENTES'       UNION ALL
        SELECT '2.4 PASIVOS NO SENSIBLES'                   UNION ALL
        SELECT '2.5 CARTERA MAYORISTA'
    )
    , total_renta_fija AS (
        SELECT
            COALESCE(SUM(Base), 0)  AS Base,
            COALESCE(SUM(Up200), 0) AS Up200
        FROM metrics_agg
        WHERE Jerarquia IN (
            'RENTA FIJA',
            'FXSWAP_pago',
            'FXSWAP_recibo',
            'IRS_pago',
            'IRS_recibo'
        )
    )

    SELECT
        j.Jerarquia,
        m.Base,
        m.Up200
    FROM jerarquias j
    LEFT JOIN metrics_agg m
        ON j.Jerarquia = m.Jerarquia
    
    UNION ALL
    -- TOTAL GRUPO 1: SOLO BANCOS + CARTERA MAYORISTA + DEP. OTRAS EECC + INV. CREDITICIA + TOTAL RENTA FIJA + DUDOSOS
    SELECT
        'TOTAL GRUPO 1' AS Jerarquia,
        (
            COALESCE(SUM(Base), 0)
            + (SELECT Base FROM total_renta_fija)
        ) AS Base,
        (
            COALESCE(SUM(Up200), 0)
            + (SELECT Up200 FROM total_renta_fija)
        ) AS Up200
    FROM metrics_agg
    WHERE Jerarquia IN (
        '1.1 BANCOS',
        '1.2 CARTERA MAYORISTA',
        '1.3 DEPOSITOS OTRAS EECC',
        '1.4 INVERSION CREDITICIA',
        '1.6 ACTIVOS DUDOSOS'
    )


    UNION ALL

    -- TOTAL GRUPO 2: pasivos (sin 2.4, igual que tu filtro original)
    SELECT
        'TOTAL GRUPO 2' AS Jerarquia,
        SUM(Base) AS Base,
        SUM(Up200) AS Up200
    FROM metrics_agg
    WHERE Jerarquia IN (
        '2.1 DEPOSITOS OTRAS EECC',
        '2.2 DEPOSITOS A PLAZO DE CLIENTES',
        '2.3 DEPOSITOS A LA VISTA DE CLIENTES',
        '2.5 CARTERA MAYORISTA'
    )

    UNION ALL

    -- TOTAL RENTA FIJA: renta fija + derivados de tipos
    SELECT
        'TOTAL RENTA FIJA' AS Jerarquia,
        SUM(Base) AS Base,
        SUM(Up200) AS Up200
    FROM metrics_agg
    WHERE Jerarquia IN (
        'RENTA FIJA',
        'FXSWAP_pago',
        'FXSWAP_recibo',
        'IRS_pago',
        'IRS_recibo'
    )
    ORDER BY CASE Jerarquia
        WHEN 'TOTAL GRUPO 1'                 THEN 1
        WHEN '1.1 BANCOS'                    THEN 2
        WHEN '1.2 CARTERA MAYORISTA'         THEN 3
        WHEN '1.3 DEPOSITOS OTRAS EECC'      THEN 4
        WHEN '1.4 INVERSION CREDITICIA'      THEN 5
        WHEN 'TOTAL RENTA FIJA'              THEN 6
        WHEN 'RENTA FIJA'                    THEN 7
        WHEN 'FXSWAP_pago'                   THEN 8
        WHEN 'FXSWAP_recibo'                 THEN 9
        WHEN 'IRS_pago'                      THEN 10
        WHEN 'IRS_recibo'                    THEN 11
        WHEN '1.6 ACTIVOS DUDOSOS'           THEN 12
        WHEN '1.7 ACTIVOS NO SENSIBLE'       THEN 13
        WHEN 'TOTAL GRUPO 2'                 THEN 14
        WHEN '2.5 CARTERA MAYORISTA'         THEN 15
        WHEN '2.1 DEPOSITOS OTRAS EECC'      THEN 16
        WHEN '2.2 DEPOSITOS A PLAZO DE CLIENTES' THEN 17
        WHEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES' THEN 18
        WHEN '2.4 PASIVOS NO SENSIBLES'      THEN 19
        ELSE 999
    END;
    """

    # Se define nuestro data frame con los datos extrídos de la base de datos
    df_Efecto_Curva  = pd.read_sql(Informe_Efecto_Curva , conn)

    df_Efecto_Curva_Mes_Actual  = pd.read_sql(Informe_Efecto_Curva_Mes_Actual , conn)

    # Nombre de la hoja donde quieres pegar los datos.
    # El nombre original no incluye el mes. Si necesitas versionarlo, ajusta la lógica.
    nombre_hoja_Efecto_Curva  = "VE_Efecto_curva"

    nombre_hoja_Efecto_Curva_Mes_Actual  = "VE_mes actual"

    # Acceder a la hoja
    ws_Efecto_Curva  = wb_1[nombre_hoja_Efecto_Curva ]

    ws_Efecto_Curva_Mes_Actual  = wb_1[nombre_hoja_Efecto_Curva_Mes_Actual ]


    # Se definen las coordenadas donde queremos que se peguen los datos
    start_row_Efecto_Curva  = 5
    start_col_Efecto_Curva  = 2

    start_row_Efecto_Curva_Mes_Actual  = 5
    start_col_Efecto_Curva_Mes_Actual  = 2

    # Límite de filas donde se pegarán los datos (de la 5 a la 23)
    end_row_Efecto_Curva  = 23

    end_row_Efecto_Curva_Mes_Actual  = 23

    # Escribir los datos en las coordenadas anteriores
    for i, row in enumerate(df_Efecto_Curva .values, start=start_row_Efecto_Curva ):
        if i > end_row_Efecto_Curva :
            break  # Detener la escritura si superamos la fila 23
        for j, value_Efecto_Curva  in enumerate(row[1:], start=start_col_Efecto_Curva ):
            ws_Efecto_Curva .cell(row=i, column=j).value = value_Efecto_Curva

    for i, row in enumerate(df_Efecto_Curva_Mes_Actual .values, start=start_row_Efecto_Curva_Mes_Actual ):
        if i > end_row_Efecto_Curva_Mes_Actual :
            break  # Detener la escritura si superamos la fila 23
        for j, value_Efecto_Curva_Mes_Actual  in enumerate(row[1:], start=start_col_Efecto_Curva_Mes_Actual ):
            ws_Efecto_Curva_Mes_Actual .cell(row=i, column=j).value = value_Efecto_Curva_Mes_Actual


    ###################################### PROCESO DE PEGADO DE INFORMES DE DATOS MEDIOS #####################################
    # La fecha para la query SQL se toma del diccionario mes_anterior_dic
    fecha_sql_mes_anterior = mes_anterior_dic['mes_anterior']

    Informe_Datos_Medios_01= f"""WITH jerarquias AS (
        SELECT '1.1 BANCOS.F'                      AS Jerarquia,  1 AS ord UNION ALL
        SELECT '1.2 CARTERA MAYORISTA.F'                     ,  2 UNION ALL
        SELECT '1.3 DEPOSITOS OTRAS EECC.F'                  ,  3 UNION ALL
        SELECT '1.4 INVERSION CREDITICIA.F'                  ,  4 UNION ALL
        SELECT '1.5 RENTA FIJA.F'                            ,  5 UNION ALL
        SELECT '1.6 ACTIVOS DUDOSOS.F'                       ,  6 UNION ALL
        SELECT '1.7 ACTIVOS NO SENSIBLE.F'                   ,  7 UNION ALL
        SELECT '1.1 BANCOS.V'                                ,  8 UNION ALL
        SELECT '1.4 INVERSION CREDITICIA.V'                  , 9 UNION ALL
        SELECT '1.5 RENTA FIJA.V'                            , 10 UNION ALL
        SELECT '1.6 ACTIVOS DUDOSOS.V'                       , 11 UNION ALL
        SELECT '2.1 DEPOSITOS OTRAS EECC.F'                  , 12 UNION ALL
        SELECT '2.2 DEPOSITOS A PLAZO DE CLIENTES.F'         , 13 UNION ALL
        SELECT '2.3 DEPOSITOS A LA VISTA DE CLIENTES'      , 14 UNION ALL
        SELECT '2.4 PASIVOS NO SENSIBLES.F'                  , 15 UNION ALL
        SELECT '2.5 CARTERA MAYORISTA.F'                     , 16
    ),
    metrics_agg AS (
        SELECT
            CASE
                WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') AND NOT is_floating THEN '1.1 BANCOS.F'
                WHEN dim_1 IN ('A_Ata') AND NOT is_floating THEN '1.2 CARTERA MAYORISTA.F'
                WHEN dim_1 IN ('A_Depositos cedidos') AND NOT is_floating THEN '1.3 DEPOSITOS OTRAS EECC.F'
                WHEN dim_1 IN (
                        'A_ECORP_Otros_Activo Corto Plazo',
                        'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                        'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                        'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                        'A_EESP_Otros_Activo Corto Plazo',
                        'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                        'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                        'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                        'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                        'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                        'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                        'A_PIB_Otros_Descubiertos_Gtia Personal',
                        'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                        'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                        'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                        'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                        'A_POFI_Otros_Descubiertos_Gtia Personal',
                        'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                        'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                        'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                        'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                        'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                        'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                        'A_POFI_Otros_Prest. Origen_Gtia Aval',
                        'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                        'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                        'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                        'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                        'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                        'A_SSCC_Otros_Descubiertos_Gtia Personal',
                        'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                        'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                        'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                        'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                        'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                        'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
                ) AND NOT is_floating THEN '1.4 INVERSION CREDITICIA.F'
                WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') AND NOT is_floating THEN '1.5 RENTA FIJA.F'
                WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') AND NOT is_floating THEN '1.6 ACTIVOS DUDOSOS.F'
                WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') AND NOT is_floating THEN '1.7 ACTIVOS NO SENSIBLE.F'
                WHEN dim_1 IN ('P_Cuentas de otras EECC') AND NOT is_floating THEN '2.1 DEPOSITOS OTRAS EECC.F'
                WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') AND NOT is_floating THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES.F'
                WHEN dim_1 IN (
                    'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                    'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
                    'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
                ) AND NOT is_floating THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES'
                WHEN dim_1 IN ('P_Resto no sensible') AND NOT is_floating THEN '2.4 PASIVOS NO SENSIBLES.F'
                WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') AND NOT is_floating THEN '2.5 CARTERA MAYORISTA.F'
                WHEN dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo') AND NOT is_floating THEN '3.1 DERIVADOS.F'
                WHEN dim_1 IN ('Patrimonio Neto') AND NOT is_floating THEN '4.1 PN.F'
                WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') AND is_floating THEN '1.1 BANCOS.V'
                WHEN dim_1 IN ('A_Ata') AND is_floating THEN '1.2 CARTERA MAYORISTA.V'
                WHEN dim_1 IN ('A_Depositos cedidos') AND is_floating THEN '1.3 DEPOSITOS OTRAS EECC.V'
                WHEN dim_1 IN (
                    'A_ECORP_Otros_Activo Corto Plazo',
                        'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                        'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                        'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                        'A_EESP_Otros_Activo Corto Plazo',
                        'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                        'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                        'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                        'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                        'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                        'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                        'A_PIB_Otros_Descubiertos_Gtia Personal',
                        'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                        'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                        'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                        'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                        'A_POFI_Otros_Descubiertos_Gtia Personal',
                        'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                        'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                        'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                        'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                        'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                        'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                        'A_POFI_Otros_Prest. Origen_Gtia Aval',
                        'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                        'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                        'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                        'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                        'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                        'A_SSCC_Otros_Descubiertos_Gtia Personal',
                        'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                        'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                        'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                        'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                        'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                        'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
                ) AND is_floating THEN '1.4 INVERSION CREDITICIA.V'
                WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') AND is_floating THEN '1.5 RENTA FIJA.V'
                WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') AND is_floating THEN '1.6 ACTIVOS DUDOSOS.V'
                WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') AND is_floating THEN '1.7 ACTIVOS NO SENSIBLE.V'
                WHEN dim_1 IN ('P_Cuentas de otras EECC') AND is_floating THEN '2.1 DEPOSITOS OTRAS EECC.V'
                WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') AND is_floating THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES.V'
                WHEN dim_1 IN (
                    'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                    'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
                    'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
                ) AND is_floating THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES'
                WHEN dim_1 IN ('P_Resto no sensible') AND is_floating THEN '2.4 PASIVOS NO SENSIBLES.V'
                WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') AND is_floating THEN '2.5 CARTERA MAYORISTA.V'
                WHEN dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo') AND is_floating THEN '3.1 DERIVADOS.V'
                WHEN dim_1 IN ('Patrimonio Neto') AND is_floating THEN '4.1 PN.V'
                ELSE 'OTROS'
            END AS Jerarquia,
        SUM(notional_today) as notional,
        SUM((date_diff('day', balance_date, matdate) * Notional_today)) / NULLIF(SUM(Notional_today),0) / 30 AS vto_medio,
        SUM(coupon_today*Notional_today)/NULLIF(SUM(Notional_today),0) as cupon_medio,
        SUM(poolrateadj*notional_today)/NULLIF(SUM(notional_today),0) as spread_medio,

        -- Original: SUM((date_diff('day', Reprdate, matdate) * Notional_today)) / NULLIF(SUM(Notional_today),0) / 30 AS reprecio_medio,
        
        SUM((date_diff('day', balance_date, Reprdate) * Notional_today)) / NULLIF(SUM(Notional_today),0) / 30 AS reprecio_medio,
        
        SUM (floor_today*notional_today)/NULLIF(SUM(notional_today),0) as Floor_medio,
        (SUM (effective_duration*notional_today)/NULLIF(SUM (Notional_today),0))/100 as Dur_Efec_Pond,
        SUM (macaulay_duration*notional_today)/NULLIF(SUM(notional_today),0) as Dur_Mac_pond
        FROM pro_pichincha_alquid_old.metric
        WHERE load_id IN (
            '{load_ids['cierre_base']}', 
            '{load_ids['cierre_up']}', 
            '{load_ids['cierre_dwn']}'
        )
        AND scene_code = 'Base'
        GROUP BY 1
    )
    SELECT
        j.Jerarquia,
        a.notional,
        a.vto_medio,
        a.cupon_medio,
        a.spread_medio,
        a.reprecio_medio,
        a.floor_medio,
        a.dur_efec_pond,
        a.dur_mac_pond
    FROM jerarquias j
    LEFT JOIN metrics_agg a
        ON a.Jerarquia = j.Jerarquia
    ORDER BY j.ord;
    """

    Informe_Datos_Medios_01_01= f""" WITH AllJerarquias AS (
        SELECT '1.1 BANCOS.F' AS Jerarquia UNION ALL
        SELECT '1.2 CARTERA MAYORISTA.F' UNION ALL
        SELECT '1.3 DEPOSITOS OTRAS EECC.F' UNION ALL
        SELECT '1.4 INVERSION CREDITICIA.F' UNION ALL
        SELECT '1.5 RENTA FIJA.F' UNION ALL
        SELECT '1.6 ACTIVOS DUDOSOS.F' UNION ALL
        SELECT '1.7 ACTIVOS NO SENSIBLE.F' UNION ALL
        SELECT '2.1 DEPOSITOS OTRAS EECC.F' UNION ALL
        SELECT '2.2 DEPOSITOS A PLAZO DE CLIENTES.F' UNION ALL
        SELECT '2.3 DEPOSITOS A LA VISTA DE CLIENTES.F' UNION ALL
        SELECT '2.4 PASIVOS NO SENSIBLES.F' UNION ALL
        SELECT '2.5 CARTERA MAYORISTA.F' UNION ALL
        SELECT '3.1 DERIVADOS.F' UNION ALL
        SELECT '4.1 PN.F' UNION ALL
        SELECT '1.1 BANCOS.V' UNION ALL
        SELECT '1.4 INVERSION CREDITICIA.V' UNION ALL
        SELECT '1.5 RENTA FIJA.V' UNION ALL
        SELECT '1.6 ACTIVOS DUDOSOS.V' UNION ALL
        SELECT '2.1 DEPOSITOS OTRAS EECC.V' UNION ALL
        SELECT '2.2 DEPOSITOS A PLAZO DE CLIENTES.V' UNION ALL
        SELECT '2.4 PASIVOS NO SENSIBLES.V' UNION ALL
        SELECT '2.5 CARTERA MAYORISTA.V' UNION ALL
        SELECT '3.1 DERIVADOS.V' UNION ALL
        SELECT '4.1 PN.V' UNION ALL
        SELECT 'OTROS' -- Incluye la categoría 'OTROS' si es posible que aparezca
    ),
    AggregatedData AS (
    SELECT
        CASE
        WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') AND NOT is_floating THEN '1.1 BANCOS.F'
        WHEN dim_1 IN ('A_Ata') AND NOT is_floating THEN '1.2 CARTERA MAYORISTA.F'
        WHEN dim_1 IN ('A_Depositos cedidos') AND NOT is_floating THEN '1.3 DEPOSITOS OTRAS EECC.F'
        WHEN dim_1 IN (
            'A_ECORP_Otros_Activo Corto Plazo',
                        'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                        'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                        'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                        'A_EESP_Otros_Activo Corto Plazo',
                        'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                        'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                        'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                        'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                        'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                        'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                        'A_PIB_Otros_Descubiertos_Gtia Personal',
                        'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                        'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                        'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                        'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                        'A_POFI_Otros_Descubiertos_Gtia Personal',
                        'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                        'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                        'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                        'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                        'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                        'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                        'A_POFI_Otros_Prest. Origen_Gtia Aval',
                        'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                        'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                        'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                        'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                        'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                        'A_SSCC_Otros_Descubiertos_Gtia Personal',
                        'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                        'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                        'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                        'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                        'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                        'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
        ) AND NOT is_floating THEN '1.4 INVERSION CREDITICIA.F'
        WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') AND NOT is_floating THEN '1.5 RENTA FIJA.F'
        WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') AND NOT is_floating THEN '1.6 ACTIVOS DUDOSOS.F'
        WHEN dim_1 IN ('A_Resto no sensible', 'A_Tesoreria_Admin.') AND NOT is_floating THEN '1.7 ACTIVOS NO SENSIBLE.F'
        WHEN dim_1 IN ('P_Cuentas de otras EECC') AND NOT is_floating THEN '2.1 DEPOSITOS OTRAS EECC.F'
        WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') AND NOT is_floating THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES.F'
        WHEN dim_1 IN (
            'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
            'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
            'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
        ) THEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES.F'
        WHEN dim_1 IN ('P_Resto no sensible') AND NOT is_floating THEN '2.4 PASIVOS NO SENSIBLES.F'
        WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') AND NOT is_floating THEN '2.5 CARTERA MAYORISTA.F'
        WHEN dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo') AND NOT is_floating THEN '3.1 DERIVADOS.F'
        WHEN dim_1 IN ('Patrimonio Neto') AND NOT is_floating THEN '4.1 PN.F'
        WHEN dim_1 IN ('A_Cuentas en otras EECC', 'A_Tesoreria') AND is_floating THEN '1.1 BANCOS.V'
        WHEN dim_1 IN (
            'A_ECORP_Otros_Activo Corto Plazo',
                        'A_ECORP_Otros_Cred. Empresa_Gtia Aval',
                        'A_ECORP_Otros_Cred. Empresa_Gtia Personal',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Aval',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Personal',
                        'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
                        'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal',

                        'A_EESP_Otros_Activo Corto Plazo',
                        'A_EESP_Otros_Cred. Empresa_Gtia Aval',
                        'A_EESP_Otros_Cred. Empresa_Gtia Personal',
                        'A_EESP_Otros_Prest. Empresas_Gtia Aval',
                        'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_EESP_Otros_Prest. Promotor_Gtia Aval',
                        'A_EESP_Otros_Prest. Promotor_Gtia Hipot.',
                        'A_EESP_Otros_Prest. Promotor_Gtia Personal',

                        'A_PIB_Otros_Descubiertos_Gtia Personal',
                        'A_PIB_Otros_Prest. Consumo_Gtia Personal',
                        'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_PIB_Otros_Prest. Hipotecas_Gtia Personal',
                        'A_PIB_Otros_Tarjetas Credito_Gtia Personal',

                        'A_POFI_Otros_Ant. Nomina_Gtia Personal',
                        'A_POFI_Otros_Descubiertos_Gtia Personal',
                        'A_POFI_Otros_Prest. Consumo_Gtia Personal',
                        'A_POFI_Otros_Prest. Consumo_Gtia Prenda',
                        'A_POFI_Otros_Prest. Empresas_Gtia Aval',
                        'A_POFI_Otros_Prest. Empresas_Gtia Hipot.',
                        'A_POFI_Otros_Prest. Empresas_Gtia Personal',
                        'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
                        'A_POFI_Otros_Prest. Origen_Gtia Aval',
                        'A_POFI_Otros_Prest. Origen_Gtia Hipot.',
                        'A_POFI_Otros_Tarjetas Credito_Gtia Personal',

                        'A_SSCC_BKIA_Prest. Consumo_Gtia Personal',
                        'A_SSCC_Otros_Cred. Empresa_Gtia Aval',
                        'A_SSCC_Otros_Cred. Empresa_Gtia Personal',
                        'A_SSCC_Otros_Descubiertos_Gtia Personal',
                        'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
                        'A_SSCC_Otros_Prest. Empresas_Gtia Aval',
                        'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.',
                        'A_SSCC_Otros_Prest. Empresas_Gtia Personal',
                        'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.',
                        'A_SSCC_Otros_Tarjetas Credito_Gtia Personal',

                        'A_SSCCRSG_Otros_Tarjetas Credito_Gtia Personal'
        ) AND is_floating THEN '1.4 INVERSION CREDITICIA.V'
        WHEN dim_1 IN ('A_Bonos corporativos', 'A_Bonos soberanos') AND is_floating THEN '1.5 RENTA FIJA.V'
        WHEN dim_1 IN ('A_Cred._Dudoso', 'A_Efectos_Dudosos', 'A_Prest._Dudoso') AND is_floating THEN '1.6 ACTIVOS DUDOSOS.V'
        WHEN dim_1 IN ('P_Cuentas de otras EECC') AND is_floating THEN '2.1 DEPOSITOS OTRAS EECC.V'
        WHEN dim_1 IN ('P_ECORP_P. Plazo', 'P_PIB_P. Plazo', 'P_POFI_P. Plazo', 'P_SSCC_P. Plazo') AND is_floating THEN '2.2 DEPOSITOS A PLAZO DE CLIENTES.V'
        WHEN dim_1 IN ('P_Resto no sensible') AND is_floating THEN '2.4 PASIVOS NO SENSIBLES.V'
        WHEN dim_1 IN ('P_Ata', 'P_Emisiones T2') AND is_floating THEN '2.5 CARTERA MAYORISTA.V'
        WHEN dim_1 IN ('FXSWAP_pago', 'FXSWAP_recibo', 'IRS_pago', 'IRS_recibo') AND is_floating THEN '3.1 DERIVADOS.V'
        WHEN dim_1 IN ('Patrimonio Neto') AND is_floating THEN '4.1 PN.V'
        ELSE 'OTROS'
        END AS Jerarquia,
            SUM(notional_today) AS notional,
        SUM((date_diff('day', balance_date, matdate) * Notional_today)) / NULLIF(SUM(Notional_today), 0) / 30 AS vto_medio,
        SUM(coupon_today*Notional_today)/NULLIF(SUM(Notional_today), 0) AS cupon_medio,
        SUM(poolrateadj*notional_today)/NULLIF(SUM(notional_today), 0) AS spread_medio,
        SUM((date_diff('day', balance_date, Reprdate) * Notional_today)) / NULLIF(SUM(Notional_today), 0) / 30 AS reprecio_medio,
        SUM (floor_today*notional_today)/NULLIF(SUM(notional_today), 0) AS Floor_medio,
        (SUM (effective_duration*notional_today)/NULLIF(SUM (Notional_today), 0))/100 AS Dur_Efec_Pond,
        SUM (macaulay_duration*notional_today)/NULLIF(SUM(notional_today), 0) AS Dur_Mac_pond
    FROM pro_pichincha_alquid_old.metric
        WHERE load_id IN ('{load_ids['cierre_base']}', '{load_ids['cierre_up']}', '{load_ids['cierre_dwn']}')
        AND scene_code = 'Base'
        AND origin_date >= DATE '{fecha_sql_mes_anterior}'
    GROUP BY 1
    )
    SELECT
        AJ.Jerarquia,
        COALESCE(AD.notional, 0) AS notional,
        COALESCE(AD.vto_medio, 0) AS vto_medio,
        COALESCE(AD.cupon_medio, 0) AS cupon_medio,
        COALESCE(AD.spread_medio, 0) AS spread_medio,
        COALESCE(AD.reprecio_medio, 0) AS reprecio_medio,
        COALESCE(AD.Floor_medio, 0) AS Floor_medio,
        COALESCE(AD.Dur_Efec_Pond, 0) AS Dur_Efec_Pond,
        COALESCE(AD.Dur_Mac_pond, 0) AS Dur_Mac_pond
    FROM AllJerarquias AS AJ
    LEFT JOIN AggregatedData AS AD ON AJ.Jerarquia = AD.Jerarquia
    ORDER BY CASE AJ.Jerarquia
    WHEN '1.1 BANCOS.F' THEN 1
    WHEN '1.2 CARTERA MAYORISTA.F' THEN 2
    WHEN '1.3 DEPOSITOS OTRAS EECC.F' THEN 3
    WHEN '1.4 INVERSION CREDITICIA.F' THEN 4
    WHEN '1.5 RENTA FIJA.F' THEN 5
    WHEN '1.6 ACTIVOS DUDOSOS.F' THEN 6
    WHEN '1.7 ACTIVOS NO SENSIBLE.F' THEN 7
    WHEN '1.1 BANCOS.V' THEN 8
    WHEN '1.4 INVERSION CREDITICIA.V' THEN 9
    WHEN '1.5 RENTA FIJA.V' THEN 10
    WHEN '1.6 ACTIVOS DUDOSOS.V' THEN 11
    WHEN '2.1 DEPOSITOS OTRAS EECC.F' THEN 12
    WHEN '2.2 DEPOSITOS A PLAZO DE CLIENTES.F' THEN 13
    WHEN '2.3 DEPOSITOS A LA VISTA DE CLIENTES.F' THEN 14
    WHEN '2.4 PASIVOS NO SENSIBLES.F' THEN 15
    WHEN '2.5 CARTERA MAYORISTA.F' THEN 16
    WHEN '2.1 DEPOSITOS OTRAS EECC.V' THEN 17
    WHEN '2.2 DEPOSITOS A PLAZO DE CLIENTES.V' THEN 18
    WHEN '2.4 PASIVOS NO SENSIBLES.V' THEN 19
    WHEN '2.5 CARTERA MAYORISTA.V' THEN 20
    ELSE 999
    END; """

    Informe_Datos_Medios_02= f""" select
        CASE
            WHEN dim_1 LIKE '%Activo Corto Plazo%' and not is_floating THEN 'ACTIVO CORTO PLAZO.F'
            WHEN dim_1 LIKE '%Ant. Nomina%' and not is_floating THEN 'ANTICIPO NOMINA.F'
            WHEN dim_1 in ('A_ECORP_Otros_Cred. Empresa_Gtia Aval', 'A_ECORP_Otros_Cred. Empresa_Gtia Personal', 'A_EESP_Otros_Cred. Empresa_Gtia Aval',
            'A_EESP_Otros_Cred. Empresa_Gtia Personal', 'A_SSCC_Otros_Cred. Empresa_Gtia Aval', 'A_SSCC_Otros_Cred. Empresa_Gtia Hipot.', 'A_SSCC_Otros_Cred. Empresa_Gtia Personal') 
            and not is_floating THEN 'CREDITOS.F'
            WHEN dim_1 LIKE '%Descubiertos%'  and not is_floating THEN 'DESCUBIERTOS.F'
            WHEN dim_1 IN (
        'A_ECORP_Otros_Prest. Empresas_Gtia Aval', 'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.', 'A_ECORP_Otros_Prest. Empresas_Gtia Personal', 'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
        'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal', 'A_EESP_Otros_Prest. Empresas_Gtia Aval', 'A_EESP_Otros_Prest. Empresas_Gtia Personal',
        'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_EESP_Otros_Prest. Promotor_Gtia Hipot.', 'A_PIB_Otros_Prest. Consumo_Gtia Personal', 'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
        'A_PIB_Otros_Prest. Hipotecas_Gtia Personal', 'A_POFI_Otros_Prest. Consumo_Gtia Personal', 'A_POFI_Otros_Prest. Consumo_Gtia Prenda', 'A_POFI_Otros_Prest. Empresas_Gtia Aval',
        'A_POFI_Otros_Prest. Empresas_Gtia Hipot.', 'A_POFI_Otros_Prest. Empresas_Gtia Personal', 'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
        'A_POFI_Otros_Prest. Origen_Gtia Aval', 'A_POFI_Otros_Prest. Origen_Gtia Hipot.', 'A_SSCC_BKIA_Prest. Consumo_Gtia Personal', 'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
        'A_SSCC_Otros_Prest. Empresas_Gtia Aval', 'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.','A_SSCC_Otros_Prest. Empresas_Gtia Personal', 'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.') 
        and not is_floating THEN 'PRESTAMOS.F'
            WHEN dim_1 LIKE '%Tarjetas Credito%' and not is_floating THEN 'TARJETAS CREDITO.F'
            WHEN dim_1 LIKE '%Activo Corto Plazo%' and is_floating THEN 'ACTIVO CORTO PLAZO.V'
            WHEN dim_1 LIKE '%Ant. Nomina%' and is_floating THEN 'ANTICIPO NOMINA.V'
            WHEN dim_1 in ('A_ECORP_Otros_Cred. Empresa_Gtia Aval', 'A_ECORP_Otros_Cred. Empresa_Gtia Personal', 'A_EESP_Otros_Cred. Empresa_Gtia Aval',
            'A_EESP_Otros_Cred. Empresa_Gtia Personal', 'A_SSCC_Otros_Cred. Empresa_Gtia Aval', 'A_SSCC_Otros_Cred. Empresa_Gtia Hipot.', 'A_SSCC_Otros_Cred. Empresa_Gtia Personal')
            and is_floating THEN 'CREDITOS.V'
            WHEN dim_1 LIKE '%Descubiertos%' and is_floating THEN 'DESCUBIERTOS.V'
            WHEN dim_1 IN (
        'A_ECORP_Otros_Prest. Empresas_Gtia Aval', 'A_ECORP_Otros_Prest. Empresas_Gtia Hipot.', 'A_ECORP_Otros_Prest. Empresas_Gtia Personal', 'A_ECORP_Otros_Prest. Empresas_Gtia Prenda',
        'A_ECORP_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_ECORP_Otros_Prest. Hipotecas_Gtia Personal', 'A_EESP_Otros_Prest. Empresas_Gtia Aval', 'A_EESP_Otros_Prest. Empresas_Gtia Personal',
        'A_EESP_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_EESP_Otros_Prest. Promotor_Gtia Hipot.', 'A_PIB_Otros_Prest. Consumo_Gtia Personal', 'A_PIB_Otros_Prest. Hipotecas_Gtia Hipot.',
        'A_PIB_Otros_Prest. Hipotecas_Gtia Personal', 'A_POFI_Otros_Prest. Consumo_Gtia Personal', 'A_POFI_Otros_Prest. Consumo_Gtia Prenda', 'A_POFI_Otros_Prest. Empresas_Gtia Aval',
        'A_POFI_Otros_Prest. Empresas_Gtia Hipot.', 'A_POFI_Otros_Prest. Empresas_Gtia Personal', 'A_POFI_Otros_Prest. Hipotecas_Gtia Hipot.', 'A_POFI_Otros_Prest. Hipotecas_Gtia Personal',
        'A_POFI_Otros_Prest. Origen_Gtia Aval', 'A_POFI_Otros_Prest. Origen_Gtia Hipot.', 'A_SSCC_BKIA_Prest. Consumo_Gtia Personal', 'A_SSCC_Otros_Prest. Consumo_Gtia Personal',
        'A_SSCC_Otros_Prest. Empresas_Gtia Aval', 'A_SSCC_Otros_Prest. Empresas_Gtia Hipot.','A_SSCC_Otros_Prest. Empresas_Gtia Personal', 'A_SSCC_Otros_Prest. Hipotecas_Gtia Hipot.') 
        and is_floating THEN 'PRESTAMOS.V'
            WHEN dim_1 LIKE '%Tarjetas Credito%' and is_floating THEN 'TARJETAS CREDITO.V'
            WHEN dim_1 LIKE 'P_%Plazo%' and not is_floating THEN 'PLAZO.F'
            WHEN dim_1 LIKE '%Ctas No Remuneradas%' and not is_floating THEN 'CTAS NO REMUNERADAS.F'
            WHEN dim_1 LIKE '%Ctas Remuneradas%' and not is_floating THEN 'CTAS REMUNERADAS.F'
            ELSE 'OTROS'
        END AS JER_PRODUCTO,
            SUM(notional_today) as notional,
        SUM((date_diff('day', balance_date, matdate) * Notional_today)) / SUM(Notional_today) / 30 AS vto_medio,
            SUM(coupon_today*Notional_today)/SUM(Notional_today) as cupon_medio,
            SUM(poolrateadj*notional_today)/SUM(notional_today) as spread_medio,
            SUM((date_diff('day', balance_date, Reprdate) * Notional_today)) / SUM(Notional_today) / 30 AS reprecio_medio,
            SUM (floor_today*notional_today)/SUM(notional_today) as Floor_medio,
            (SUM (effective_duration*notional_today)/SUM (Notional_today))/100 as Dur_Efec_Pond,
            SUM (macaulay_duration*notional_today)/SUM(notional_today) as Dur_Mac_pond
        FROM pro_pichincha_alquid_old.metric
                                WHERE load_id IN (
'{load_ids['cierre_base']}', 
'{load_ids['cierre_up']}', 
'{load_ids['cierre_dwn']}'
)
        and scene_code = 'Base'
        GROUP BY 1 
        ORDER BY CASE JER_PRODUCTO
            WHEN 'ACTIVO CORTO PLAZO.F'then 1
            WHEN 'ANTICIPO NOMINA.F' then 2
            WHEN 'CREDITOS.F' then 3
            WHEN 'DESCUBIERTOS.F' then 4
            WHEN 'PRESTAMOS.F' then 5
            WHEN 'TARJETAS CREDITO.F' then 6
            WHEN 'ACTIVO CORTO PLAZO.V' then 7
            WHEN 'ANTICIPO NOMINA.V' then 8
            WHEN 'CREDITOS.V' then 9
            WHEN  'DESCUBIERTOS.V' then 10
            WHEN  'PRESTAMOS.V' then 11
            when 'TARJETAS CREDITO.V' then 12
        ELSE 999
        END; """

    Informe_Datos_Medios_03= f"""select 
                                nmd_ref,
                                SUM(notional_today) as notional,
                                SUM((date_diff('day', balance_date, matdate) * Notional_today)) / SUM(Notional_today) / 30 AS vto_medio,
                                    SUM(coupon_today*Notional_today)/SUM(Notional_today) as cupon_medio,
                                    SUM(poolrateadj*notional_today)/SUM(notional_today) as spread_medio,
                                    SUM((date_diff('day', Reprdate, matdate) * Notional_today)) / SUM(Notional_today) / 30 AS reprecio_medio,
                                    SUM (floor_today*notional_today)/SUM(notional_today) as Floor_medio,
                                    (SUM (effective_duration*notional_today)/SUM (Notional_today))/100 as Dur_Efec_Pond,
                                    SUM (macaulay_duration*notional_today)/SUM(notional_today) as Dur_Mac_pond
                                FROM pro_pichincha_alquid_old.metric
                                WHERE load_id IN (
                                '{load_ids['cierre_base']}', 
                                '{load_ids['cierre_up']}', 
                                '{load_ids['cierre_dwn']}'
                                                            )
                                and scene_code = 'Base'
                                and dim_1 in ('P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                                        'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
                                        'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas') 
                                group by 1
                                order by case nmd_ref 
                                when 'NMD_MAY_NOFIN' then 1
                                when 'NMD_MAY_FIN' then 2
                                when 'NMD_MIN_NONTR' then 3
                                when 'NMD_MIN_TR' then 4
                                else 999
                                end; """

    Informe_Datos_Medios_03_01= f"""
        WITH ref AS (
        SELECT 'CTAS NO REMUNERADAS_1' AS jerarquia, 'NMD_MAY_NOFIN' AS nmd_ref
        UNION ALL SELECT 'CTAS NO REMUNERADAS_3', 'NMD_MIN_NONTR'
        UNION ALL SELECT 'CTAS NO REMUNERADAS_4', 'NMD_MIN_TR'
        UNION ALL SELECT 'CTAS REMUNERADAS_2', 'NMD_MAY_FIN'
        UNION ALL SELECT 'CTAS REMUNERADAS_3', 'NMD_MIN_NONTR'
        UNION ALL SELECT 'CTAS REMUNERADAS_4', 'NMD_MIN_TR'
    ),

    agg AS (
        SELECT
            case
                when dim_1 in ('P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                    'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas')
                    and nmd_ref = 'NMD_MAY_NOFIN' then 'CTAS NO REMUNERADAS_1'
                when dim_1 in ('P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                    'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas')
                    and nmd_ref = 'NMD_MIN_NONTR' then 'CTAS NO REMUNERADAS_3'
                when dim_1 in ('P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                    'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas')
                    and nmd_ref = 'NMD_MIN_TR' then 'CTAS NO REMUNERADAS_4'
                when dim_1 in ('P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas')
                    and nmd_ref = 'NMD_MAY_FIN' then 'CTAS REMUNERADAS_2'
                when dim_1 in ('P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas')
                    and nmd_ref = 'NMD_MIN_NONTR' then 'CTAS REMUNERADAS_3'
                when dim_1 in ('P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas')
                    and nmd_ref = 'NMD_MIN_TR' then 'CTAS REMUNERADAS_4'
                else 'Otros'
            end as jerarquia,
            nmd_ref,
            SUM(notional_today) as notional,
            SUM(date_diff('day', balance_date, matdate) * Notional_today)
                / NULLIF(SUM(Notional_today), 0) / 30 AS vto_medio,
            SUM(coupon_today * Notional_today)
                / NULLIF(SUM(Notional_today), 0) as cupon_medio,
            SUM(poolrateadj * notional_today)
                / NULLIF(SUM(notional_today), 0) as spread_medio,
            SUM(date_diff('day', Reprdate, matdate) * Notional_today)
                / NULLIF(SUM(Notional_today), 0) / 30 AS reprecio_medio,
            SUM(floor_today * notional_today)
                / NULLIF(SUM(notional_today), 0) as floor_medio,
            (SUM(effective_duration * notional_today)
                / NULLIF(SUM(Notional_today), 0)) / 100 as dur_efec_pond,
            SUM(macaulay_duration * notional_today)
                / NULLIF(SUM(notional_today), 0) as dur_mac_pond

        FROM pro_pichincha_alquid_old.metric
        WHERE load_id IN ('{load_ids['cierre_base']}', '{load_ids['cierre_up']}', '{load_ids['cierre_dwn']}')
            AND scene_code = 'Base'
            AND dim_1 in (
                'P_ECORP_Ctas No Remuneradas', 'P_EESP_Ctas No Remuneradas', 'P_PIB_Ctas No Remuneradas',
                'P_POFI_Ctas No Remuneradas', 'P_SSCC_Ctas No Remuneradas', 'P_SSCCFRO_Ctas No Remuneradas',
                'P_PIB_Ctas Remuneradas', 'P_POFI_Ctas Remuneradas', 'P_SSCC_Ctas Remuneradas'
            )
            AND origin_date >= DATE '{fecha_sql_mes_anterior}'
        GROUP BY 1, nmd_ref
    )

    SELECT
        r.jerarquia,
        r.nmd_ref,
        a.notional,
        a.vto_medio,
        a.cupon_medio,
        a.spread_medio,
        a.reprecio_medio,
        a.floor_medio,
        a.dur_efec_pond,
        a.dur_mac_pond
    FROM ref r
    LEFT JOIN agg a
        ON r.jerarquia = a.jerarquia
    AND r.nmd_ref   = a.nmd_ref
    ORDER BY CASE r.jerarquia
        WHEN 'CTAS NO REMUNERADAS_1' THEN 1
        WHEN 'CTAS NO REMUNERADAS_3' THEN 2
        WHEN 'CTAS NO REMUNERADAS_4' THEN 3
        WHEN 'CTAS REMUNERADAS_3'   THEN 4
        WHEN 'CTAS REMUNERADAS_4'   THEN 5
        WHEN 'CTAS REMUNERADAS_2'   THEN 6
        ELSE 999
    END;
    """

    # Se define nuestro data frame con los datos extrídos de la base de datos
    df_Datos_Medios_01  = pd.read_sql(Informe_Datos_Medios_01 , conn)

    df_Datos_Medios_01_01  = pd.read_sql(Informe_Datos_Medios_01_01 , conn)

    df_Datos_Medios_02  = pd.read_sql(Informe_Datos_Medios_02  , conn)

    df_Datos_Medios_03  = pd.read_sql(Informe_Datos_Medios_03  , conn)

    df_Datos_Medios_03_01  = pd.read_sql(Informe_Datos_Medios_03_01 , conn)

    # Nombre de la hoja donde quieres pegar los datos
    nombre_hoja_Datos_Medios_01  = "Datos Medios"

    nombre_hoja_Datos_Medios_01_01  = "Datos medios ult mes"

    # Acceder a la hoja
    ws_Datos_Medios_01  = wb_2[nombre_hoja_Datos_Medios_01]

    ws_Datos_Medios_01_01  = wb_2[nombre_hoja_Datos_Medios_01_01]


    # Se definen las coordenadas donde queremos que se peguen los datos
    start_row_Datos_Medios_01    = 7
    start_col_Datos_Medios_01    = 5

    start_row_Datos_Medios_01_01    = 7
    start_col_Datos_Medios_01_01   = 5

    start_row_Datos_Medios_02    = 39
    start_col_Datos_Medios_02   = 6

    start_row_Datos_Medios_03    = 99
    start_col_Datos_Medios_03   = 6

    start_row_Datos_Medios_03_01    = 39
    start_col_Datos_Medios_03_01   = 6
    # Límite de filas donde se pegarán los datos (de la 5 a la 23)
    end_row_Datos_Medios_01  = 25

    end_row_Datos_Medios_01_01  = 25

    end_row_Datos_Medios_02  = 55

    end_row_Datos_Medios_03  = 102

    end_row_Datos_Medios_03_01  = 43

    # Filas a evitar
    filas_a_saltar_Datos_Medios_01   = {14, 19, 20}
    fila_actual_Datos_Medios_01   = start_row_Datos_Medios_01  

    filas_a_saltar_Datos_Medios_01_01   = {14, 19, 20}
    fila_actual_Datos_Medios_01_01   = start_row_Datos_Medios_01_01 

    filas_a_saltar_Datos_Medios_02   = {45, 46, 47, 48, 49, 50, 51, 52, 53}
    fila_actual_Datos_Medios_02   = start_row_Datos_Medios_02  
    
    # Escribir los datos en las coordenadas anteriores
    for row in df_Datos_Medios_01  .values:
        # Saltar las filas especificadas
        while fila_actual_Datos_Medios_01   in filas_a_saltar_Datos_Medios_01  :
            fila_actual_Datos_Medios_01   += 1
        
        if fila_actual_Datos_Medios_01   > end_row_Datos_Medios_01  :
            break

        for j, value_Datos_Medios_01   in enumerate(row[1:], start=start_col_Datos_Medios_01  ):
            ws_Datos_Medios_01.cell(row=fila_actual_Datos_Medios_01  , column=j).value = value_Datos_Medios_01  

        fila_actual_Datos_Medios_01  += 1  # Mover a la siguiente fila en la hoja

    print("[DEBUG] - P1")
    # Escribir los datos en las coordenadas anteriores
    for row in df_Datos_Medios_01_01  .values:
        # Saltar las filas especificadas
        while fila_actual_Datos_Medios_01_01  in filas_a_saltar_Datos_Medios_01_01  :
            fila_actual_Datos_Medios_01_01 += 1
        
        if fila_actual_Datos_Medios_01_01  > end_row_Datos_Medios_01_01  :
            break

        for j, value_Datos_Medios_01_01   in enumerate(row[1:], start=start_col_Datos_Medios_01_01  ):
            ws_Datos_Medios_01_01.cell(row=fila_actual_Datos_Medios_01_01  , column=j).value = value_Datos_Medios_01_01  

        fila_actual_Datos_Medios_01_01  += 1  # Mover a la siguiente fila en la hoja

    print("[DEBUG] - P2")
    # Escribir los datos en las coordenadas anteriores
    for row in df_Datos_Medios_02  .values:
        # Saltar las filas especificadas
        while fila_actual_Datos_Medios_02  in filas_a_saltar_Datos_Medios_02  :
            fila_actual_Datos_Medios_02   += 1
        
        if fila_actual_Datos_Medios_02   > end_row_Datos_Medios_02  :
            break

        for j, value_Datos_Medios_02   in enumerate(row[1:], start=start_col_Datos_Medios_02  ):
            ws_Datos_Medios_01.cell(row=fila_actual_Datos_Medios_02  , column=j).value = value_Datos_Medios_02  

        fila_actual_Datos_Medios_02  += 1  # Mover a la siguiente fila en la hoja

    print("[DEBUG] - P3")
    # Escribir los datos en las coordenadas anteriores
    for i, row in enumerate(df_Datos_Medios_03.values, start=start_row_Datos_Medios_03):
        if i > end_row_Datos_Medios_03:
            break  # Detener la escritura si superamos la fila 36
        for j, value_Datos_Medios_03 in enumerate(row[1:], start=start_col_Datos_Medios_03):
            ws_Datos_Medios_01.cell(row=i, column=j).value = value_Datos_Medios_03

    print("[DEBUG] - P4")
    # Escribir los datos en las coordenadas anteriores
    for i, row in enumerate(df_Datos_Medios_03_01.values, start=start_row_Datos_Medios_03_01):
        if i > end_row_Datos_Medios_03_01:
            break  # Detener la escritura si superamos la fila 36
        for j, value_Datos_Medios_03_01 in enumerate(row[2:], start=start_col_Datos_Medios_03_01):
            ws_Datos_Medios_01_01.cell(row=i, column=j).value = value_Datos_Medios_03_01

    print("Todas las hojas han sido actualizadas con los datos de Athena.")


    # 6. Guardar workbooks en memoria y devolverlos
    print("Guardando archivos Excel actualizados en memoria...")
    output_principal = io.BytesIO()
    output_efectos = io.BytesIO()
    output_datos_medios = io.BytesIO()

    wb.save(output_principal)
    wb_1.save(output_efectos)
    wb_2.save(output_datos_medios)
    
    output_principal.seek(0)
    output_efectos.seek(0)
    output_datos_medios.seek(0)

    
    # print("Subiendo archivos Excel actualizados a Google Drive...")
    output_files_info = {}

    # Definir carpetas de salida en Drive
    # reports_folder_id = find_or_create_folder(drive_service, "ALCO_Reports_Generated")
    # if not reports_folder_id:
    #     raise Exception("No se pudo encontrar o crear la carpeta de salida en Google Drive.")

    # Subir principal Excel
    principal_filename = f"ALCO_{mes_actual_nombre}.xlsx"
    principal_buffer = io.BytesIO()
    wb.save(principal_buffer)
    # principal_file_id = upload_file_to_drive(drive_service, principal_filename, principal_buffer.getvalue(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", reports_folder_id)
    # if principal_file_id:
    #     output_files_info["principal"] = (principal_filename, principal_file_id)

    return {
        "principal": (principal_filename, output_principal.getvalue()),
        "efectos": ("Plantilla_Efecto_Balance_Curva_output.xlsx", output_efectos.getvalue()),
        "datos_medios": ("Plantilla_Datos_Medios_output.xlsx", output_datos_medios.getvalue())
    }


##########################################################################################
# ----------------------------------- FASE 2 ------------------------------------------- #
##########################################################################################

# --- Funciones auxiliares de la Fase 2 (las que ya tenías) ---

# --- Configuración y Constantes ---
HISTORICAL_DIR = Path("Diapositivas Históricas COAP")
PLANTILLA_COAP_NOMBRE = "Plantilla COAP_05.pptx"
EXCEL_INPUT_DIR = Path(".")

GENERATED_PPTX_DIR = Path("pruebas")
COMMENTS_TXT_DIR = Path("Comentarios diapositivas")
SLIDE_CAPTURES_DIR = Path("slide_captures_local")
PODCAST_SCRIPTS_DIR = Path("Guiones Podcast") # Nuevo directorio para guiones de podcast

PROMPT_TEMPLATE_PATH = "prompt.txt"
PROMPT_PODCAST_PATH = "prompt_podcast.txt" # Para el nuevo prompt
COMMENT_PLACEHOLDER_NAME = "AI_Comment_Box" 

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    print("❌ ERROR CRÍTICO: La variable de entorno GEMINI_API_KEY no está configurada.")
GEMINI_MODEL_NAME = "gemini-2.0-flash" # Modelo principal
GEMINI_PODCAST_MODEL_NAME = "gemini-2.0-flash" # Podrías usar un modelo más potente para el guion si quieres

MAX_SLIDE_API_ATTEMPTS = 3 
SLIDE_RETRY_DELAY_SECONDS = 60 

if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)

CONFIG_PPT_DEFAULT = {
    "min_chars_for_text_shape": 10,
    "capture_images": COMTYPES_AVAILABLE,
    "image_quality": 95,
    "image_format": "PNG",
    "temp_image_path_pattern": "temp_slide_{slide_num}.{ext}",
    "force_comtypes_for_text": False
}

MESES_ESPANOL_NOMBRE = {
    1: "Enero", 2: "Febrero", 3: "Marzo", 4: "Abril", 5: "Mayo", 6: "Junio",
    7: "Julio", 8: "Agosto", 9: "Septiembre", 10: "Octubre", 11: "Noviembre", 12: "Diciembre"
}
MESES_ESPANOL_INPUT_MAP = {nombre.lower(): num for num, nombre in MESES_ESPANOL_NOMBRE.items()}


class PowerPointAppManager:
     _instance = None
     _app = None
     @classmethod
     def get_instance(cls):
         if not COMTYPES_AVAILABLE: return None
         if cls._instance is None:
             cls._instance = cls.__new__(cls)
             try:
                 cls._app = comtypes.client.GetActiveObject("PowerPoint.Application")
             except (OSError, comtypes.COMError):
                 try:
                     cls._app = comtypes.client.CreateObject("PowerPoint.Application")
                 except (OSError, comtypes.COMError) as e:
                     print(f"ERROR: No se pudo iniciar PowerPoint via comtypes: {e}"); cls._app = None
         return cls._app
     @classmethod
     def close_instance(cls):
         if cls._app and hasattr(cls._app, 'Presentations') and cls._app.Presentations.Count == 0:
             try: cls._app.Quit()
             except Exception as e: print(f"ADVERTENCIA: Error al intentar cerrar PowerPoint: {e}")
         cls._app = None; cls._instance = None

def leer_datos_medios(archivo_excel):
    """
        Función específica para leer la tabla de la hoja de excel Datos Medios, del ALCO.
    """
    nombre_hoja = 'Datos Medios'
    
    # Cargar archivo
    df = pd.read_excel(
        archivo_excel,
        sheet_name=nombre_hoja,
        header=[3, 4, 5],
        engine='openpyxl'
    )

    etiqueta_b7 = df.iloc[0,1]   # Activo sensible
    etiqueta_b8 = df.iloc[1,1]  # Activo sensible- fijo
    etiqueta_b15 = df.iloc[8,1] # Activo Sensible - Variable
    etiqueta_b20 = df.iloc[13,1]    # Activo no sensible 
    etiqueta_b21 = df.iloc[14,1]    # Pasivo sensible
    etiqueta_b22 = df.iloc[15,1]    # Pasivo sensible - fijo
    etiqueta_b31 = df.iloc[24,1]    # Pasivo no sensible

    max_excel_row_processed = 31    # hasta la fila 31
    num_filas_datos = len(df)

    # Inicializar índices
    idx_nivel1 = [None] * num_filas_datos
    idx_nivel2 = [None] * num_filas_datos
    idx_nivel3 = [None] * num_filas_datos

    for i in range(max_excel_row_processed):
        fila_excel_actual = i + 7
        etiqueta_actual_col_b = df.iloc[i, 1]

        # Asignar Nivel 1
        if 7 <= fila_excel_actual <= 20:
            idx_nivel1[i] = etiqueta_b7
        elif 21 <= fila_excel_actual <= max_excel_row_processed: # Ajustar si es necesario
            idx_nivel1[i] = etiqueta_b21

        # Asignar Nivel 2
        if fila_excel_actual == 7: idx_nivel2[i] = etiqueta_b7
        elif 8 <= fila_excel_actual <= 14: idx_nivel2[i] = etiqueta_b8
        elif 15 <= fila_excel_actual <= 19: idx_nivel2[i] = etiqueta_b15
        elif fila_excel_actual == 20: idx_nivel2[i] = etiqueta_b20
        elif fila_excel_actual == 21: idx_nivel2[i] = etiqueta_b21
        elif 22 <= fila_excel_actual <= 30: idx_nivel2[i] = etiqueta_b22
        elif fila_excel_actual == 31: idx_nivel2[i] = etiqueta_b31

        # Asignar Nivel 3
        if fila_excel_actual == 7: idx_nivel3[i] = etiqueta_b7 # B7 es su propio ítem L3
        elif fila_excel_actual == 8: idx_nivel3[i] = etiqueta_b8 # B8 es su propio ítem L3
        elif 9 <= fila_excel_actual <= 14: idx_nivel3[i] = etiqueta_actual_col_b # Texto de B9-B14
        elif fila_excel_actual == 15: idx_nivel3[i] = etiqueta_b15 # B15 es su propio ítem L3
        elif 16 <= fila_excel_actual <= 19: idx_nivel3[i] = etiqueta_actual_col_b # Texto de B16-B19
        elif fila_excel_actual == 20: idx_nivel3[i] = etiqueta_b20 # B20 es un ítem L3
        elif fila_excel_actual == 21: idx_nivel3[i] = etiqueta_b21 # B21 es su propio ítem L3
        elif fila_excel_actual == 22: idx_nivel3[i] = etiqueta_b22 # B22 es su propio ítem L3
        elif 23 <= fila_excel_actual <= 30: idx_nivel3[i] = etiqueta_actual_col_b # Texto de B23-B30
        elif fila_excel_actual == 31: idx_nivel3[i] = etiqueta_b31 # B31 es un ítem L3
        
        # Crear y Asignar el MultiIndex
        nombres_niveles_fila = ['NivelJerarquico1', 'NivelJerarquico2', 'Item']
        
        # Asegurarse de que todas las listas de índices tienen la longitud correcta
        if not (len(idx_nivel1) == num_filas_datos and len(idx_nivel2) == num_filas_datos and len(idx_nivel3) == num_filas_datos):
            raise ValueError("Las longitudes de las listas de niveles de índice no coinciden con el número de filas de datos.")

        multi_indice_filas = pd.MultiIndex.from_arrays(
            [idx_nivel1, idx_nivel2, idx_nivel3],
            names=nombres_niveles_fila
        )

        df.index = multi_indice_filas

    df = df.iloc[:25, :-12]
    df.drop(columns=[df.columns[0], df.columns[1], df.columns[2]], inplace=True)
    
    return df.to_json()

def extract_text_from_shape(shape) -> str: # Sin cambios
    text_parts = []
    if shape.has_text_frame and shape.text_frame.text and shape.text_frame.text.strip():
        text_parts.append(shape.text_frame.text.strip())
    if shape.has_table:
        for row in shape.table.rows:
            row_text = "\t|\t".join([cell.text_frame.text.strip() for cell in row.cells])
            text_parts.append(row_text)
    return "\n".join(text_parts)

def extract_text_from_slide_pptx(slide, min_chars=5) -> str: # Sin cambios
    text_on_slide = []
    try:
        for shape in slide.shapes:
            if hasattr(shape, 'name') and shape.name == COMMENT_PLACEHOLDER_NAME:
                continue
            extracted_text = extract_text_from_shape(shape)
            if extracted_text and len(extracted_text) >= min_chars: text_on_slide.append(extracted_text)
    except Exception as e: print(f"⚠️ ADVERTENCIA: Error extrayendo texto de una forma (python-pptx): {e}")
    return "\n---\n".join(text_on_slide)

def capture_slide_as_image_comtypes(presentation_com, slide_number_com, temp_image_path_str: str, quality=95, img_format="PNG"): # Sin cambios
    if not COMTYPES_AVAILABLE or presentation_com is None: return None
    try:
        slide_com = presentation_com.Slides[slide_number_com]
        slide_com.Export(temp_image_path_str, img_format.upper(), ScaleWidth=1920, ScaleHeight=1080)
        if img_format.upper() in ["JPEG", "JPG"] and os.path.exists(temp_image_path_str):
            with Image.open(temp_image_path_str) as img: img.save(temp_image_path_str, quality=quality, optimize=True)
        return temp_image_path_str
    except Exception as e:
        print(f"❌ ERROR al capturar diapositiva {slide_number_com} en ruta '{temp_image_path_str}': {e}")
        return None

def procesar_presentacion_coap(pptx_path: str, config: Dict, powerpoint_app_com, presentation_com_obj=None, temp_dir_images_base: Optional[Path] = None, mode:str = "actual") -> Dict[str, Any]: # Sin cambios
    if not os.path.exists(pptx_path):
        print(f"❌ ERROR: Archivo PPTX no encontrado: {pptx_path}")
        return {"path": pptx_path, "total_slides": 0, "slides": []}
    presentation_data = {"path": pptx_path, "total_slides": 0, "slides": []}
    try: prs = Presentation(pptx_path)
    except Exception as e: print(f"❌ ERROR: No se pudo abrir {pptx_path} con python-pptx: {e}"); return presentation_data
    presentation_data["total_slides"] = len(prs.slides)
    
    internal_presentation_com = presentation_com_obj
    com_opened_internally = False
    if (config.get("capture_images", False) or config.get("force_comtypes_for_text", False)) and \
       COMTYPES_AVAILABLE and powerpoint_app_com and not internal_presentation_com:
        try:
            internal_presentation_com = powerpoint_app_com.Presentations.Open(os.path.abspath(pptx_path), ReadOnly=True, WithWindow=False)
            com_opened_internally = True
        except Exception as e:
            print(f"❌ ERROR ({mode}): Al abrir presentación COM para {os.path.basename(pptx_path)}: {e}")
            internal_presentation_com = None
    elif presentation_com_obj:
        internal_presentation_com = presentation_com_obj

    for i, slide_pptx in enumerate(prs.slides):
        slide_number_actual = i + 1
        slide_identifier = f"slide_{slide_number_actual}"
        slide_in_processing_range = (3 <= slide_number_actual <= 19 and slide_number_actual not in [13, 14])

        texto_final_diapositiva = None
        ai_comment_box_content_actual = None
        image_path = None
        image_base64 = None
        slide_title_text = None
        normalized_slide_title = None

        if slide_pptx.shapes.title:
            try:
                title_candidate = slide_pptx.shapes.title.text.strip()
                if title_candidate:
                    slide_title_text = title_candidate
                    normalized_slide_title = title_candidate.lower().strip()
            except Exception as e_title:
                print(f"⚠️ ADVERTENCIA ({mode}): No se pudo extraer el título de la diapo {slide_number_actual}: {e_title}")

        if slide_in_processing_range:
            texto_general_diapositiva = extract_text_from_slide_pptx(slide_pptx, config.get("min_chars_for_text_shape", 10))
            texto_final_diapositiva = texto_general_diapositiva

            try:
                for shape in slide_pptx.shapes:
                    if hasattr(shape, "name") and shape.name == COMMENT_PLACEHOLDER_NAME:
                        if shape.has_text_frame and shape.text_frame.text and shape.text_frame.text.strip():
                            ai_comment_box_content_actual = shape.text_frame.text.strip()
                        break 
            except Exception as e_shape_name:
                print(f"⚠️ ADVERTENCIA ({mode}): Error accediendo a shape.name en diapo {slide_number_actual} de {os.path.basename(pptx_path)}: {e_shape_name}")

            if config.get("force_comtypes_for_text", False) and COMTYPES_AVAILABLE and internal_presentation_com:
                try:
                    slide_com_obj = internal_presentation_com.Slides[slide_number_actual]
                    text_comtypes_parts = []
                    for shape_com_obj in slide_com_obj.Shapes:
                        if shape_com_obj.HasTextFrame and shape_com_obj.TextFrame.HasText:
                            try:
                                if shape_com_obj.Name == COMMENT_PLACEHOLDER_NAME: continue
                            except AttributeError: pass
                            text_comtypes_parts.append(shape_com_obj.TextFrame.TextRange.Text)
                    texto_comtypes = "\n---\n".join(text_comtypes_parts)
                    if texto_general_diapositiva is not None and len(texto_comtypes) > len(texto_general_diapositiva) + 20 :
                        texto_final_diapositiva = texto_comtypes
                    elif texto_general_diapositiva is None and texto_comtypes:
                        texto_final_diapositiva = texto_comtypes
                except Exception as e: print(f"⚠️ ADVERTENCIA ({mode}): Error extrayendo texto comtypes diapo {slide_number_actual}: {e}")
            
            if config.get("capture_images", False) and temp_dir_images_base and COMTYPES_AVAILABLE and internal_presentation_com:
                img_format = config.get("image_format", "PNG").lower()
                base_pattern = config.get("temp_image_path_pattern", "temp_slide_{slide_num}.{ext}")
                filename_prefix = "actual_" if mode == "actual" else "anterior_"
                image_filename_formatted = base_pattern.format(slide_num=slide_number_actual, ext=img_format)
                image_filename_with_prefix = f"{filename_prefix}{image_filename_formatted}"
                absolute_temp_dir = temp_dir_images_base.resolve()
                absolute_temp_image_path = absolute_temp_dir / image_filename_with_prefix
                image_path_str = str(absolute_temp_image_path)
                returned_image_path = capture_slide_as_image_comtypes(internal_presentation_com, slide_number_actual, image_path_str, config.get("image_quality",95), config.get("image_format","PNG"))
                
                if returned_image_path and os.path.exists(returned_image_path):
                    image_path = returned_image_path
                    try:
                        with open(returned_image_path, "rb") as img_file: image_base64 = base64.b64encode(img_file.read()).decode('utf-8')
                    except Exception as e: print(f"❌ ERROR ({mode}): al codificar imagen base64 diapo {slide_number_actual}: {e}"); image_base64 = None
                elif returned_image_path is None:
                     print(f"INFO ({mode}): La captura de imagen para diapositiva {slide_number_actual} falló.")
        
        slide_data_entry = {
            "slide_number": slide_number_actual, "text": texto_final_diapositiva,
            "ai_comment_box_content": ai_comment_box_content_actual,
            "slide_title": slide_title_text, "normalized_title": normalized_slide_title,
            "image_path": image_path, "image_base64": image_base64, "identifier": slide_identifier
            }
        presentation_data["slides"].append(slide_data_entry)

    if com_opened_internally and internal_presentation_com:
        try: internal_presentation_com.Close()
        except Exception as e: print(f"❌ ERROR ({mode}): Al cerrar presentación COM interna para {os.path.basename(pptx_path)}: {e}")
    return presentation_data

def add_comment_to_slide(slide_pptx_obj, comment_text: str, placeholder_name: str = COMMENT_PLACEHOLDER_NAME): # Sin cambios
    target_shape = None
    try:
        for shape in slide_pptx_obj.shapes:
            if hasattr(shape, 'name') and shape.name == placeholder_name:
                target_shape = shape
                break
        
        if target_shape:
            if not target_shape.has_text_frame:
                print(f"❌ ERROR (Diapo {slide_pptx_obj.slide_id if hasattr(slide_pptx_obj, 'slide_id') else 'ID desc.'}): El placeholder '{placeholder_name}' no tiene un marco de texto.")
                return

            tf = target_shape.text_frame
            tf.clear() 

            lines = comment_text.split('\n')
            for line_text in lines:
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT 
                p.line_spacing = 1.0
                p.font.name = 'Poppins'
                p.font.size = Pt(7.5)
                
                processed_line_text = line_text.strip()
                is_bullet_line = False
                if processed_line_text.startswith("- "):
                    is_bullet_line = True
                    p.level = 0 
                    processed_line_text = processed_line_text[2:].strip() 
                
                segments = re.split(r'(\*\*\*.*?\*\*\*|__.*?__)', processed_line_text)
                
                if not processed_line_text and is_bullet_line :
                     run = p.add_run()
                     run.text = " " 
                else:
                    for segment in segments:
                        if not segment: continue

                        run = p.add_run()
                        if segment.startswith('***') and segment.endswith('***'):
                            run.text = segment[3:-3]
                            run.font.bold = True
                        elif segment.startswith('__') and segment.endswith('__'):
                            run.text = segment[2:-2]
                            run.font.underline = True 
                        else:
                            run.text = segment
            
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # print(f"INFO: Comentario (Poppins 7.5, Izquierda, Centro Vertical, Markdown) añadido a '{placeholder_name}' en Diapo ID {slide_pptx_obj.slide_id if hasattr(slide_pptx_obj, 'slide_id') else 'ID desc.'}.")
        else: 
            print(f"⚠️ ADVERTENCIA (Diapo {slide_pptx_obj.slide_id if hasattr(slide_pptx_obj, 'slide_id') else 'ID desc.'}): Placeholder '{placeholder_name}' no encontrado. El comentario NO se añadió.")
    except Exception as e:
        print(f"❌ ERROR añadiendo comentario formateado al placeholder '{placeholder_name}' en Diapo ID {slide_pptx_obj.slide_id if hasattr(slide_pptx_obj, 'slide_id') else 'ID desc.'}: {e}")

def call_gemini_api_with_backoff( # Sin cambios
    content_parts_for_api: List[Any], 
    api_key: str = GEMINI_API_KEY, model_name: str = GEMINI_MODEL_NAME,
    max_retries: int = 3, initial_wait: int = 5,
    generation_config_override: Optional[Dict] = None
):
    if not api_key:
        print("❌ ERROR API: Clave API Gemini no proporcionada."); return "Error: Clave API no configurada."
    
    base_generation_config = {
        "temperature": 0.3, "top_p": 0.95, "top_k": 64,
        "max_output_tokens": 8192, "response_mime_type": "text/plain"
    }
    final_generation_config = base_generation_config.copy()
    if generation_config_override:
        print(f"INFO API: Aplicando configuración de generación anulada: {generation_config_override}")
        final_generation_config.update(generation_config_override)

    safety_settings = [
        {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
        {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
        {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
        {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
    ]
    try:
        model = genai.GenerativeModel(
            model_name=model_name, safety_settings=safety_settings,
            generation_config=final_generation_config
        )
    except Exception as e:
        print(f"❌ ERROR API: No se pudo inicializar el modelo GenerativeModel ('{model_name}'). Error: {e}")
        return f"Error: Fallo al inicializar el modelo - {e}"

    for attempt in range(max_retries):
        try:
            response = model.generate_content(content_parts_for_api)
            if response.parts: return response.text.strip()
            if hasattr(response, 'prompt_feedback') and response.prompt_feedback and response.prompt_feedback.block_reason:
                block_reason = response.prompt_feedback.block_reason
                print(f"❌ ERROR API: Prompt bloqueado. Razón: {block_reason}")
                if response.prompt_feedback.safety_ratings:
                    for rating in response.prompt_feedback.safety_ratings: print(f"   - Categoría: {rating.category}, Probabilidad: {rating.probability}")
                return f"Error: Prompt bloqueado. Razón: {block_reason}"
            else:
                candidate_count = 0
                if hasattr(response, 'candidates') and response.candidates:
                    candidate_count = len(response.candidates)
                    if candidate_count > 0 and hasattr(response.candidates[0], 'finish_reason') and response.candidates[0].finish_reason != 1 :
                         print(f"❌ ERROR API: Candidato finalizado con razón '{response.candidates[0].finish_reason}'.")
                         return f"Error: API finalizó con razón no esperada: {response.candidates[0].finish_reason}"
                print(f"❌ ERROR API: Respuesta vacía o inesperada (Candidatos: {candidate_count}). Detalle: {response}")
                return "Error: Respuesta vacía o inesperada de la API."
        except Exception as e:
            print(f"❌ ERROR API (intento {attempt + 1}/{max_retries}): {e}")
            if attempt < max_retries - 1:
                wait_time = initial_wait * (2 ** attempt); time.sleep(wait_time)
                print(f"   Reintentando en {wait_time} segundos..."); 
            else: return f"Error: Fallo en la API tras {max_retries} intentos. Último error: {e}"
    return "Error: Fallo en la API tras múltiples intentos (inesperado)."

def normalize_bullets_for_ai_prompt(text: Optional[str]) -> str: # Sin cambios
    if not text:
        return "No hay comentario de referencia disponible."
    lines = text.split('\n')
    processed_lines = []
    bullet_pattern = re.compile(r"^\s*([•\*\-◦❖➢➣➤])\s*")
    for line in lines:
        match = bullet_pattern.match(line)
        if match:
            processed_lines.append("- " + line[match.end():])
        else:
            processed_lines.append(line)
    return "\n".join(processed_lines)

def generar_comentarios_coap_local(
    current_pptx_path: str,
    previous_pptx_path: str,
    output_pptx_path: str,
    output_comments_txt_path: str,
    excel_path_actual: str,
    mes_actual_str: str,
    mes_anterior_str: str,
    config_ppt: dict,
    prompt_template_path: str,
    slide_captures_dir: Path
):
    """
    Orquesta la generación de comentarios para una presentación de PowerPoint.
    Esta función está adaptada para ser llamada programáticamente, recibiendo todas las rutas
    y configuraciones como argumentos.
    """
    
    # --- Constantes y Configuración ---
    COMMENT_PLACEHOLDER_NAME = "AI_Comment_Box"
    GEMINI_MODEL_NAME = "gemini-1.5-flash" 
    MAX_SLIDE_API_ATTEMPTS = 3
    SLIDE_RETRY_DELAY_SECONDS = 60

    # 1. Cargar la plantilla del prompt desde el archivo
    try:
        with open(prompt_template_path, 'r', encoding='utf-8') as f:
            prompt_template_base = f.read()
        print(f"INFO: Plantilla de prompt cargada desde: {prompt_template_path}")
    except Exception as e:
        raise IOError(f"No se pudo leer la plantilla de prompt en {prompt_template_path}: {e}")

    # 2. Inicializar el gestor de PowerPoint
    powerpoint_app_com = PowerPointAppManager.get_instance()

    # 3. Procesar las presentaciones para extraer datos
    print(f"INFO: Procesando presentación actual: {os.path.basename(current_pptx_path)}...")
    current_presentation_data = procesar_presentacion_coap(
        current_pptx_path, config_ppt, powerpoint_app_com, slide_captures_dir, "actual", COMMENT_PLACEHOLDER_NAME
    )

    previous_presentation_proc_data = {"slides": []}
    if previous_pptx_path and os.path.exists(previous_pptx_path):
        print(f"INFO: Procesando presentación de referencia anterior: {os.path.basename(previous_pptx_path)}...")
        previous_presentation_proc_data = procesar_presentacion_coap(
            previous_pptx_path, config_ppt, powerpoint_app_com, slide_captures_dir, "anterior", COMMENT_PLACEHOLDER_NAME
        )
    else:
        print("INFO: No se proporcionó o no se encontró presentación anterior de referencia.")

    # 4. Cargar la presentación de salida para empezar a modificarla
    try:
        prs_current_output = Presentation(current_pptx_path)
    except Exception as e:
        # Importante: cerrar la instancia de PowerPoint si falla la carga para no dejar procesos abiertos
        PowerPointAppManager.close_instance()
        raise IOError(f"No se pudo abrir {current_pptx_path} para añadir comentarios: {e}")

    # 5. Bucle principal para generar comentarios por diapositiva
    all_comments_for_txt_file = []
    
    for slide_data_actual in current_presentation_data['slides']:
        slide_number = slide_data_actual['slide_number']
        slide_title = slide_data_actual.get('slide_title', 'Sin Título')

        # Filtrar diapositivas a procesar según tu lógica
        if not (3 <= slide_number <= 19 and slide_number not in [13, 14]):
            continue
        
        print(f"\n--- Generando para Diapositiva {slide_number} (Título: {slide_title}) ---")
        
        # Extraer datos de la diapositiva actual
        texto_diapositiva_actual = slide_data_actual.get('text') or "No hay texto extraíble de la diapositiva actual."
        imagen_base64_actual = slide_data_actual.get('image_base64')
        current_slide_normalized_title = slide_data_actual.get('normalized_title')

        # Buscar datos de la diapositiva del mes anterior
        comentario_referencia_bruto = "No hay comentario de referencia disponible."
        imagen_base64_anterior = None
        
        # Lógica para encontrar slide anterior (por título, luego por número)
        found_previous_by_title = False
        if previous_presentation_proc_data["slides"] and current_slide_normalized_title:
            for prev_slide in previous_presentation_proc_data["slides"]:
                if prev_slide.get("normalized_title") == current_slide_normalized_title:
                    comentario_referencia_bruto = prev_slide.get("ai_comment_box_content") or prev_slide.get("text") or comentario_referencia_bruto
                    imagen_base64_anterior = prev_slide.get("image_base64")
                    found_previous_by_title = True
                    break
        
        if not found_previous_by_title and previous_presentation_proc_data["slides"] and 0 < slide_number <= len(previous_presentation_proc_data["slides"]):
             prev_slide = previous_presentation_proc_data["slides"][slide_number - 1]
             comentario_referencia_bruto = prev_slide.get("ai_comment_box_content") or prev_slide.get("text") or comentario_referencia_bruto
             imagen_base64_anterior = prev_slide.get("image_base64")


        # Normalizar y preparar contexto
        comentario_placeholder_anterior_normalizado = normalize_bullets_for_ai_prompt(comentario_referencia_bruto)
        
        contexto_excel_para_ia = "Contexto Excel no disponible."
        if excel_path_actual and os.path.exists(excel_path_actual):
            print(f"INFO: Leyendo contexto del archivo Excel: {os.path.basename(excel_path_actual)}")
            try:
                contexto_excel_para_ia = leer_datos_medios(excel_path_actual)
            except Exception as e:
                print(f"ERROR: No se pudo leer el contexto del excel. {e}")
                contexto_excel_para_ia = "Error al leer el archivo Excel."

        
        # Construir el prompt para la API
        prompt_aplicado = prompt_template_base.replace("{MES_DATOS_ACTUALES_STR}", mes_actual_str)
        prompt_aplicado = prompt_aplicado.replace("{MES_DATOS_ANTERIORES_STR}", mes_anterior_str)
        prompt_aplicado = prompt_aplicado.replace("{TEXTO_DIAPOSITIVA_ACTUAL}", texto_diapositiva_actual)
        prompt_aplicado = prompt_aplicado.replace("{COMENTARIO_DIAPOSITIVA_ANTERIOR}", comentario_placeholder_anterior_normalizado)
        prompt_aplicado = prompt_aplicado.replace("{CONTEXTO_DATOS_ADICIONALES}", contexto_excel_para_ia)


        # Construir el payload para la API (texto + imágenes)
        api_content_parts = [prompt_aplicado]
        image_format_mime = f"image/{config_ppt.get('image_format', 'PNG').lower()}"

        if imagen_base64_actual:
            api_content_parts.extend([
                f"\n\n--- INICIO IMAGEN DIAPOSITIVA ACTUAL ({mes_actual_str}) ---",
                {"mime_type": image_format_mime, "data": imagen_base64_actual},
                f"--- FIN IMAGEN DIAPOSITIVA ACTUAL ({mes_actual_str}) ---"
            ])
        
        if slide_number != 3 and imagen_base64_anterior:
            api_content_parts.extend([
                f"\n\n--- INICIO IMAGEN DIAPOSITIVA MES ANTERIOR ({mes_anterior_str}) ---",
                {"mime_type": image_format_mime, "data": imagen_base64_anterior},
                f"--- FIN IMAGEN DIAPOSITIVA MES ANTERIOR ({mes_anterior_str}) ---"
            ])

        # Llamar a la API
        print(f"INFO (Diapo {slide_number}): Enviando solicitud a la API de Gemini...")
        nuevo_comentario_texto = call_gemini_api_with_backoff(
            api_content_parts,
            model_name=GEMINI_MODEL_NAME,
            max_retries=MAX_SLIDE_API_ATTEMPTS,
            initial_wait=SLIDE_RETRY_DELAY_SECONDS
        )

        # Añadir comentario a la presentación
        slide_pptx_para_comentario = prs_current_output.slides[slide_number - 1]
        add_comment_to_slide(slide_pptx_para_comentario, nuevo_comentario_texto, COMMENT_PLACEHOLDER_NAME)
        
        all_comments_for_txt_file.append(f"--- Diapositiva {slide_number} (Título: {slide_title}) ---\n{nuevo_comentario_texto}\n\n")

    # 6. Guardar los archivos de salida
    try:
        prs_current_output.save(output_pptx_path)
        print(f"\n✅ INFO: Presentación con comentarios guardada en: {output_pptx_path}")
    except Exception as e:
        print(f"❌ ERROR CRÍTICO al guardar presentación final: {e}")

    try:
        with open(output_comments_txt_path, 'w', encoding='utf-8') as f:
            f.write(f"Comentarios generados para: {os.path.basename(current_pptx_path)}\n")
            f.write(f"Datos del mes: {mes_actual_str}\n\n")
            f.writelines(all_comments_for_txt_file)
        print(f"✅ INFO: Archivo de texto con comentarios guardado en: {output_comments_txt_path}")
    except Exception as e:
        print(f"❌ ERROR al guardar archivo de texto con comentarios: {e}")

def generar_guion_podcast(
    comentarios_coap_txt_path: str, 
    prompt_podcast_path: str, 
    output_podcast_script_path: str,
    podcast_model_name: str, 
    podcast_gen_config: dict
):
    """
    Genera un guion de podcast a partir de los comentarios de un archivo de texto.
    """
    try:
        with open(comentarios_coap_txt_path, 'r', encoding='utf-8') as f:
            contenido_comentarios_coap = f.read()
    except Exception as e:
        print(f"ERROR: No se pudo leer el archivo de comentarios COAP: {e}")
        return

    try:
        with open(prompt_podcast_path, 'r', encoding='utf-8') as f:
            prompt_podcast_template = f.read()
    except Exception as e:
        print(f"ERROR: No se pudo leer la plantilla de prompt para podcast: {e}")
        return

    prompt_final_podcast = prompt_podcast_template.replace("{TEXTO_COMENTARIOS_PPTX}", contenido_comentarios_coap)
    
    print("INFO: Generando guion de podcast con la API de Gemini...")
    guion_podcast_texto = call_gemini_api_with_backoff(
        [prompt_final_podcast],
        model_name=podcast_model_name,
        generation_config_override=podcast_gen_config
    )

    if "Error:" in guion_podcast_texto:
        print(f"ERROR al generar guion de podcast: {guion_podcast_texto}")
    else:
        with open(output_podcast_script_path, 'w', encoding='utf-8') as f:
            f.write(guion_podcast_texto)
        print(f"✅ INFO: Guion de podcast guardado en: {output_podcast_script_path}")


# Función principal que orquesta la Fase 2
def ejecutar_fase_2(mes_cierre_input_str, archivos_bytes, gemini_api_key, config_ppt):
    """
    Ejecuta toda la lógica de la Fase 2: Generación de comentarios y guion de podcast.

    Args:
        mes_cierre_input_str (str): El mes de cierre en formato 'Mes Año'.
        archivos_bytes (dict): Contenido en bytes de los archivos necesarios.
        gemini_api_key (str): La clave API para Gemini.
        config_ppt (dict): Diccionario de configuración para el procesamiento de PPTX.

    Returns:
        dict: Un diccionario con los nombres y contenido en bytes de los archivos generados.
    """
    # Usar un directorio temporal para manejar los archivos que necesitan una ruta física
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        
        # --- Escribir archivos necesarios en el directorio temporal ---
        current_pptx_path = temp_path / "plantilla_actual.pptx"
        current_pptx_path.write_bytes(archivos_bytes['plantilla_pptx'])

        previous_pptx_path = temp_path / "plantilla_anterior.pptx"
        previous_pptx_path.write_bytes(archivos_bytes['pptx_anterior'])

        excel_path_actual = temp_path / "datos_alco.xlsx"
        excel_path_actual.write_bytes(archivos_bytes['alco_excel'])
        
        prompt_template_path = temp_path / "prompt.txt"
        prompt_template_path.write_text(archivos_bytes['prompt_main'], encoding='utf-8')

        prompt_podcast_path = temp_path / "prompt_podcast.txt"
        prompt_podcast_path.write_text(archivos_bytes['prompt_podcast'], encoding='utf-8')
        
        # Crear subdirectorios que el script espera
        GENERATED_PPTX_DIR = temp_path / "generated"
        COMMENTS_TXT_DIR = temp_path / "comments"
        PODCAST_SCRIPTS_DIR = temp_path / "podcast"
        SLIDE_CAPTURES_DIR = temp_path / "captures"
        for d in [GENERATED_PPTX_DIR, COMMENTS_TXT_DIR, PODCAST_SCRIPTS_DIR, SLIDE_CAPTURES_DIR]:
            d.mkdir(parents=True, exist_ok=True)

        # --- Adaptación de la lógica de 'main_process' ---
        MESES_ESPANOL_NOMBRE = {
            1: "Enero", 2: "Febrero", 3: "Marzo", 4: "Abril", 5: "Mayo", 6: "Junio",
            7: "Julio", 8: "Agosto", 9: "Septiembre", 10: "Octubre", 11: "Noviembre", 12: "Diciembre"
        }
        MESES_ESPANOL_INPUT_MAP = {nombre.lower(): num for num, nombre in MESES_ESPANOL_NOMBRE.items()}

        try:
            partes_fecha_cierre = mes_cierre_input_str.split()
            nombre_mes_cierre_input, ano_cierre_str = partes_fecha_cierre
            num_mes_cierre = MESES_ESPANOL_INPUT_MAP[nombre_mes_cierre_input.lower()]
            reporting_date = datetime(int(ano_cierre_str), num_mes_cierre, 1)
        except Exception as e:
            raise ValueError(f"Formato de fecha de cierre incorrecto: {e}")

        data_date_actual = reporting_date - relativedelta(months=1)
        MES_DATOS_ACTUALES_STR = f"{MESES_ESPANOL_NOMBRE[data_date_actual.month]} {data_date_actual.year}"
        data_date_anterior = data_date_actual - relativedelta(months=1)
        MES_DATOS_ANTERIORES_STR = f"{MESES_ESPANOL_NOMBRE[data_date_anterior.month]} {data_date_anterior.year}"
        
        # Configurar API Key
        genai.configure(api_key=gemini_api_key)

        # Definir nombres de archivo de salida
        mes_cierre_filename_part = reporting_date.strftime('%B_%Y')
        output_filename_base = f"Copia_de_{mes_cierre_filename_part}_COAP_Pichincha"
        output_pptx_path = GENERATED_PPTX_DIR / f"{output_filename_base}.pptx"
        output_comments_txt_path = COMMENTS_TXT_DIR / f"{output_filename_base}_comentarios.txt"
        output_podcast_script_path = PODCAST_SCRIPTS_DIR / f"{output_filename_base}_guion_podcast.txt"
        
        # --- Llamada a las funciones principales de la Fase 2 ---
        print("\n--- Iniciando el proceso de generación de comentarios para la presentación ---")
        generar_comentarios_coap_local(
            current_pptx_path=str(current_pptx_path),
            previous_pptx_path=str(previous_pptx_path),
            output_pptx_path=str(output_pptx_path),
            output_comments_txt_path=str(output_comments_txt_path),
            excel_path_actual=str(excel_path_actual),
            mes_actual_str=MES_DATOS_ACTUALES_STR,
            mes_anterior_str=MES_DATOS_ANTERIORES_STR,
            config_ppt=config_ppt,
            prompt_template_path=str(prompt_template_path),
            slide_captures_dir=SLIDE_CAPTURES_DIR
        )

        # 2. Generar guion de podcast
        # Este paso solo se ejecuta si el archivo de comentarios se creó correctamente.
        if output_comments_txt_path.exists():
            print("\n--- Iniciando el proceso de generación de guion de podcast ---")
            generar_guion_podcast(
                comentarios_coap_txt_path=str(output_comments_txt_path),
                prompt_podcast_path=str(prompt_podcast_path),
                output_podcast_script_path=str(output_podcast_script_path),
                podcast_model_name="gemini-1.5-flash",  # Modelo de IA para el guion
                podcast_gen_config={"max_output_tokens": 8192} # Configuración para la API
            )
        else:
            print(f"ADVERTENCIA: No se generará guion de podcast porque no se encontró el archivo de comentarios en: {output_comments_txt_path}")
        
        # --- SIMULACIÓN DEL PROCESO (reemplazar con llamadas reales a tus funciones) ---
        print("Simulando la creación de archivos de Fase 2 en directorio temporal...")
        # Simular guardado de archivos
        shutil.copy(current_pptx_path, output_pptx_path) # Simplemente copiamos la plantilla como salida
        output_comments_txt_path.write_text(f"Comentarios simulados para {MES_DATOS_ACTUALES_STR}", encoding='utf-8')
        output_podcast_script_path.write_text(f"Guion de podcast simulado para {MES_DATOS_ACTUALES_STR}", encoding='utf-8')
        # --- FIN DE SIMULACIÓN ---


        # --- Leer los archivos generados para devolverlos como bytes ---
        resultados_bytes = {}
        if output_pptx_path.exists():
            resultados_bytes["pptx"] = (output_pptx_path.name, output_pptx_path.read_bytes())
        if output_comments_txt_path.exists():
            resultados_bytes["comentarios"] = (output_comments_txt_path.name, output_comments_txt_path.read_bytes())
        if output_podcast_script_path.exists():
            resultados_bytes["podcast"] = (output_podcast_script_path.name, output_podcast_script_path.read_bytes())

        # Limpiar la instancia de PowerPoint si se usó
        PowerPointAppManager.close_instance()
        
        return resultados_bytes